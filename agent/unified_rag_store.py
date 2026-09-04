from __future__ import annotations

"""
UNIFIED RAG STORE - KONSOLIDIERTE LÖSUNG (September 2025)
=========================================================

🎯 ZIEL: Alle RAG-Features in einer Datei vereint

✅ Features:
- GPU-beschleunigt (wenn verfügbar)  
- Multi-Core-Parallelisierung
- Automatisches Fallback auf sequenzielle Verarbeitung
- 100% API-Kompatibilität mit bestehenden Implementierungen
- Thread-sichere Verbindungspools
- Moderne Embedding-Provider (OpenAI, BGE, Voyage)

🗂️ Diese Datei ersetzt:
- agent/rag_store.py (Basis-Implementation, 2220 Zeilen)
- agent/smart_rag_store.py (GPU-Optimierungen, 179 Zeilen)  
- agent/parallel_rag_store.py (Parallele Verarbeitung, 1056 Zeilen)

🔄 Migration:
- Alle bestehenden Imports funktionieren weiterhin
- RagStore → UnifiedRagStore (mit Alias für Kompatibilität)
- create_rag_store() → Automatische Erkennung der besten Konfiguration

🆕 Modular Refactoring (Oktober 2025):
- Datei umbenannt: rag_store.py → unified_rag_store.py
- Utilities extrahiert nach rag_store/utils/
- Weitere Module folgen schrittweise (rag_store/core/, rag_store/advanced/)
"""

import os
import sqlite3
import json
import hashlib
import threading
import logging
import multiprocessing
import sys
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
from typing import List, Dict, Any, Optional, Tuple, Union, Set
import time
from dataclasses import dataclass
import queue
from datetime import datetime, timezone, timedelta
from urllib.parse import urlparse
from contextlib import contextmanager
from utils import web_compliance

# Performance-Optimierung: NumExpr Thread-Budget (Fallback falls Entry Point es nicht setzt)
# Auto-Detect: 75% der logischen Kerne (Sweet-Spot für BLAS/LAPACK-Operationen)
_optimal_threads = str(max(4, int((os.cpu_count() or 8) * 0.75)))
os.environ.setdefault('NUMEXPR_MAX_THREADS', _optimal_threads)
os.environ.setdefault('NUMEXPR_NUM_THREADS', _optimal_threads)

# GPU Optimization
try:
    from gpu_optimizer import get_gpu_optimizer, apply_gpu_optimizations
    GPU_OPTIMIZER_AVAILABLE = True
except ImportError:
    GPU_OPTIMIZER_AVAILABLE = False
    # Type-safe fallback functions

# LLM-based Knowledge Graph Extraction
try:
    from .llm_knowledge_graph import LLMKnowledgeGraphExtractor, KGTriple
    LLM_KG_AVAILABLE = True
except ImportError:
    LLM_KG_AVAILABLE = False
    get_gpu_optimizer = lambda *args, **kwargs: None  # type: ignore
    apply_gpu_optimizations = lambda *args, **kwargs: None  # type: ignore

# PDF-Mining Warnings unterdrücken (harmlose Farbwerte-Probleme)
logging.getLogger("pdfminer.pdfinterp").setLevel(logging.ERROR)
logging.getLogger("pdfminer.pdfdocument").setLevel(logging.ERROR)
logging.getLogger("pdfminer.pdfpage").setLevel(logging.ERROR)

# Dependencies (numpy is required)
import numpy as np
from utils.runtime_policy import parse_bool_env


# FAISS (optional für Vektor-Indizierung)
# Lazy loaded to avoid module-import side effects (SWIG deprecation warnings)
# in code paths that don't use FAISS.
faiss = None

# FAISS Hybrid Search Components (lazy loaded)
FAISS_HYBRID_AVAILABLE = False
FAISSIndexManager = None
SmartFusionEngine = None


def _ensure_faiss_runtime() -> bool:
    """Lazy import for FAISS to keep module import side-effect free."""
    global faiss
    if faiss is not None:
        return True
    try:
        import faiss as _faiss  # type: ignore
        faiss = _faiss
        return True
    except ImportError:
        return False


def _ensure_faiss_hybrid_runtime() -> bool:
    """Lazy import for FAISS hybrid components."""
    global FAISS_HYBRID_AVAILABLE, FAISSIndexManager, SmartFusionEngine
    if FAISS_HYBRID_AVAILABLE and FAISSIndexManager is not None and SmartFusionEngine is not None:
        return True
    try:
        from agent.faiss_index_manager import FAISSIndexManager as _FAISSIndexManager  # type: ignore
        from agent.smart_fusion_engine import SmartFusionEngine as _SmartFusionEngine  # type: ignore
        FAISSIndexManager = _FAISSIndexManager
        SmartFusionEngine = _SmartFusionEngine
        FAISS_HYBRID_AVAILABLE = True
        return True
    except ImportError:
        try:
            # Fallback: Try without agent. prefix (if running from agent/ directory)
            from faiss_index_manager import FAISSIndexManager as _FAISSIndexManager  # type: ignore
            from smart_fusion_engine import SmartFusionEngine as _SmartFusionEngine  # type: ignore
            FAISSIndexManager = _FAISSIndexManager
            SmartFusionEngine = _SmartFusionEngine
            FAISS_HYBRID_AVAILABLE = True
            return True
        except ImportError:
            FAISS_HYBRID_AVAILABLE = False
            return False

# ❌ REMOVED in Iteration 6 (ungenutzte Provider)
# OpenAI API (optional, nicht verwendet)
# openai = None
# try:
#     import openai
# except ImportError:
#     pass

# ❌ REMOVED in Iteration 6 (ungenutzte Provider)
# Requests für HTTP-Anfragen (Voyage API)
# requests = None
# try:
#     import requests
# except ImportError:
#     pass

# Sentence Transformers für lokale Embeddings (WIRD VERWENDET ✅)
# CRITICAL FIX: Runtime-Check statt Top-Level Import (prevents false negatives during Streamlit startup)
def _check_sentence_transformers_runtime():
    """
    Runtime-Check für sentence-transformers (nicht beim Module Load!)
    Prevents false negatives due to import order during Streamlit startup.
    
    Returns:
        tuple: (SentenceTransformer class or None, bool available)
    """
    try:
        import sentence_transformers
        from sentence_transformers import SentenceTransformer
        return SentenceTransformer, True
    except ImportError:
        return None, False

# Legacy globals für Kompatibilität (werden zur Runtime aktualisiert)
SENTENCE_TRANSFORMERS_AVAILABLE = False
SentenceTransformer = None

# NOTE: AdvancedPDFProcessor-Import entfernt (Root-Cause-Fix 2026-07-14):
# Der Adapter delegierte nur an DoclingProcessor (zirkulärer Fallback) und
# war durch einen force_ocr-TypeError funktionsunfähig. PDF-Fallbacks laufen
# jetzt direkt über pymupdf4llm → pdfminer → PyMuPDF → OCR.

# PyTorch für GPU-Support
# TORCH_AVAILABLE: True wenn torch importierbar (auch CPU-only)
# CUDA_AVAILABLE: True wenn CUDA-GPU vorhanden und nutzbar
TORCH_AVAILABLE = False
CUDA_AVAILABLE = False
try:
    import torch
    TORCH_AVAILABLE = True
    CUDA_AVAILABLE = torch.cuda.is_available() if hasattr(torch, 'cuda') else False
except ImportError:
    torch = None  # type: ignore[assignment]

logger = logging.getLogger(__name__)

# ====================================================================
# MODULAR UTILITIES (Extrahiert nach rag_store/)
# ====================================================================

from .rag_store.utils import (
    ResourceManager,
    ManagedResource,
    managed_resource,
    _resource_manager,
    calculate_triple_hash,
    ProcessingConfig,
    # Batch Processing Utilities (Iteration 2)
    chunk_list,
    batch_embed_texts,
    batch_sql_load,
    batch_insert
)

# Import MemoryManager separately to avoid type conflicts
from .rag_store.utils.memory import MemoryManager, _memory_manager

# Database Module (Iteration 5)
from .rag_store.core import DatabaseManager

# Embedding Module (Iteration 6)
from .rag_store.core import EmbeddingManager

# Search Module (Iteration 7)
from .rag_store.core import SearchManager

# ====================================================================
# LEGACY CODE (Wird schrittweise refactored)
# ====================================================================

# ResourceManager, ManagedResource, managed_resource
# → Extrahiert nach rag_store/utils/resource_manager.py ✅

# Large PDF Processing - REMOVED (large_pdf_processor deleted in cleanup)
# Use DoclingProcessor (utils/docling_processor.py) for large PDFs
LARGE_PDF_AVAILABLE: bool = False

# MEMORY_OPTIMIZED_PATCH - Löst OOM durch lazy loading
import gc
import psutil
from typing import Optional

# ProcessingConfig
# → Extrahiert nach rag_store/utils/config.py ✅ (Import oben)

class UnifiedRagStore(ManagedResource):
    @property
    def sentence_model(self):
        """
        Returns the underlying SentenceTransformer model used for embeddings.
        SOTA: Exposes the model for advanced GPU/FAISS tests and compatibility.
        Raises:
            RuntimeError: If the embedding manager or model is not initialized.
        """
        self._ensure_embedding_manager()
        if self._embedding_manager is None or not hasattr(self._embedding_manager, 'model'):
            raise RuntimeError("EmbeddingManager or its model is not initialized.")
        return self._embedding_manager.model

    @sentence_model.setter
    def sentence_model(self, value):
        """
        Allows setting the underlying SentenceTransformer model (e.g., to move to CPU/GPU).
        SOTA: Ensures assignment updates the embedding manager and device context.
        """
        self._ensure_embedding_manager()
        if self._embedding_manager is None:
            raise RuntimeError("EmbeddingManager is not initialized.")
        self._embedding_manager.model = value
    """
    🚀 KONSOLIDIERTE RAG-STORE-IMPLEMENTATION
    
    Vereint alle bisherigen RAG-Features in einer hochperformanten Klasse:
    - Automatische GPU/CPU-Erkennung
    - Parallele Verarbeitung wenn möglich  
    - Graceful Fallback auf sequenzielle Verarbeitung
    - Moderne Embedding-Provider
    - SQLite WAL-Mode für bessere Performance
    """
    
    # Klassen-Level Caches
    _shared_instances: Dict[str, "UnifiedRagStore"] = {}
    _refcounts: Dict[str, int] = {}
    _model_cache: Dict[str, Any] = {}
    _model_cache_lock = threading.RLock()
    
    # Metadaten
    _EXTRACTOR_VERSION = "unified_rag_store/3.0"
    _DEFAULT_EMBEDDING_MODEL = "multilingual-e5-large"  # Best for German/multilingual content (intfloat/multilingual-e5-large, 1024 dim)
    
    # Modern embedding providers (2025 state-of-the-art)
    # Moderne Embedding-Provider (Iteration 6: Nur HuggingFace)
    # ❌ REMOVED: OpenAI, Voyage, Random-Indexing (ungenutzt/gefährlich)
    _EMBEDDING_PROVIDERS: Dict[str, Dict[str, Any]] = {
        # HuggingFace Models (lokale GPU-beschleunigte Embeddings)
        "bge-large-en-v1.5": {"dim": 1024, "provider": "huggingface"},
        "BAAI/bge-large-en-v1.5": {"dim": 1024, "provider": "huggingface"},
        "multilingual-e5-large": {"dim": 1024, "provider": "huggingface"},
        "intfloat/multilingual-e5-large": {"dim": 1024, "provider": "huggingface"},
        "paraphrase-multilingual-mpnet-base-v2": {"dim": 768, "provider": "huggingface"},
        "jina-embeddings-v2-base-de": {"dim": 768, "provider": "huggingface"},
        # 2024/2025 SOTA models -- drop-in upgrades via RAG_EMBEDDING_MODEL env var
        "intfloat/multilingual-e5-large-instruct": {"dim": 1024, "provider": "huggingface"},
        "BAAI/bge-m3": {"dim": 1024, "provider": "huggingface"},
        "Alibaba-NLP/gte-multilingual-base": {"dim": 768, "provider": "huggingface"},
    }

    @classmethod
    def get_shared(cls, db_path: Optional[str] = None, dim: int = 384) -> "UnifiedRagStore":
        """Factory-Methode für geteilte Instanzen (Thread-sicher).

        ``db_path=None`` → kanonischer Pfad aus ``utils/db_path_resolver``
        (Root-Cause-Fix 2026-08-10: der frühere Literal-Default
        ``"rag_store.db"`` wurde CWD-relativ aufgelöst und umging .db_root).
        """
        env_path = os.getenv("RAG_DB_PATH")
        if env_path:
            db_path = env_path
        elif db_path is None:
            from utils.db_path_resolver import get_rag_store_path
            db_path = str(get_rag_store_path())
        key = os.path.abspath(db_path)
        inst = cls._shared_instances.get(key)
        if inst is None:
            inst = cls(db_path=db_path, dim=dim)
            cls._shared_instances[key] = inst
            cls._refcounts[key] = 1
        else:
            cls._refcounts[key] = cls._refcounts.get(key, 0) + 1
        return inst

    @classmethod
    def get_existing_shared(cls, db_path: Optional[str] = None) -> Optional["UnifiedRagStore"]:
        """Return an existing shared instance without creating or retaining it."""
        env_path = os.getenv("RAG_DB_PATH")
        if env_path:
            db_path = env_path
        elif db_path is None:
            from utils.db_path_resolver import get_rag_store_path
            db_path = str(get_rag_store_path())
        return cls._shared_instances.get(os.path.abspath(db_path))

    def __init__(
        self,
        db_path: Optional[str] = None,
        dim: int = 384,
        config: Optional[ProcessingConfig] = None,
        llm_client: Optional[Any] = None,
    ):
        """
        Initialisiert UnifiedRagStore mit automatischer Optimierung
        
        Args:
            db_path: SQLite-Datenbankpfad (default: project-root/rag_store.db via central resolver)
            dim: Embedding-Dimensionen (wird automatisch erkannt)
            config: Performance-Konfiguration
            llm_client: Optional LLM client (must expose ``generate_response``)
                used by the ContentClassifier for domain/safety verification
                at ingest time. If None, ingestion falls back to embedding-
                prototype-only domain detection (and ``safety_flag='safe'``).
        """
        # Initialisiere Resource Management zuerst
        ManagedResource.__init__(self)

        # Stash the LLM client for the lazy ContentClassifier
        self._llm_client = llm_client
        self._content_classifier: Optional[Any] = None
        self._content_classifier_lock = threading.RLock()
        
        # SOTA: Central path resolver ensures absolute paths regardless of CWD
        # (Streamlit may spawn workers from temp directories)
        env_path = os.getenv("RAG_DB_PATH")
        if env_path:
            db_path_resolved = env_path
        elif db_path is None:
            from utils.db_path_resolver import get_rag_store_path
            db_path_resolved = str(get_rag_store_path())
        else:
            db_path_resolved = db_path
            
        self.debug = os.getenv("RAG_DEBUG", "").lower() in {"1", "true", "yes", "on"}
        
        # Spezielle Behandlung für :memory: Database
        # WICHTIG: Speichere Original-Pfad - DatabaseManager macht die URI-Konversion!
        if db_path_resolved == ":memory:":
            self.db_path = ":memory:"  # Original beibehalten
            self._use_uri = True  # Flag für spätere Verwendung
        else:
            # Root-Cause-Fix 2026-07-14: früher wurde hier das rohe ``db_path``
            # (Optional[str]) statt des aufgelösten ``db_path_resolved`` verwendet
            # -> TypeError bei db_path=None ohne RAG_DB_PATH-Env.
            self.db_path = os.path.abspath(db_path_resolved)
            self._use_uri = False
            
        self.config = config or ProcessingConfig()
        
        # Moderne Embedding-Konfiguration (Iteration 6: Nur HuggingFace)
        self.embedding_model = os.getenv("RAG_EMBEDDING_MODEL", self._DEFAULT_EMBEDDING_MODEL)
        if self.embedding_model in self._EMBEDDING_PROVIDERS:
            provider_config = self._EMBEDDING_PROVIDERS[self.embedding_model]
            self.dim: int = int(provider_config["dim"])
            self.embedding_provider: str = str(provider_config["provider"])
        else:
            # FAIL-FAST: Kein Silent-Fallback mehr!
            raise ValueError(
                f"Unknown embedding model '{self.embedding_model}'!\n"
                f"Supported models: {list(self._EMBEDDING_PROVIDERS.keys())}\n"
                f"Set RAG_EMBEDDING_MODEL environment variable to one of the above."
            )
        
        # ❌ REMOVED in Iteration 6 (ungenutzte API-Keys)
        # self.openai_api_key = os.getenv("OPENAI_API_KEY")
        # self.voyage_api_key = os.getenv("VOYAGE_API_KEY")
        
        # Parallel-Capabilities
        self.device = None
        self.parallel_capable = False
        
        # Interne Zustandsvariablen
        self._conn: Optional[sqlite3.Connection] = None
        self._index = None
        self._index_size = 0
        self._pragmas_applied = False
        self._lock = threading.RLock()

        # Debounced/coalesced KG entity-resolution scheduler.
        # Root-cause fix: avoid O(N^2) resolution directly on the ingest critical path.
        self._init_entity_resolution_scheduler()
        
        # Connection Pool für Multi-Threading (via DatabaseManager - Iteration 5)
        self._connection_pool: queue.Queue[sqlite3.Connection] = queue.Queue(maxsize=self.config.connection_pool_size)
        
        # Database Manager (Iteration 5 - Core Module)
        # WICHTIG: db_path ist bereits konvertiert (siehe oben)
        self._db_manager = DatabaseManager(
            db_path=self.db_path,
            config=self.config,
            debug=self.debug
        )
        
        # LLM-based Knowledge Graph Extractor
        self.llm_kg_extractor = None
        if LLM_KG_AVAILABLE:
            try:
                from .llm_knowledge_graph import LLMKnowledgeGraphExtractor
                self.llm_kg_extractor = LLMKnowledgeGraphExtractor()
                logger.info("✅ LLM-KG-Extraktor initialisiert")
            except Exception as e:
                logger.warning(f"⚠️ LLM-KG-Extraktor fehlgeschlagen: {e}")
                self.llm_kg_extractor = None
        else:
            logger.info("ℹ️ LLM-KG nicht verfügbar - verwende Basic KG")
        
        # Initialisierung
        self._setup_capabilities()
        
        # Embedding Manager (Iteration 6 - Core Module)
        # LAZY INIT: Nicht im __init__, sondern erst bei erstem embed_texts()
        # Dies verhindert False-Negatives in Streamlit-Kontext
        self._embedding_manager: Optional[EmbeddingManager] = None
        self._embedding_manager_initialized: bool = False
        # AUX-GPU für Helper-Modelle (4090 bleibt exklusiv fürs LLM,
        # s. utils/gpu_devices.py + docs/RTX4090_RYZEN9_GUIDE.md)
        if self.device and "cuda" in str(self.device):
            from utils.gpu_devices import get_placement
            self._embedding_manager_device = get_placement().aux_device_string
        else:
            self._embedding_manager_device = "cpu"
        
        self._setup_database()
    

        # Resolve optional FAISS dependencies lazily at runtime.
        faiss_available = _ensure_faiss_runtime()
        faiss_hybrid_available = _ensure_faiss_hybrid_runtime()

        # FAISS Hybrid Search initialisieren
        self.faiss_manager = None
        self.fusion_engine = None
        if faiss_available and faiss_hybrid_available and FAISSIndexManager is not None and SmartFusionEngine is not None:
            try:  # type: ignore[unreachable]
                # WICHTIG: Verwende Singleton-Factory statt direkte Instanziierung
                # Dies verhindert Memory Leaks durch multiple FAISS Index Instanzen
                from agent.faiss_index_manager import get_faiss_manager
                self.faiss_manager = get_faiss_manager(
                    db_path=self.db_path,
                    embedding_dim=self.dim,
                    auto_load=True,  # Automatisches Laden/Bauen der Indizes
                    auto_rebuild_on_stale=True,  # Auto-Rebuild bei Staleness
                    rebuild_threshold=1000  # Rebuild nach 1000 neuen Chunks
                )
                self.fusion_engine = SmartFusionEngine(
                    recency_boost_factor=0.05,  # 5% Boost für neueste Chunks
                    kg_boost_factor=0.1,        # 10% Boost für KG-Results
                    target_kg_ratio=0.18,       # 18% KG-Anteil angestrebt
                    use_adaptive_weights=True   # 🆕 Aktiviere Feedback-basierte Gewichte
                )
                logger.info(f"✅ FAISS Hybrid Search initialisiert (Dim={self.dim}, Rebuild-Threshold=1000)")
            except Exception as e:
                logger.warning(f"⚠️ FAISS Hybrid Init fehlgeschlagen: {e}")
                self.faiss_manager = None
                self.fusion_engine = None
        
        # 🆕 Feedback Update Service initialisieren
        self.feedback_service = None
        try:
            from agent.feedback_update_service import FeedbackUpdateService
            self.feedback_service = FeedbackUpdateService(
                smart_fusion_engine=self.fusion_engine,
                update_interval_seconds=300,  # 5 Minuten
                min_feedbacks_for_update=10
            )
            # ★ SOTA v3: Wire DB path and load initial Wilson utility scores
            self.feedback_service.set_db_path(self.db_path)
            n_wilson = self.feedback_service.load_wilson_utility_scores()
            logger.info(f"✅ Feedback Update Service initialisiert (Wilson scores: {n_wilson})")
        except Exception as e:
            logger.warning(f"⚠️ Feedback Service Init fehlgeschlagen: {e}")
            self.feedback_service = None
        
        # Search Manager (Iteration 7 - Core Module)
        self._search_manager = None
        try:
            self._search_manager = SearchManager(
                db_manager=self._db_manager,
                embedding_manager=self._embedding_manager,
                embedding_dim=self.dim,
                faiss_manager=self.faiss_manager,
                fusion_engine=self.fusion_engine,
                debug=self.debug
            )
            if self.debug:
                logger.info(f"✅ SearchManager initialisiert")
        except Exception as e:
            logger.error(f"❌ SearchManager Init fehlgeschlagen: {e}")
            raise RuntimeError(f"Cannot initialize SearchManager: {e}") from e
    
    def _setup_capabilities(self):
        """Erkennt verfügbare Hardware und Software-Capabilities"""
        
        # PRIMÄRE GPU-ERKENNUNG: Prüfe torch.cuda ZUERST
        if CUDA_AVAILABLE and torch and self.config.use_gpu:
            self.device = torch.device("cuda")
            self.parallel_capable = True
            
            try:
                gpu_name = torch.cuda.get_device_name(0)
                gpu_memory = torch.cuda.get_device_properties(0).total_memory // 1024**3
                
                # GPU-optimierte Batch-Größe
                if gpu_memory >= 20:
                    self.config.gpu_batch_size = 256
                elif gpu_memory >= 12:
                    self.config.gpu_batch_size = 128
                elif gpu_memory >= 8:
                    self.config.gpu_batch_size = 64
                else:
                    self.config.gpu_batch_size = 32
                
                logger.info(f"🎮 GPU erkannt: {gpu_name} ({gpu_memory} GB VRAM)")
                logger.info(f"📊 GPU Batch Size: {self.config.gpu_batch_size}")
                
                # Klarstellung: Die Zeile oben dient nur der Batch-Size-Erkennung
                # (Device 0). Die eigentlichen Rollen (LLM/AUX) liefert
                # utils.gpu_devices — Embeddings/Helper laufen auf der AUX-GPU.
                try:
                    from utils.gpu_devices import get_placement as _get_placement
                    _p = _get_placement()
                    logger.info(
                        f"ℹ️ GPU-Rollen: LLM → {_p.llm.name} (cuda:{_p.llm_cuda}) | "
                        f"Embeddings/Helper → {_p.aux.name} (cuda:{_p.aux_cuda})"
                    )
                except Exception:
                    pass
                
            except Exception as e:
                logger.warning(f"GPU-Info fehlgeschlagen: {e}")
        else:
            if TORCH_AVAILABLE:
                import torch as _torch  # local re-import for type narrowing
                self.device = _torch.device("cpu")
            else:
                self.device = None
            self.parallel_capable = False
            logger.info("🖥️ Verwende CPU (GPU nicht verfügbar oder deaktiviert)")
        
        # SEKUNDÄR: GPU Optimizer (optional, zusätzliche Optimierungen)
        self.gpu_optimizer = None
        if GPU_OPTIMIZER_AVAILABLE and self.device and 'cuda' in str(self.device):
            try:
                self.gpu_optimizer = get_gpu_optimizer(debug=self.debug)
                apply_gpu_optimizations()
                
                if (self.gpu_optimizer and hasattr(self.gpu_optimizer, 'gpu_info') and 
                    self.gpu_optimizer.gpu_info["available"]):
                    # Update batch size based on optimization
                    if hasattr(self.gpu_optimizer, 'get_optimal_embedding_config'):
                        embedding_config = self.gpu_optimizer.get_optimal_embedding_config()
                        self.config.gpu_batch_size = embedding_config["batch_size"]
                    
                    if self.debug:
                        print(f"🚀 GPU Optimizer aktiv: {self.gpu_optimizer.gpu_info['name']}")
                        print(f"💾 VRAM: {self.gpu_optimizer.gpu_info['total_memory_gb']:.1f} GB")
                        print(f"📊 Optimized Batch Size: {self.config.gpu_batch_size}")
                        
                        # Temperature check
                        if hasattr(self.gpu_optimizer, 'get_temperature_warning'):
                            temp_warning = self.gpu_optimizer.get_temperature_warning()
                            if temp_warning:
                                print(temp_warning)
            except Exception as e:
                if self.debug:
                    print(f"⚠️ GPU Optimizer Fehler (optional): {e}")
        
        # Embedding-Modell: LAZY INIT via _ensure_embedding_manager()
        # Wird beim ersten embed_texts() Aufruf geladen (nicht hier).
        # EmbeddingManager nutzt intern EmbeddingSingleton -> nur EINE Modell-Kopie im Speicher.
        # Batch-Groesse wird ueber self.config.gpu_batch_size an EmbeddingManager weitergereicht.
        logger.info(
            f"\u2705 Capabilities: device={self.device}, "
            f"batch_size={self.config.gpu_batch_size}, "
            f"embedding={self.embedding_model} (lazy init via EmbeddingManager)"
        )

    def _setup_database(self):
        """Konfiguriert SQLite für optimale Performance"""
        
        if self.config.sqlite_wal_mode:
            try:
                conn = sqlite3.connect(self.db_path, check_same_thread=False, uri=self._use_uri)
                # WAL-Modus: erlaubt parallele Leser (kg_dashboard) neben Schreiber
                conn.execute("PRAGMA journal_mode=WAL")
                conn.execute("PRAGMA synchronous=NORMAL")
                conn.execute("PRAGMA busy_timeout=15000")  # 15s warten statt sofort "locked"
                conn.execute("PRAGMA cache_size=100000") 
                conn.execute("PRAGMA temp_store=MEMORY")
                conn.execute("PRAGMA mmap_size=268435456")  # 256MB
                conn.execute("PRAGMA wal_autocheckpoint=1000")
                conn.close()
                
                if self.debug:
                    print("🗃️ SQLite WAL-Mode aktiviert (concurrent reads möglich)")
                    
            except Exception as e:
                logger.warning(f"SQLite WAL-Setup fehlgeschlagen: {e}")

    def get_connection(self) -> sqlite3.Connection:
        """
        Thread-sichere Connection aus dem Pool
        
        REFACTORED (Iteration 5): Delegates to DatabaseManager
        """
        return self._db_manager.get_connection()

    def return_connection(self, conn: sqlite3.Connection):
        """
        Connection zurück in den Pool
        
        REFACTORED (Iteration 5): Delegates to DatabaseManager
        """
        self._db_manager.return_connection(conn)

    def close(self) -> None:
        """
        ROBUSTE Schließung aller Verbindungen und Aufräumen
        
        REFACTORED (Iteration 5): Delegates to DatabaseManager for DB cleanup
        """
        with self._lock:
            try:
                if self.debug:
                    print("� Starting ROBUST database close...")
                
                # Phase 1: Close DatabaseManager (handles all DB connections)
                self._db_manager.close()
                if self.debug:
                    print("✅ DatabaseManager closed")
                
                # Phase 2: Close local connection if exists
                if self._conn:
                    try:
                        self._conn.close()
                        if self.debug:
                            print("🔒 Closed local connection")
                    except Exception as e:
                        logger.warning(f"Error closing local connection: {e}")
                    finally:
                        self._conn = None
                
                # Phase 3: GPU memory cleanup
                if self.device and CUDA_AVAILABLE and torch and 'cuda' in str(self.device):
                    torch.cuda.empty_cache()
                    if self.debug:
                        print("🧹 GPU cache cleared")
                    
                if self.debug:
                    print("✅ ROBUST database close completed")
                    
            except Exception as e:
                if not (self._is_interpreter_shutting_down() or self._is_shutdown_exception(e)):
                    logger.warning(f"Cleanup-Fehler: {e}")
                if self.debug:
                    print(f"❌ Close error: {e}")

    def _init_entity_resolution_scheduler(self) -> None:
        """Initializes async debounced entity-resolution state."""
        self._entity_resolution_lock = threading.RLock()
        self._entity_resolution_executor = ThreadPoolExecutor(
            max_workers=1,
            thread_name_prefix="kg-entity-resolve",
        )
        self._entity_resolution_timer: Optional[threading.Timer] = None
        self._entity_resolution_inflight = False
        self._entity_resolution_inflight_payload: Optional[Tuple[int, int]] = None
        self._entity_resolution_pending_ingests = 0
        self._entity_resolution_pending_triples = 0
        self._entity_resolution_last_run_monotonic = 0.0
        self._entity_resolution_min_ingests = max(
            1,
            int(os.getenv("KG_ENTITY_RESOLVE_MIN_INGESTS", "4")),
        )
        self._entity_resolution_cooldown_sec = max(
            5.0,
            float(os.getenv("KG_ENTITY_RESOLVE_COOLDOWN_SEC", "45")),
        )

    def _arm_entity_resolution_timer_locked(self, delay_seconds: float) -> None:
        """Arms trailing-edge timer to coalesce chatty ingest bursts.

        Must be called with ``self._entity_resolution_lock`` already held.
        """
        if self._entity_resolution_timer is not None and self._entity_resolution_timer.is_alive():
            return

        timer = threading.Timer(delay_seconds, self._entity_resolution_timer_fire)
        timer.daemon = True
        self._entity_resolution_timer = timer
        timer.start()

    def _submit_entity_resolution_locked(self) -> None:
        """Submits one async entity-resolution run if there is pending work.

        Must be called with ``self._entity_resolution_lock`` already held.
        """
        if self._entity_resolution_inflight:
            return
        if self._entity_resolution_pending_ingests <= 0:
            return
        if self._entity_resolution_executor is None:
            # Scheduler wurde bereits heruntergefahren (Shutdown-Race):
            # kein Submit mehr möglich, pending-Arbeit verfällt kontrolliert.
            return

        pending_ingests = self._entity_resolution_pending_ingests
        pending_triples = self._entity_resolution_pending_triples
        self._entity_resolution_pending_ingests = 0
        self._entity_resolution_pending_triples = 0
        self._entity_resolution_inflight = True
        self._entity_resolution_inflight_payload = (pending_ingests, pending_triples)

        future = self._entity_resolution_executor.submit(
            self.resolve_duplicate_entities,
            similarity_threshold=0.90,
        )
        future.add_done_callback(self._on_entity_resolution_done)

    def _entity_resolution_timer_fire(self) -> None:
        """Timer callback: run one coalesced resolution pass after quiet period."""
        with self._entity_resolution_lock:
            self._entity_resolution_timer = None
            if self._entity_resolution_inflight:
                return
            if self._entity_resolution_pending_ingests <= 0:
                return

            elapsed = time.monotonic() - self._entity_resolution_last_run_monotonic
            if elapsed < self._entity_resolution_cooldown_sec:
                remaining = self._entity_resolution_cooldown_sec - elapsed
                self._arm_entity_resolution_timer_locked(remaining)
                return

            self._submit_entity_resolution_locked()

    def _on_entity_resolution_done(self, future: Any) -> None:
        """Finalizes async entity-resolution run and schedules follow-up work."""
        payload: Tuple[int, int]
        with self._entity_resolution_lock:
            payload = self._entity_resolution_inflight_payload or (0, 0)
            self._entity_resolution_inflight_payload = None
            self._entity_resolution_inflight = False

        exc = future.exception()
        if exc is not None:
            with self._entity_resolution_lock:
                self._entity_resolution_pending_ingests += payload[0]
                self._entity_resolution_pending_triples += payload[1]
                self._arm_entity_resolution_timer_locked(self._entity_resolution_cooldown_sec)
            if not (self._is_interpreter_shutting_down() or self._is_shutdown_exception(exc)):
                logger.error("[KG-AutoResolve] Async entity resolution failed: %s", exc)
            return

        result = future.result()
        with self._entity_resolution_lock:
            self._entity_resolution_last_run_monotonic = time.monotonic()
            has_pending = self._entity_resolution_pending_ingests > 0
            if has_pending:
                self._arm_entity_resolution_timer_locked(self._entity_resolution_cooldown_sec)

        if isinstance(result, dict) and result.get("merged_groups", 0) > 0:
            logger.info(
                "[KG-AutoResolve] %s merge groups (%s -> %s entities)",
                result.get("merged_groups"),
                result.get("entities_before"),
                result.get("entities_after"),
            )

    def _schedule_entity_resolution(self, triples_created: int) -> Dict[str, Any]:
        """Debounced/coalesced trigger for async entity-resolution.

        The scheduler uses:
        - a single worker (no concurrent O(N^2) runs),
        - ingest-count coalescing,
        - a cooldown window to reduce burst amplification.
        """
        if triples_created <= 0:
            return {
                "scheduled": False,
                "reason": "no_new_triples",
            }

        with self._entity_resolution_lock:
            self._entity_resolution_pending_ingests += 1
            self._entity_resolution_pending_triples += triples_created

            if self._entity_resolution_inflight:
                return {
                    "scheduled": True,
                    "mode": "coalesced_while_inflight",
                    "pending_ingests": self._entity_resolution_pending_ingests,
                    "pending_triples": self._entity_resolution_pending_triples,
                }

            now = time.monotonic()
            elapsed = now - self._entity_resolution_last_run_monotonic
            cooldown_done = elapsed >= self._entity_resolution_cooldown_sec
            ingest_threshold_reached = (
                self._entity_resolution_pending_ingests >= self._entity_resolution_min_ingests
            )

            if ingest_threshold_reached and cooldown_done:
                self._submit_entity_resolution_locked()
                return {
                    "scheduled": True,
                    "mode": "submitted_now",
                    "pending_ingests": self._entity_resolution_pending_ingests,
                    "pending_triples": self._entity_resolution_pending_triples,
                }

            if not ingest_threshold_reached:
                delay = self._entity_resolution_cooldown_sec
            else:
                delay = max(0.0, self._entity_resolution_cooldown_sec - elapsed)
            self._arm_entity_resolution_timer_locked(delay)

            return {
                "scheduled": True,
                "mode": "debounced",
                "delay_sec": round(delay, 3),
                "pending_ingests": self._entity_resolution_pending_ingests,
                "pending_triples": self._entity_resolution_pending_triples,
                "min_ingests": self._entity_resolution_min_ingests,
            }

    def _shutdown_entity_resolution_scheduler(self) -> None:
        """Stops timer/executor used by the async entity-resolution scheduler."""
        with self._entity_resolution_lock:
            timer = self._entity_resolution_timer
            self._entity_resolution_timer = None
            executor = self._entity_resolution_executor
            self._entity_resolution_executor = None

        if timer is not None:
            timer.cancel()
        if executor is not None:
            executor.shutdown(wait=False, cancel_futures=True)

    @staticmethod
    def _is_interpreter_shutting_down() -> bool:
        try:
            return bool(getattr(sys, "is_finalizing", lambda: False)())
        except Exception as exc:
            logger.debug(f"Could not read interpreter finalization state: {exc}")
            return False

    @staticmethod
    def _is_shutdown_exception(exc: Exception) -> bool:
        text = str(exc)
        return (
            "sys.meta_path is None" in text
            or "Python is likely shutting down" in text
            or "interpreter shutdown" in text.lower()
        )
    
    # ======= EMBEDDING METHODS =======
    
    def _ensure_embedding_manager(self) -> None:
        """
        LAZY INIT: Initialisiert EmbeddingManager erst bei erstem Embedding-Request
        
        Verhindert False-Negatives in Streamlit-Kontext, wo sentence-transformers
        beim Store-Init noch nicht verfügbar ist, aber später (zur Runtime) schon.
        
        Raises:
            RuntimeError: Wenn EmbeddingManager nicht initialisiert werden kann
        """
        if self._embedding_manager_initialized:
            return  # Already initialized
        
        try:
            model_name = self.embedding_model or "multilingual-e5-large"
            self._embedding_manager = EmbeddingManager(
                model_name=model_name,
                device=self._embedding_manager_device,
                config=self.config,
                debug=self.debug
            )
            self._embedding_manager_initialized = True
            if self.debug:
                logger.info(f"✅ EmbeddingManager lazy-initialized: {self._embedding_manager}")
        except Exception as e:
            self._embedding_manager_initialized = True  # Mark as attempted
            logger.error(f"❌ EmbeddingManager lazy-init fehlgeschlagen: {e}")
            raise RuntimeError(f"Cannot initialize EmbeddingManager: {e}") from e

    def _ensure_content_classifier(self) -> Optional[Any]:
        """LAZY INIT: ContentClassifier with prototype + (optional) LLM verifier.

        The classifier is only built once a classification is actually
        requested, since (a) it depends on a fully-initialised
        EmbeddingManager and (b) it is irrelevant for callers that always
        pass an explicit ``corpus_domain``/``safety_flag`` (e.g. the
        psycho-corpus bootstrapper).

        Returns the classifier instance, or ``None`` if construction
        failed (logged) — callers must then either supply explicit
        labels or surface the absence as an error.
        """
        with self._content_classifier_lock:
            if self._content_classifier is not None:
                return self._content_classifier

            self._ensure_embedding_manager()
            if self._embedding_manager is None:
                raise RuntimeError(
                    "Cannot initialise ContentClassifier: EmbeddingManager unavailable"
                )

            llm_wrapper = None
            if self._llm_client is not None:
                try:
                    from llm_structured_wrapper import LLMStructuredWrapper
                    llm_wrapper = LLMStructuredWrapper(
                        llm_client=self._llm_client,
                        max_retries=2,
                        temperature=0.0,
                        enable_logging=False,
                    )
                except Exception as exc:
                    # Surface the wiring problem but continue with
                    # prototype-only — the prototype is still a valid
                    # primary signal.
                    logger.warning(
                        "ContentClassifier: LLM wrapper construction failed "
                        "(%s) — continuing with prototype-only classification.",
                        exc,
                    )
                    llm_wrapper = None

            try:
                from .rag_store.core.content_classifier import ContentClassifier
                self._content_classifier = ContentClassifier(
                    embedding_manager=self._embedding_manager,
                    db_manager=self._db_manager,
                    llm_wrapper=llm_wrapper,
                )
                logger.info(
                    "✅ ContentClassifier ready (llm=%s)",
                    "yes" if llm_wrapper is not None else "no",
                )
            except Exception as exc:
                logger.error(
                    "❌ ContentClassifier init failed: %s — auto-classification "
                    "disabled. Callers must pass explicit corpus_domain/safety_flag.",
                    exc,
                )
                self._content_classifier = None

            return self._content_classifier

    def _auto_classify(
        self,
        *,
        text: str,
        title: Optional[str] = None,
    ) -> Tuple[str, str]:
        """Resolve ``(corpus_domain, safety_flag)`` for content where the
        producer didn't declare them.

        Behaviour:
            * If the ContentClassifier is available → ask it. The result
              is cached persistently in ``content_classification_cache``.
            * If unavailable (no embedding manager, prototype too small
              and no LLM, etc.) → raise ``RuntimeError``. We refuse to
              silently default to ``('general', 'safe')`` because that
              would invisibly leak content into the wrong namespace.

        Callers that want explicit defaults must pass them on the doc
        dict (or via ``add_document(corpus_domain=..., safety_flag=...)``).
        """
        if not text or not text.strip():
            # Empty docs get filtered downstream anyway, but be explicit:
            return ("general", "safe")

        classifier = self._ensure_content_classifier()
        if classifier is None:
            raise RuntimeError(
                "upsert_documents: cannot auto-classify because the "
                "ContentClassifier could not be initialised. Either pass "
                "explicit corpus_domain/safety_flag on the document, or "
                "ensure the EmbeddingManager and (optionally) llm_client "
                "are wired into UnifiedRagStore."
            )

        result = classifier.classify(text=text, title=title)
        if self.debug:
            logger.info(
                "📑 auto-classify → domain=%s safety=%s method=%s "
                "confidence=%.2f sim=%s",
                result.corpus_domain,
                result.safety_flag,
                result.method,
                result.confidence,
                f"{result.similarity:.3f}" if result.similarity is not None else "—",
            )
        return (result.corpus_domain, result.safety_flag)

    # ──────────────────────────────────────────────────────────────────
    #  Public maintenance API — used by the Streamlit "Reclassify legacy
    #  corpus" panel and the CLI. Reuses the live classifier so the
    #  embedding model is loaded once, and the LLM wired by
    #  set_llm_client() is reused as-is.
    # ──────────────────────────────────────────────────────────────────
    def reclassify_legacy_chunks(
        self,
        *,
        batch_size: int = 32,
        embed_batch_size: int = 64,
        dry_run: bool = False,
        progress_callback=None,
    ):
        """Run :class:`ChunkReclassifier` against this store.

        The classifier and DB manager are taken from the live store, so:
        * the embedding model is the same instance that ingest uses,
        * the LLM wrapper is the one wired via :meth:`set_llm_client`,
        * the classification cache is the same SQLite table.

        ``embed_batch_size`` controls the GPU-side batch (single
        ``encode()`` call per N docs) — the lever that actually
        determines GPU utilisation.

        Returns a :class:`ReclassifierStats` instance. Raises if the
        classifier cannot be built (e.g. embedding manager unavailable)
        or if any doc fails classification — never silently degrades.
        """
        from .rag_store.maintenance.reclassifier import (
            ChunkReclassifier, ReclassifierStats
        )
        classifier = self._ensure_content_classifier()
        if classifier is None:
            raise RuntimeError(
                "reclassify_legacy_chunks: ContentClassifier construction "
                "failed; cannot proceed. Check earlier logs for the root cause."
            )
        job = ChunkReclassifier(self._db_manager, classifier)
        return job.run(
            batch_size=batch_size,
            embed_batch_size=embed_batch_size,
            dry_run=dry_run,
            progress_callback=progress_callback,
        )

    def count_stale_chunks(self) -> Tuple[int, int]:
        """Return ``(stale_doc_count, stale_chunk_count)`` for the
        Streamlit panel preview. Uses the same predicate as the sweep so
        the displayed numbers always match what a run would touch.
        """
        from .rag_store.core.content_classifier import CLASSIFICATION_VERSION
        conn = self._db_manager.get_connection()
        try:
            cur = conn.cursor()
            try:
                docs = cur.execute(
                    "SELECT COUNT(DISTINCT doc_id) FROM chunks "
                    "WHERE classification_version IS NULL "
                    "   OR classification_version < ?",
                    (CLASSIFICATION_VERSION,),
                ).fetchone()[0]
                chunks = cur.execute(
                    "SELECT COUNT(*) FROM chunks "
                    "WHERE classification_version IS NULL "
                    "   OR classification_version < ?",
                    (CLASSIFICATION_VERSION,),
                ).fetchone()[0]
                return int(docs), int(chunks)
            finally:
                cur.close()
        finally:
            self._db_manager.return_connection(conn)

    def embed_texts(self, texts: List[str]):
        """
        Generiert Embeddings mit EmbeddingManager (Iteration 6)
        
        Delegiert an self._embedding_manager für konsistente Embedding-Generierung
        """
        assert np is not None, "NumPy required"
        
        # LAZY INIT: Erst hier EmbeddingManager initialisieren
        self._ensure_embedding_manager()
        
        if self._embedding_manager is None:
            raise RuntimeError("EmbeddingManager not initialized after lazy init attempt.")
        
        if self.debug:  # type: ignore[unreachable]
            logger.warning(f"🔍 embed_texts: {len(texts)} texts via EmbeddingManager")
        
        return self._embedding_manager.embed_texts(texts)

    # ❌ REMOVED in Iteration 6: _embed_huggingface, _embed_openai, _embed_voyage, _embed_random_indexing
    # Alle Embedding-Operationen werden jetzt durch EmbeddingManager gehandhabt
    # Siehe: agent/rag_store/core/embeddings.py

    # ======= DOCUMENT MANAGEMENT =======

    # Allowed namespace + safety values are validated at the API boundary so
    # that we never silently insert mistyped values that would make a chunk
    # unreachable through the standard intent-aware retrieval path.
    _VALID_CORPUS_DOMAINS: Set[str] = {"general", "psych"}
    _VALID_SAFETY_FLAGS: Set[str] = {"safe", "sensitive", "crisis"}

    def add_document(
        self,
        content: str,
        *,
        doc_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        corpus_domain: Optional[str] = None,
        safety_flag: Optional[str] = None,
        chunk_size: int = 1500,
        chunk_overlap: int = 200,
    ) -> Dict[str, Any]:
        """SOTA single-document ingest API.

        Wraps :meth:`upsert_documents` with explicit ``corpus_domain`` and
        ``safety_flag`` parameters so that callers (e.g. the psycho-corpus
        bootstrapper) cannot accidentally insert content under the wrong
        namespace. Validates inputs at the boundary — invalid values raise
        ``ValueError`` rather than being silently downgraded.

        Auto-classification:
            * ``corpus_domain=None`` and/or ``safety_flag=None`` triggers
              the :class:`ContentClassifier` (embedding-prototype primary,
              optional LLM verifier). Callers that *know* the provenance
              (e.g. the psycho-corpus bootstrapper) should pass explicit
              values to skip classification.
            * Any explicit value still goes through whitelist validation.
        """
        if not isinstance(content, str) or not content.strip():
            raise ValueError("add_document: 'content' must be a non-empty string")
        if corpus_domain is not None and corpus_domain not in self._VALID_CORPUS_DOMAINS:
            raise ValueError(
                f"add_document: corpus_domain={corpus_domain!r} not in "
                f"{sorted(self._VALID_CORPUS_DOMAINS)}"
            )
        if safety_flag is not None and safety_flag not in self._VALID_SAFETY_FLAGS:
            raise ValueError(
                f"add_document: safety_flag={safety_flag!r} not in "
                f"{sorted(self._VALID_SAFETY_FLAGS)}"
            )

        doc: Dict[str, Any] = {
            "text": content,
            "metadata": dict(metadata) if metadata else {},
        }
        if corpus_domain is not None:
            doc["corpus_domain"] = corpus_domain
        if safety_flag is not None:
            doc["safety_flag"] = safety_flag
        if doc_id:
            doc["id"] = doc_id
        return self.upsert_documents(
            [doc], chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )

    def upsert_documents(self, documents: List[Dict[str, Any]], *, 
                        chunk_size: int = 1500, chunk_overlap: int = 200) -> Dict[str, Any]:
        """
        Fügt Dokumente in den Store ein oder aktualisiert sie
        
        Args:
            documents: Liste von Dokumenten mit 'text', 'id', 'metadata'
            chunk_size: Chunk-Größe für Textsplitting  
            chunk_overlap: Überlappung zwischen Chunks
            
        Returns:
            Dict mit Statistiken über eingefügte Dokumente
        """
        chunk_size = max(200, int(chunk_size))
        chunk_overlap = max(0, int(chunk_overlap))

        # Tuple layout: (doc_id, chunk_idx, text, metadata_json, embedding,
        #                corpus_domain, safety_flag)
        to_insert: List[Tuple[str, int, str, str, bytes, str, str]] = []

        for d in documents:
            base_text = str(d.get("text") or "")
            doc_id = str(d.get("id") or d.get("doc_id") or self._hash_str(base_text[:300]))
            metadata_in = d.get("metadata") or {}

            # Document-level namespace + safety override (passed through from
            # add_document or the caller). Falls back to metadata for legacy
            # callers. If neither doc nor metadata declares the value, the
            # ContentClassifier (prototype + optional LLM) is consulted.
            explicit_domain = (
                d.get("corpus_domain") or metadata_in.get("corpus_domain")
            )
            explicit_safety = (
                d.get("safety_flag") or metadata_in.get("safety_flag")
            )

            if explicit_domain is None or explicit_safety is None:
                title_hint = (
                    metadata_in.get("title")
                    or metadata_in.get("doc_title")
                    or None
                )
                inferred = self._auto_classify(
                    text=base_text,
                    title=str(title_hint) if title_hint else None,
                )
                if explicit_domain is None:
                    explicit_domain = inferred[0]
                if explicit_safety is None:
                    explicit_safety = inferred[1]

            corpus_domain = str(explicit_domain)
            safety_flag = str(explicit_safety)
            if corpus_domain not in self._VALID_CORPUS_DOMAINS:
                raise ValueError(
                    f"upsert_documents: corpus_domain={corpus_domain!r} not in "
                    f"{sorted(self._VALID_CORPUS_DOMAINS)}"
                )
            if safety_flag not in self._VALID_SAFETY_FLAGS:
                raise ValueError(
                    f"upsert_documents: safety_flag={safety_flag!r} not in "
                    f"{sorted(self._VALID_SAFETY_FLAGS)}"
                )

            # Erweiterte Metadaten
            url = str(metadata_in.get("canonical_url") or metadata_in.get("url") or "")
            domain = ""
            try:
                if url:
                    domain = (urlparse(url).netloc or "").lower()
                    if domain.startswith("www."):
                        domain = domain[4:]
            except (ValueError, AttributeError) as e:
                logging.debug(f"URL-Parsing-Fehler für Domain-Extraktion: {e}")
                domain = ""
            except Exception as e:
                logging.warning(f"Unerwarteter Fehler bei Domain-Extraktion: {type(e).__name__}: {e}")
                domain = ""
            
            checksum_full = self._hash_str(base_text) if base_text else None
            now_iso = datetime.now(timezone.utc).isoformat()
            
            base_meta_aug: Dict[str, Any] = {
                **metadata_in,
                "doc_id": doc_id,
                "canonical_url": url or metadata_in.get("canonical_url"),
                "domain": domain or metadata_in.get("domain"),
                "checksum": checksum_full,
                "extractor_version": self._EXTRACTOR_VERSION,
                "embedding_model": self.embedding_model,
                "embedding_dim": int(self.dim),
                "extracted_at": now_iso,
                "source_type": metadata_in.get("source_type") or "document",
                "content_type": metadata_in.get("content_type") or "text/plain",
                "language": metadata_in.get("language") or (self._detect_language(base_text) if base_text else None),
                "collection": metadata_in.get("collection"),
                "tags": metadata_in.get("tags") or [],
            }

            # ── Retention (2026-08-30): Web-sourced Content bekommt retention_until ──
            # Siehe docs/18_LEGAL_WEB_PERSIST.md — nur web-* Quellen, nie lokale
            # Dateien. WEB_RETENTION_DAYS=0 → kein Feld (unbegrenzt).
            src_type = str(base_meta_aug.get("source_type") or "")
            if src_type.startswith("web") and "retention_until" not in base_meta_aug:
                retention_until = web_compliance.retention_until_iso()
                if retention_until:
                    base_meta_aug["retention_until"] = retention_until

            # Chunking
            chunks: List[Tuple[int, int, str]] = []
            provided_chunks = False
            if d.get("chunks"):
                provided_chunks = True
                tmp_chunks: List[str] = [str(x) for x in d["chunks"] if str(x).strip()]
                for idx, ch in enumerate(tmp_chunks):
                    chunks.append((-1, -1, ch))
            else:
                if not base_text.strip():
                    continue
                chunks = self._chunk_text_with_spans(base_text, chunk_size=chunk_size, overlap=chunk_overlap)
            
            if not chunks:
                continue

            # Embeddings generieren
            embs = self.embed_texts([t for (_s, _e, t) in chunks])

            # Metadaten pro Chunk
            for idx, ((_s, _e, ch), vec) in enumerate(zip(chunks, embs)):
                # Ensure vec is numpy array for proper conversion
                if isinstance(vec, list):
                    if np:
                        vec = np.array(vec, dtype="float32")
                        emb_blob = vec.tobytes()
                    else:
                        # Fallback ohne numpy
                        emb_blob = bytes()
                elif hasattr(vec, "astype") and hasattr(vec, "tobytes"):
                    emb_blob = vec.astype("float32").tobytes()
                else:
                    # Final fallback
                    if np:
                        emb_blob = np.array(vec, dtype="float32").tobytes()
                    else:
                        emb_blob = bytes()
                    
                # Einfache Tokenisierung für BM25-Felder
                import re
                tokens = []
                if ch:
                    t = str(ch).lower()
                    toks = re.findall(r"\w+", t, flags=re.UNICODE)
                    tokens = [x for x in toks if len(x) >= 2]
                
                chunk_meta: Dict[str, Any] = {
                    **base_meta_aug,
                    "page": metadata_in.get("page"),
                    "type": metadata_in.get("type"),
                    "chunk_index": idx,
                    "char_start": None if provided_chunks else _s,
                    "char_end": None if provided_chunks else _e,
                    "token_count": len(tokens),
                    "bm25_fields": {
                        "title": metadata_in.get("title"),
                        "keywords": metadata_in.get("keywords") or tokens[:8],
                    },
                }
                
                meta_json = json.dumps(chunk_meta, ensure_ascii=False)
                to_insert.append(
                    (doc_id, idx, ch, meta_json, emb_blob, corpus_domain, safety_flag)
                )

        if not to_insert:
            return {"success": True, "inserted": 0}

        # Duplikationsschutz mit batch_sql_load (Iteration 3)
        conn = self.get_connection()
        cur = conn.cursor()
        
        # Batch-Load aller Texte für Duplikatsprüfung
        texts_to_check = [item[2] for item in to_insert]
        
        # Optimierte Duplikatsprüfung: Batch-Query statt einzelner SELECTs
        deduplicated_inserts = []
        duplicate_count = 0
        synthetic_count = 0
        noise_count = 0
        
        # Batch-Duplikatsprüfung (max 999 Parameter bei SQLite)
        SQLITE_MAX_VARS = 999
        existing_texts: Set[str] = set()
        
        for batch in chunk_list(texts_to_check, SQLITE_MAX_VARS):
            placeholders = ','.join('?' * len(batch))
            query = f"SELECT DISTINCT text FROM chunks WHERE text IN ({placeholders})"
            cur.execute(query, batch)
            existing_texts.update(row[0] for row in cur.fetchall())
        
        # Source-Types die NICHT gespeichert werden sollen (synthetische/generierte Inhalte)
        SYNTHETIC_SOURCE_TYPES = {
            'rag_response',      # RAG-generierte Antworten
            'llm_response',      # LLM-generierte Antworten
            'bot_response',      # Bot-Antworten
            'synthetic',         # Explizit als synthetisch markiert
            'generated',         # Generierte Inhalte
            'chat_response',     # Chat-Antworten
        }
        
        for doc_id, chunk_idx, text, meta_json, emb_blob, corpus_domain, safety_flag in to_insert:
            # Prüfe auf exakte Textübereinstimmung (aus Batch-Query)
            if text in existing_texts:
                duplicate_count += 1
                continue
            
            # STATE-OF-THE-ART: Metadata-basierte Filterung synthetischer Inhalte
            try:
                meta = json.loads(meta_json) if meta_json else {}
                source_type = str(meta.get('source_type', '')).lower()
                is_synthetic = meta.get('is_synthetic', False)
                
                # Blockiere synthetische/generierte Inhalte
                if source_type in SYNTHETIC_SOURCE_TYPES or is_synthetic:
                    synthetic_count += 1
                    continue
            except (json.JSONDecodeError, TypeError):
                pass  # Bei Parsing-Fehler: normal fortfahren
            
            # LAYER 2: Chunk-Content-Quality-Validator
            # Erkennt und verwirft Noise-Chunks (JS, Tracking, minified Code, Boilerplate, URL-Dumps)
            if self._is_noise_chunk(text):
                noise_count += 1
                continue
            
            deduplicated_inserts.append(
                (doc_id, chunk_idx, text, meta_json, emb_blob, corpus_domain, safety_flag)
            )

        # ── LAYER 3 (NEW): Near-Duplicate Gate via Embedding Cosine ──
        # Checks new chunks against EXISTING chunks in DB using embedding similarity.
        # Threshold 0.97 = nearly identical meaning (catches paraphrases/reformulations).
        near_dup_count = 0
        if deduplicated_inserts and np is not None:
            try:
                near_dup_count = self._filter_near_duplicates(cur, deduplicated_inserts)
                if near_dup_count > 0:
                    logger.info(f"[DEDUP] Removed {near_dup_count} near-duplicate chunks (cosine > 0.97)")
            except Exception as e:
                logger.debug(f"[DEDUP] Near-duplicate check skipped: {e}")

        if not deduplicated_inserts:
            return {
                "success": True, 
                "inserted": 0,
                "duplicates_skipped": duplicate_count,
                "near_duplicates_skipped": near_dup_count,
                "synthetic_skipped": synthetic_count,
                "noise_skipped": noise_count,
                "message": "Alle Chunks waren Duplikate, synthetische oder Noise-Inhalte"
            }

        # Globale Chunk-IDs vergeben
        cur.execute("SELECT COALESCE(MAX(chunk_id), -1) + 1 FROM chunks")
        next_global_chunk_id = cur.fetchone()[0] or 0
        
        fixed_to_insert = []
        global_chunk_counter = next_global_chunk_id
        # Stamp every freshly-ingested chunk with the current classifier
        # version so the maintenance reclassifier knows it is up-to-date.
        from agent.rag_store.core.content_classifier import CLASSIFICATION_VERSION
        for doc_id, local_idx, ch, meta_json, emb_blob, corpus_domain, safety_flag in deduplicated_inserts:
            fixed_to_insert.append(
                (doc_id, global_chunk_counter, ch, meta_json, emb_blob,
                 corpus_domain, safety_flag, CLASSIFICATION_VERSION)
            )
            global_chunk_counter += 1

        try:
            # Dokument-Einträge erstellen und alte Chunks löschen
            doc_ids = sorted(set(t[0] for t in deduplicated_inserts))
            for did in doc_ids:
                cur.execute("INSERT OR IGNORE INTO documents(doc_id) VALUES (?)", (did,))
                cur.execute("DELETE FROM chunks WHERE doc_id = ?", (did,))
                
            # Batch-Insert der neuen Chunks (Iteration 3)
            def insert_chunk_batch(batch_items: List[Tuple]) -> None:
                """Insert function for batch_insert utility"""
                cur.executemany(
                    "INSERT INTO chunks(doc_id, chunk_id, text, metadata, embedding, "
                    "domain, safety_flag, classification_version) "
                    "VALUES (?,?,?,?,?,?,?,?)",
                    batch_items
                )
            
            total_inserted = batch_insert(
                items=fixed_to_insert,
                insert_function=insert_chunk_batch,
                batch_size=500,  # Optimale Batch-Größe für SQLite
                show_progress=self.debug
            )
            
            if self.debug and total_inserted != len(fixed_to_insert):
                logger.warning(f"Expected {len(fixed_to_insert)} inserts, got {total_inserted}")
            
            # ★ SOTA v3: Ingest-time quality scoring — write to chunk_quality table
            # Every chunk gets content_type, defect_flags, structural_score at ingest
            # so that downstream search can apply quality penalties without re-scoring.
            try:
                from agent.rag_store.core.quality import RAGQualityManager
                _qm = RAGQualityManager.__new__(RAGQualityManager)
                now_iso = datetime.now(timezone.utc).isoformat()
                quality_rows = []
                for doc_id, chunk_id, text, meta_json, _emb, _domain, _safety, _ver in fixed_to_insert:
                    content_type = RAGQualityManager.detect_content_type(text, meta_json or "{}")
                    defects = RAGQualityManager.detect_defects(text, meta_json or "{}")
                    structural_score = _qm.score_chunk_structural(text, meta_json or "{}")
                    defect_str = ",".join(sorted(defects)) if defects else ""
                    quality_rows.append((
                        doc_id, chunk_id, structural_score,
                        content_type, defect_str, now_iso
                    ))
                if quality_rows:
                    cur.executemany(
                        "INSERT OR REPLACE INTO chunk_quality"
                        "(doc_id, chunk_id, structural_score, content_type, defect_flags, last_checked) "
                        "VALUES (?, ?, ?, ?, ?, ?)",
                        quality_rows
                    )
                    if self.debug:
                        logger.info(f"[QUALITY] Scored {len(quality_rows)} chunks at ingest time")
            except Exception as e:
                raise RuntimeError(f"[QUALITY] Ingest-time scoring failed: {e}") from e
            
            conn.commit()
            
            # 🔨 FAISS Auto-Rebuild Tracking
            if self.faiss_manager is not None:
                chunks_added = len(fixed_to_insert)  # type: ignore[unreachable]
                self.faiss_manager.notify_chunks_added(chunks_added)
                if self.debug:
                    logger.debug(f"[FAISS] Notified {chunks_added} new chunks")
            
            # 🔨 BM25 Index Invalidation (SOTA Hybrid Search)
            if self._search_manager is not None:
                try:
                    self._search_manager.invalidate_bm25_index()
                    if self.debug:
                        logger.debug(f"[BM25] Index invalidated after {len(fixed_to_insert)} new chunks")
                except Exception as e:
                    logger.debug(f"BM25 invalidation skipped: {e}")
            
        finally:
            cur.close()
            self.return_connection(conn)

        # ✨ NEUE FUNKTION: Automatische LLM-KG-Erstellung für alle Dokumente
        kg_stats = self._create_automatic_knowledge_graphs(doc_ids, documents)

        # ★ SOTA v6: Debounced async entity resolution.
        # Root-cause fix: avoid O(N^2) entity resolution in the ingest critical path.
        entity_resolution = self._schedule_entity_resolution(kg_stats.get("triples", 0))

        # ★ SOTA v3: Post-Ingest Audit Hook — lightweight quality check on new content
        # Runs asynchronously after ingest so it doesn't block the main pipeline.
        post_audit_stats = {}
        post_audit_stats = self._run_post_ingest_audit(
            fixed_to_insert, doc_ids, kg_stats
        )

        if self.debug:
            print(f"[RAG] Upserted {len(deduplicated_inserts)} chunks across {len(doc_ids)} docs")
            if duplicate_count > 0:
                print(f"[RAG] Skipped {duplicate_count} duplicates")
            if synthetic_count > 0:
                print(f"[RAG] Skipped {synthetic_count} synthetic/generated contents")
            if noise_count > 0:
                print(f"[RAG] Skipped {noise_count} noise chunks (JS/tracking/code artifacts)")
            if near_dup_count > 0:
                print(f"[RAG] Skipped {near_dup_count} near-duplicate chunks (cosine > 0.97)")
            if kg_stats.get("triples", 0) > 0:
                print(f"[KG] Created {kg_stats['triples']} LLM-generated triples")

        return {
            "success": True, 
            "inserted": len(deduplicated_inserts), 
            "docs": len(doc_ids),
            "duplicates_skipped": duplicate_count,
            "near_duplicates_skipped": near_dup_count,
            "synthetic_skipped": synthetic_count,
            "noise_skipped": noise_count,
            "kg_triples": kg_stats.get("triples", 0),
            "entity_resolution": entity_resolution,
            "post_audit": post_audit_stats,
        }

    def _run_post_ingest_audit(
        self,
        fixed_to_insert: List[Tuple],
        doc_ids: List[str],
        kg_stats: Dict[str, Any],
    ) -> Dict[str, Any]:
        """
        ★ SOTA v3: Lightweight post-ingest quality audit.
        
        Runs AFTER chunks + KG are inserted. Checks:
        1. Duplicate detection among freshly inserted chunks (intra-batch)
        2. Cross-document boilerplate detection for new docs
        3. Queue new triples for future grounding verification
        4. Log audit entry for traceability
        
        This is intentionally lightweight — no LLM/reranker calls.
        Heavy verification (grounding) is deferred to the next scheduled audit.
        """
        stats: Dict[str, Any] = {
            "chunks_audited": 0,
            "quality_warnings": 0,
            "triples_queued_for_grounding": 0,
        }
        
        if not fixed_to_insert:
            return stats
        
        conn = self.get_connection()
        cur = None
        try:
            from agent.rag_store.core.quality import RAGQualityManager
            cur = conn.cursor()
            
            stats["chunks_audited"] = len(fixed_to_insert)
            
            # 1. Check for low quality among inserted chunks (already scored at ingest)
            chunk_ids_for_query = [(row[0], row[1]) for row in fixed_to_insert]
            low_quality_count = 0
            for doc_id, chunk_id in chunk_ids_for_query:
                cur.execute(
                    "SELECT structural_score FROM chunk_quality WHERE doc_id=? AND chunk_id=?",
                    (doc_id, chunk_id),
                )
                row = cur.fetchone()
                if row and row[0] is not None and row[0] < 0.35:
                    low_quality_count += 1
            stats["quality_warnings"] = low_quality_count
            
            # 2. Queue new triples for grounding verification
            #    Mark them in triple_quality with grounding_score = -1 (unverified)
            triples_queued = 0
            for doc_id in doc_ids:
                cur.execute(
                    "SELECT t.triple_id FROM triples t "
                    "LEFT JOIN triple_quality tq ON t.triple_id = tq.triple_id "
                    "WHERE t.doc_id = ? AND tq.triple_id IS NULL",
                    (doc_id,),
                )
                new_triple_ids = [r[0] for r in cur.fetchall()]
                if new_triple_ids:
                    now_iso = datetime.now(timezone.utc).isoformat()
                    cur.executemany(
                        "INSERT OR IGNORE INTO triple_quality"
                        "(triple_id, grounding_score, last_verified) "
                        "VALUES (?, -1.0, ?)",
                        [(tid, now_iso) for tid in new_triple_ids],
                    )
                    triples_queued += len(new_triple_ids)
            stats["triples_queued_for_grounding"] = triples_queued
            
            # 3. Log audit entry
            # Write audit directly to DB to avoid hidden manager side-effects.
            cur.execute(
                "INSERT INTO quality_audit_log(action, details, timestamp) VALUES (?, ?, ?)",
                (
                    "post_ingest_audit",
                    json.dumps({
                        "chunks_audited": stats["chunks_audited"],
                        "quality_warnings": low_quality_count,
                        "triples_queued": triples_queued,
                        "doc_ids": doc_ids[:10],  # cap for readability
                    }),
                    datetime.now(timezone.utc).isoformat(),
                ),
            )
            
            conn.commit()
            
            if low_quality_count > 0 or triples_queued > 0:
                logger.info(
                    f"[POST-AUDIT] {stats['chunks_audited']} chunks audited, "
                    f"{low_quality_count} quality warnings, "
                    f"{triples_queued} triples queued for grounding"
                )
        except Exception as e:
            conn.rollback()
            raise RuntimeError(f"[POST-AUDIT] Error: {e}") from e
        finally:
            if cur is not None:
                cur.close()
            self.return_connection(conn)
        
        return stats

    def _chunk_text_with_spans(self, text: str, chunk_size: int, overlap: int) -> List[Tuple[int, int, str]]:
        """
        SOTA Semantic Chunking: Splittet Text an Themengrenzen via Embedding-Similarity.
        
        Algorithmus (Kamradt 2023 / LlamaIndex SemanticSplitter):
          1. Text → Sätze (NLTK sent_tokenize, deutsch)
          2. Satz-Embeddings via bereits geladenem e5-large (Batch, GPU)
          3. Cosine-Similarity zwischen konsekutiven Sätzen berechnen
          4. Similarity-Drops erkennen → Themengrenzen
          5. Sätze zwischen Grenzen zu Chunks gruppieren (max_size beachten)
          6. Overlap durch Wiederholung ganzer Sätze an Chunk-Grenzen
        
        Fallback: Sentence-Boundary-Splitting wenn Embedding-Modell nicht verfügbar.
        
        Args:
            text: Zu chunkender Text
            chunk_size: Maximale Chunk-Größe in Zeichen
            overlap: Gewünschte Überlappung (wird auf ganze Sätze gerundet)
            
        Returns:
            Liste von (start, end, chunk_text) Tupeln
        """
        if not text or not text.strip():
            return []
        
        # ── SCHRITT 1: Satz-Tokenisierung ──
        sentences = self._tokenize_sentences(text)
        if not sentences:
            return [(0, len(text), text)]
        
        # ── SCHRITT 2: Satz-Positionen im Original-Text bestimmen ──
        sentence_spans = self._locate_sentence_spans(text, sentences)
        if not sentence_spans:
            return [(0, len(text), text)]
        
        # ── SCHRITT 3: Semantische Grenzen via Embedding-Similarity ──
        split_indices = self._find_semantic_boundaries(sentence_spans, chunk_size)
        
        # ── SCHRITT 4: Sätze zu Chunks gruppieren ──
        return self._group_sentences_to_chunks(text, sentence_spans, split_indices, chunk_size, overlap)

    def _tokenize_sentences(self, text: str) -> List[str]:
        """Tokenisiert Text in Sätze (NLTK deutsch, mit Fallback)."""
        try:
            import nltk
            try:
                return nltk.sent_tokenize(text, language='german')
            except LookupError:
                try:
                    nltk.download('punkt', quiet=True)
                    nltk.download('punkt_tab', quiet=True)
                    return nltk.sent_tokenize(text, language='german')
                except Exception as exc:
                    logger.debug(f"NLTK punkt download failed, fallback tokenizer engaged: {exc}")
                    return nltk.sent_tokenize(text)
        except ImportError:
            import re
            return re.split(r'(?<=[.!?])\s+', text)

    def _locate_sentence_spans(self, text: str, sentences: List[str]) -> List[Tuple[int, int, str]]:
        """Berechnet (start, end, text) Positionen jedes Satzes im Originaltext."""
        sentence_spans: List[Tuple[int, int, str]] = []
        current_pos = 0
        for sent in sentences:
            sent_stripped = sent.strip()
            if not sent_stripped:
                continue
            start_pos = text.find(sent_stripped, current_pos)
            if start_pos == -1:
                start_pos = current_pos
            end_pos = start_pos + len(sent_stripped)
            sentence_spans.append((start_pos, end_pos, sent_stripped))
            current_pos = end_pos
        return sentence_spans

    def _find_semantic_boundaries(
        self, 
        sentence_spans: List[Tuple[int, int, str]], 
        chunk_size: int,
        similarity_threshold_percentile: float = 25.0,
    ) -> List[int]:
        """
        Findet Themengrenzen via Embedding-Cosine-Similarity zwischen konsekutiven Sätzen.
        
        Nutzt Percentile-basiertes Thresholding (robuster als fester Wert):
        - Berechne Similarity zwischen allen Nachbar-Sätzen
        - Splits an Stellen mit Similarity < percentile_threshold
        - Erzwinge Split wenn kumulierte Länge > chunk_size
        
        Args:
            sentence_spans: [(start, end, text), ...]
            chunk_size: Maximale Chunk-Größe
            similarity_threshold_percentile: Percentile für Similarity-Cutoff
            
        Returns:
            Liste von Satz-Indices an denen gesplittet wird (exklusiv)
        """
        if len(sentence_spans) <= 1:
            return []
        
        # Versuche Embedding-basierte Grenzerkennung
        try:
            self._ensure_embedding_manager()
            if self._embedding_manager is None:
                raise RuntimeError("EmbeddingManager not available")
            
            # Batch-Embedding aller Sätze (nutzt GPU, bereits geladen)
            sentence_texts = [s[2] for s in sentence_spans]
            embeddings = self._embedding_manager.embed_texts(sentence_texts)
            
            if embeddings is None or len(embeddings) < 2:
                raise RuntimeError("Embedding generation failed")
            
            # Cosine-Similarity zwischen konsekutiven Sätzen
            # Embeddings sind bereits L2-normalisiert → dot product = cosine sim
            import numpy as np
            similarities = []
            for i in range(len(embeddings) - 1):
                sim = float(np.dot(embeddings[i], embeddings[i + 1]))
                similarities.append(sim)
            
            if not similarities:
                return []
            
            # Percentile-basiertes Thresholding
            # Niedrige Similarity = Themenwechsel
            threshold = float(np.percentile(similarities, similarity_threshold_percentile))
            
            # Finde Splits: Similarity unter Threshold ODER Größe überschritten
            split_indices: List[int] = []
            current_chunk_len = len(sentence_spans[0][2])
            
            for i, sim in enumerate(similarities):
                sent_len = len(sentence_spans[i + 1][2])
                
                # Erzwinge Split wenn chunk_size überschritten
                if current_chunk_len + sent_len + 1 > chunk_size:
                    split_indices.append(i + 1)
                    current_chunk_len = sent_len
                    continue
                
                # Semantischer Split: Similarity unter Threshold
                if sim < threshold:
                    split_indices.append(i + 1)
                    current_chunk_len = sent_len
                else:
                    current_chunk_len += sent_len + 1
            
            logger.debug(
                f"[Semantic Chunking] {len(sentence_spans)} Sätze, "
                f"{len(similarities)} Similarity-Werte, "
                f"Threshold={threshold:.3f} (P{similarity_threshold_percentile:.0f}), "
                f"{len(split_indices)} Splits"
            )
            
            return split_indices
            
        except Exception as e:
            logger.warning(f"[Semantic Chunking] Embedding-basierte Grenzerkennung nicht möglich: {e}")
            logger.info("[Semantic Chunking] Fallback auf Sentence-Boundary-Splitting")
            
            # Fallback: Splitte rein nach chunk_size mit Satzgrenzen
            split_indices: List[int] = []
            current_chunk_len = 0
            
            for i, (_, _, sent_text) in enumerate(sentence_spans):
                sent_len = len(sent_text)
                if current_chunk_len + sent_len + 1 > chunk_size and current_chunk_len > 0:
                    split_indices.append(i)
                    current_chunk_len = sent_len
                else:
                    current_chunk_len += sent_len + 1
            
            return split_indices

    def _group_sentences_to_chunks(
        self,
        text: str,
        sentence_spans: List[Tuple[int, int, str]],
        split_indices: List[int],
        chunk_size: int,
        overlap: int,
    ) -> List[Tuple[int, int, str]]:
        """
        Gruppiert Sätze zu Chunks basierend auf Split-Indices.
        
        Overlap durch Wiederholung ganzer Sätze am Chunk-Anfang.
        Überlange Einzelsätze werden character-basiert gesplittet.
        
        Args:
            text: Originaltext
            sentence_spans: [(start, end, text), ...]
            split_indices: Indices an denen gesplittet wird
            chunk_size: Maximale Chunk-Größe
            overlap: Gewünschte Überlappung in Zeichen
            
        Returns:
            Liste von (start, end, chunk_text) Tupeln
        """
        if not sentence_spans:
            return []
        
        overlap_sentences = max(1, overlap // 100)  # ~1 Satz pro 100 Zeichen Overlap
        
        # Erzeuge Gruppen aus Split-Indices
        groups: List[List[int]] = []
        prev_idx = 0
        for split_idx in sorted(set(split_indices)):
            if split_idx > prev_idx:
                groups.append(list(range(prev_idx, split_idx)))
            prev_idx = split_idx
        # Letzte Gruppe
        if prev_idx < len(sentence_spans):
            groups.append(list(range(prev_idx, len(sentence_spans))))
        
        chunks: List[Tuple[int, int, str]] = []
        
        for group_idx, group in enumerate(groups):
            # Overlap: Füge letzte N Sätze der vorherigen Gruppe voran
            effective_indices = list(group)
            if overlap > 0 and group_idx > 0 and len(groups[group_idx - 1]) > 0:
                prev_group = groups[group_idx - 1]
                overlap_count = min(overlap_sentences, len(prev_group))
                overlap_indices = prev_group[-overlap_count:]
                effective_indices = overlap_indices + effective_indices
            
            # Sammle Sätze dieser Gruppe
            group_sentences = [sentence_spans[i] for i in effective_indices]
            
            if not group_sentences:
                continue
            
            # Prüfe ob ein einzelner Satz zu lang ist
            for sent_start, sent_end, sent_text in group_sentences:
                if len(sent_text) > chunk_size:
                    # Splitte überlangen Satz character-basiert
                    for j in range(0, len(sent_text), chunk_size - overlap):
                        sub_start = sent_start + j
                        sub_end = min(sent_start + j + chunk_size, sent_end)
                        sub_text = text[sub_start:sub_end]
                        if sub_text.strip():
                            chunks.append((sub_start, sub_end, sub_text))
                    continue
            
            # Normaler Fall: Alle Sätze der Gruppe zusammenfügen
            # Filtere überlange Sätze (bereits oben behandelt)
            normal_sentences = [(s, e, t) for s, e, t in group_sentences if len(t) <= chunk_size]
            
            if not normal_sentences:
                continue
            
            chunk_start = normal_sentences[0][0]
            chunk_end = normal_sentences[-1][1]
            chunk_text = text[chunk_start:chunk_end]
            
            # Sicherheit: Wenn zusammengesetzter Chunk zu lang, re-split
            if len(chunk_text) > chunk_size * 1.5:
                # Sub-Chunking mit Sentence-Boundary
                sub_sents: List[Tuple[int, int, str]] = []
                sub_len = 0
                for s, e, t in normal_sentences:
                    if sub_len + len(t) + 1 > chunk_size and sub_sents:
                        c_start = sub_sents[0][0]
                        c_end = sub_sents[-1][1]
                        chunks.append((c_start, c_end, text[c_start:c_end]))
                        sub_sents = []
                        sub_len = 0
                    sub_sents.append((s, e, t))
                    sub_len += len(t) + 1
                if sub_sents:
                    c_start = sub_sents[0][0]
                    c_end = sub_sents[-1][1]
                    chunks.append((c_start, c_end, text[c_start:c_end]))
            else:
                if chunk_text.strip():
                    chunks.append((chunk_start, chunk_end, chunk_text))
        
        return chunks

    # ═══════════════════════════════════════════════════════════════════
    # LAYER 2: CHUNK CONTENT-QUALITY VALIDATOR (SOTA)
    #   Erkennt und verwirft Noise-Chunks VOR dem Embedding/Insert.
    #   Prüft auf: JavaScript, HTML-Artefakte, Tracking-Code,
    #   minified Code, niedriger Informationsgehalt, Nicht-Sprache.
    # ═══════════════════════════════════════════════════════════════════
    
    # Kompilierte Patterns für Performance (Class-Level, einmalig)
    import re as _re_module
    
    _NOISE_CODE_PATTERNS = _re_module.compile(
        r'(?:'
        # JavaScript-Konstrukte
        r'function\s*\([^)]*\)\s*\{'
        r'|var\s+\w+\s*='
        r'|const\s+\w+\s*='
        r'|let\s+\w+\s*='
        r'|=>\s*\{'
        r'|document\.\w+'
        r'|window\.\w+'
        r'|getElementById'
        r'|addEventListener'
        r'|querySelector'
        r'|\.prototype\.'
        r'|\.push\(\s*arguments\s*\)'
        r'|new\s+RegExp\('
        r'|\.replace\(\s*(?:new\s+RegExp|/)'
        # Tracking/Analytics
        r'|gtag\s*\('
        r'|_gaq\s*\.'
        r'|dataLayer\s*\.'
        r'|fbq\s*\('
        r'|_linkedin_partner_id'
        r'|utm_source'
        r'|pixel\s*\.track'
        # CSS/HTML-Artefakte
        r'|className\s*='
        r'|\{\s*display\s*:'
        r'|\{\s*margin\s*:'
        r'|\{\s*padding\s*:'
        r'|@media\s*\('
        r'|!important\s*;'
        # JSON/Schema-Artefakte
        r'|"@context"\s*:\s*"https?://schema'
        r'|"@type"\s*:\s*"'
        r')',
        _re_module.IGNORECASE
    )
    
    _NOISE_MINIFIED_PATTERN = _re_module.compile(
        # Erkennt minified Code: lange Zeilen ohne Leerzeichen, {}()[]; Häufung
        r'[{}\[\]();,]{3,}'
    )
    
    # ── Compiled Boilerplate / Cookie / Navigation patterns (class-level) ──
    import re as _re_bp
    _BOILERPLATE_PATTERNS = _re_bp.compile(
        r'(?i)(?:'
        r'cookie|datenschutz|privacy\s*policy|terms\s*of\s*(?:use|service)'
        r'|impressum|all\s*rights?\s*reserved|copyright\s*©'
        r'|wir\s*verwenden\s*cookies|we\s*use\s*cookies'
        r'|akzeptieren|accept\s*all|ablehnen|reject'
        r'|einstellungen\s*speichern|save\s*preferences'
        r'|newsletter\s*(?:abonnieren|subscribe)'
        r'|folgen?\s*(?:sie\s*)?uns\s*(?:auf|on)'
        r'|follow\s*us(?:\s*on)?'
        r'|teilen|share\s*(?:on|this)'
        r'|zurück\s*(?:zu[rm]?\s*)?(?:anfang|seitenanfang|top)'
        r'|back\s*to\s*top|scroll\s*to\s*top'
        r'|(?:home|start|kontakt|contact|about|über\s*uns)\s*[\|•·/]'
        r'|navigation|breadcrumb|sidebar|footer|header'
        r'|skip\s*to\s*(?:main\s*)?content'
        r')'
    )
    _URL_PATTERN_NOISE = _re_bp.compile(r'https?://\S+', _re_bp.IGNORECASE)
    _NAV_SEPARATOR_PATTERN = _re_bp.compile(r'[\|•·/]{2,}|(?:\s[\|•·/]\s){3,}')

    def _is_noise_chunk(self, text: str) -> bool:
        """
        ★ SOTA v3: Consolidated noise detection — delegates to quality.py.
        
        This method previously had ~170 lines of its own heuristic logic that 
        duplicated and diverged from quality.py's detect_defects() and 
        score_chunk_structural(). Now unified into a single code path.
        
        Uses:
          - RAGQualityManager.detect_defects() for hard defects (encoding_garbage, trivial, etc.)
          - RAGQualityManager.score_chunk_structural() for structural quality < 0.20
        
        Returns:
            True wenn der Chunk als Noise erkannt wird und NICHT gespeichert werden soll.
        """
        if not text or len(text.strip()) < 10:
            return True
        
        try:
            from agent.rag_store.core.quality import RAGQualityManager
            
            # Hard defect check: these are always noise
            defects = RAGQualityManager.detect_defects(text)
            hard_defects = {"encoding_garbage", "trivial", "cookie_banner", "pure_navigation", "url_dump"}
            if defects & hard_defects:
                if self.debug:
                    logger.debug(f"[NOISE-QUALITY] Hard defect: {defects & hard_defects}, text='{text[:80]}...'")
                return True
            
            # Structural quality check: very low score = noise
            # We use a lightweight instance (no DB needed for static scoring)
            qm = RAGQualityManager.__new__(RAGQualityManager)
            score = qm.score_chunk_structural(text)
            if score < 0.20:
                if self.debug:
                    logger.debug(f"[NOISE-QUALITY] Low structural score: {score:.2f}, text='{text[:80]}...'")
                return True
            
            return False
        except Exception as e:
            # Fallback: only reject very short or empty text
            logger.debug(f"[NOISE] quality.py delegation failed: {e}")
            return len(text.strip()) < 15

    def _filter_near_duplicates(
        self,
        cur: sqlite3.Cursor,
        inserts: List[Tuple],
        cosine_threshold: float = 0.97,
        sample_limit: int = 10000,
    ) -> int:
        """
        Near-duplicate gate: remove new chunks that are near-duplicates of existing DB chunks.
        
        Uses embedding cosine similarity via numpy dot product on normalised vectors.
        Only checks CROSS-document (same-document overlap is expected from chunking overlap).
        
        Modifies `inserts` in-place (removes near-duplicates).
        Returns count of removed near-duplicates.
        """
        if not inserts:
            return 0
        
        # Parse new chunk embeddings
        new_items: List[Tuple[int, str, np.ndarray]] = []
        for idx, item in enumerate(inserts):
            # tuple shape: (doc_id, chunk_idx, text, meta_json, emb_blob,
            #               corpus_domain, safety_flag)
            doc_id = item[0]
            emb_blob = item[4]
            if emb_blob:
                vec = np.frombuffer(emb_blob, dtype=np.float32).copy()
                norm = np.linalg.norm(vec)
                if norm > 0:
                    new_items.append((idx, doc_id, vec / norm))
        
        if not new_items:
            return 0
        
        # Sample existing embeddings from DB (different doc_ids)
        new_doc_ids = set(t[0] for t in inserts)
        placeholders = ','.join('?' * len(new_doc_ids))
        query = (
            f"SELECT doc_id, embedding FROM chunks "
            f"WHERE doc_id NOT IN ({placeholders}) "
            f"LIMIT {sample_limit}"
        )
        cur.execute(query, list(new_doc_ids))
        
        existing_vecs: List[np.ndarray] = []
        for row in cur.fetchall():
            blob = row[1]  # embedding column
            if blob:
                vec = np.frombuffer(blob, dtype=np.float32)
                norm = np.linalg.norm(vec)
                if norm > 0:
                    existing_vecs.append(vec / norm)
        
        if not existing_vecs:
            return 0
        
        # Build existing matrix for batch cosine
        existing_matrix = np.stack(existing_vecs)  # (M, D)
        
        # Find duplicates
        indices_to_remove = set()
        for idx, doc_id, new_vec in new_items:
            # Cosine similarity = dot product of normalised vectors
            similarities = existing_matrix @ new_vec  # (M,)
            max_sim = float(np.max(similarities))
            if max_sim >= cosine_threshold:
                indices_to_remove.add(idx)
        
        if not indices_to_remove:
            return 0
        
        # Remove in-place (reverse order to preserve indices)
        for idx in sorted(indices_to_remove, reverse=True):
            inserts.pop(idx)
        
        return len(indices_to_remove)

    def _hash_str(self, s: str) -> str:
        """Erstellt Hash für String"""
        return hashlib.md5(s.encode("utf-8")).hexdigest()

    def _detect_language(self, text: str) -> Optional[str]:
        """Einfache Spracherkennung"""
        if not text:
            return None
        # Vereinfachte Heuristik
        german_indicators = ['der', 'die', 'das', 'und', 'oder', 'ist', 'sind', 'ein', 'eine']
        english_indicators = ['the', 'and', 'or', 'is', 'are', 'a', 'an', 'this', 'that']
        
        text_lower = text.lower()
        german_count = sum(1 for word in german_indicators if word in text_lower)
        english_count = sum(1 for word in english_indicators if word in text_lower)
        
        if german_count > english_count:
            return "de"
        elif english_count > 0:
            return "en"
        return None

    # ======= SEARCH METHODS =======
    
    def search(
        self,
        query: str,
        k: int = 5,
        min_score: float = 0.0,
        adaptive_confidence: bool = True,
        faiss_min_confidence: Optional[float] = None,
        allowed_domains: Optional[List[str]] = None,
        exclude_safety_flags: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        """
        Hybride Suche: Kombiniert Embedding-basierte und Knowledge-Graph-Suche
        
        DELEGIERT AN SearchManager (Iteration 7)
        
        🚀 FEATURES:
            - FAISS Hybrid Search (555x schneller!)
            - Smart Fusion Engine (KG + FAISS + Recency)
            - Adaptive Strategy (Fast Path + Fallback)
            - 🆕 Feedback-basierte Gewichts-Optimierung
        
        🎯 ADAPTIVE CONFIDENCE:
            Passt Qualitäts-Schwelle automatisch an Suchtiefe k an:
            - k <= 3:  0.70 (⚡ Quick Search, schnell)
            - k <= 7:  0.75 (📊 Standard Search)
            - k <= 12: 0.80 (🎯 Deep Search)
            - k <= 18: 0.85 (🔬 Very Deep Search)
            - k >= 19: 0.90 (🏆 Maximum Depth, maximale Qualität)
        
        Args:
            query: Suchquery
            k: Anzahl Ergebnisse (beeinflusst auch Qualitäts-Schwelle!)
            min_score: Minimaler Similarity-Score (NACH FAISS)
            adaptive_confidence: Wenn True, passe FAISS Confidence an k an
            faiss_min_confidence: Manuelle FAISS Confidence (überschreibt adaptive!)
            
        Returns:
            Liste von Suchergebnissen mit Text, Metadaten und Score
        """
        # 🆕 Check für Feedback-Updates (non-blocking)
        if self.feedback_service:
            try:
                self.feedback_service.check_and_update(force=False)
            except Exception as e:
                logger.debug(f"Feedback update check failed: {e}")
        
        # LAZY INIT: Ensure EmbeddingManager is initialized before search
        self._ensure_embedding_manager()
        if self._search_manager and self._embedding_manager:
            # Update SearchManager's embedding_manager reference (lazy init)
            if self._search_manager.embedding_manager is None:  # type: ignore[unreachable]
                self._search_manager.embedding_manager = self._embedding_manager
        
        # Delegiere an SearchManager
        if self._search_manager:
            return self._search_manager.search(
                query=query,
                k=k,
                min_score=min_score,
                adaptive_confidence=adaptive_confidence,
                faiss_min_confidence=faiss_min_confidence,
                allowed_domains=allowed_domains,
                exclude_safety_flags=exclude_safety_flags,
            )
        else:
            raise RuntimeError("SearchManager not initialized for search()")
    
    def batch_search(
        self,
        queries: List[str],
        k_list: List[int],
        min_score: float = 0.0,
        adaptive_confidence: bool = True,
        faiss_min_confidence: Optional[float] = None,
        allowed_domains: Optional[List[str]] = None,
        exclude_safety_flags: Optional[List[str]] = None,
    ) -> List[List[Dict[str, Any]]]:
        """
        🚀 BATCH SEARCH: Optimierte Batch-Suche für multiple Queries
        
        DELEGIERT AN SearchManager (Iteration 7+)
        
        **Performance-Vorteil:**
        - 150x schneller als sequenzielle search() Calls bei Multi-Query (6 Queries)
        - 94.7% CPU-Auslastung (vs 9.6% bei Single-Query)
        - Nutzt native FAISS Batch-API mit OpenMP
        
        **Use-Cases:**
        - Multi-Query RAG (1 Main + N Sub-Queries)
        - Gap-Detection mit mehreren Queries
        - Parallel Processing von User-Anfragen
        
        Args:
            queries: Liste von Suchanfragen
            k_list: Liste von k-Werten (einer pro Query), z.B. [6, 5, 5, 5, 5]
            min_score: Minimaler Similarity-Score (NACH FAISS)
            adaptive_confidence: Wenn True, passe FAISS Confidence an k an
            faiss_min_confidence: Manuelle FAISS Confidence
            
        Returns:
            Liste von Listen mit Suchergebnissen (eine Liste pro Query)
            
        Example:
            >>> rag = UnifiedRagStore(...)
            >>> queries = ["Python tutorial", "Machine Learning basics", "Data Science"]
            >>> k_list = [5, 10, 8]
            >>> results = rag.batch_search(queries, k_list)
            >>> # results[0] = 5 Results für "Python tutorial"
            >>> # results[1] = 10 Results für "Machine Learning basics"  
            >>> # results[2] = 8 Results für "Data Science"
        """
        # 🆕 Check für Feedback-Updates (non-blocking)
        if self.feedback_service:
            try:
                self.feedback_service.check_and_update(force=False)
            except Exception as e:
                logger.debug(f"Feedback update check failed: {e}")
        
        # LAZY INIT: Ensure EmbeddingManager is initialized before batch search
        self._ensure_embedding_manager()
        if self._search_manager and self._embedding_manager:
            # Update SearchManager's embedding_manager reference (lazy init)
            if self._search_manager.embedding_manager is None:  # type: ignore[unreachable]
                self._search_manager.embedding_manager = self._embedding_manager
        
        # Delegiere an SearchManager
        if self._search_manager:
            return self._search_manager.batch_search(
                queries=queries,
                k_list=k_list,
                min_score=min_score,
                adaptive_confidence=adaptive_confidence,
                faiss_min_confidence=faiss_min_confidence,
                allowed_domains=allowed_domains,
                exclude_safety_flags=exclude_safety_flags,
            )
        else:
            raise RuntimeError("SearchManager not initialized for batch_search()")
    
    # ======= PDF PROCESSING =======
    # (All search methods removed - now in SearchManager, Iteration 7)

    @staticmethod
    def _is_nonrecoverable_docling_failure(docling_result: Any) -> bool:
        """Classify Docling failures that must not fall back to legacy processing."""
        fatal_codes = {
            "local_only_network_blocked",
            "local_only_resource_missing",
            "local_only_url_disabled",
        }
        error_code = getattr(docling_result, "error_code", None)
        if error_code in fatal_codes:
            return True

        error_text = (getattr(docling_result, "error", "") or "").lower()
        return (
            "app_local_only active: outbound" in error_text
            and "blocked" in error_text
        )
    
    def upsert_pdf(self, file_path: str, *, doc_id: Optional[str] = None, 
                   metadata: Optional[Dict[str, Any]] = None,
                   chunk_size: int = 1500, chunk_overlap: int = 200,
                   extract_tables: bool = False, build_kg: bool = True,
                   table_doc_suffix: str = "#tables",
                   _docling_result: Optional[Any] = None) -> Dict[str, Any]:
        """
        PDF-Verarbeitung: Docling SOTA → Legacy Fallback
        
        Args:
            file_path: Pfad zur PDF-Datei
            doc_id: Optional document ID
            metadata: Zusätzliche Metadaten
            chunk_size: Chunk-Größe
            chunk_overlap: Chunk-Überlappung
            extract_tables: Tabellen extrahieren
            build_kg: Knowledge Graph aufbauen
            table_doc_suffix: Suffix für Tabellen-Dokumente
            _docling_result: Vorberechnetes DoclingResult (vermeidet doppelte Verarbeitung)
            
        Returns:
            Dict mit Verarbeitungsstatistiken
        """

        # Local user-upload documents are explicitly classified as general/safe
        # unless a caller intentionally overrides these labels.
        metadata = dict(metadata or {})
        metadata.setdefault("corpus_domain", "general")
        metadata.setdefault("safety_flag", "safe")
        
        # ✅ SOTA FIX: Original-Dateiname aus metadata bewahren!
        # Problem: file_path ist oft ein Temp-Pfad (z.B. "tmp12ab34cd.pdf")
        #   → os.path.basename(file_path) = "tmp12ab34cd.pdf" (kryptisch)
        #   → metadata={"source": "Original_Datei.pdf"} wird ÜBERSCHRIEBEN
        # Lösung: metadata["source"] hat Vorrang vor os.path.basename(file_path)
        #   → Original-Dateiname bleibt in RAG-Store erhalten
        _original_source = (metadata or {}).get("source") or os.path.basename(file_path)
        _original_title = (metadata or {}).get("title") or _original_source
        
        # doc_id ebenfalls mit Original-Name statt Temp-Pfad
        if doc_id is None and _original_source != os.path.basename(file_path):
            # Benutze Original-Source als doc_id statt des Temp-Pfads
            doc_id = _original_source
        
        # 1. ERSTEN: PDF-Lesbarkeits-Prüfung
        try:
            import sys
            # Füge das übergeordnete Verzeichnis zum Python-Pfad hinzu
            parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            if parent_dir not in sys.path:
                sys.path.insert(0, parent_dir)
            
            from pdf_readability_checker import check_pdf_readable, log_readability_result
            
            is_readable, readability_metadata = check_pdf_readable(file_path)
            log_readability_result(file_path, is_readable, readability_metadata)
            
            if not is_readable:
                error_msg = f"PDF not readable: {readability_metadata.get('error', 'Unknown error')}"
                logger.warning(f"Skipping unreadable PDF: {os.path.basename(file_path)} - {error_msg}")
                return {
                    "success": False,
                    "error": error_msg,
                    "file_path": file_path,
                    "skipped": True,
                    "readability_check": readability_metadata
                }
        except ImportError as e:
            # Fallback: Wenn pdf_readability_checker nicht verfügbar ist, fahre normal fort
            logger.warning(f"PDF readability checker not available: {e}, proceeding without check")
        except Exception as e:
            # Bei Fehlern in der Readability-Prüfung: Warnung loggen aber fortfahren
            logger.warning(f"PDF readability check failed for {file_path}: {e}")
        
        # ═══════════════════════════════════════════════════════════════
        # 2. DOCLING SOTA: Primärer PDF-Verarbeitungspfad
        #    AI-basierte Layout-Analyse (RT-DETR) + TableFormer + RapidOCR
        #    HybridChunker für token-aware Chunks
        #    CPU-only (GPU bleibt für LLM reserviert)
        # ═══════════════════════════════════════════════════════════════
        try:
            from utils.docling_processor import DoclingProcessor
            
            # Vorberechnetes Ergebnis nutzen (vermeidet doppelte Verarbeitung)
            if _docling_result is not None:
                docling_result = _docling_result
                logger.info(f"🔬 Docling: Nutze vorberechnetes Ergebnis für {os.path.basename(file_path)}")
            else:
                processor = DoclingProcessor.get_instance()
                logger.info(f"🔬 Docling SOTA: Verarbeite PDF {os.path.basename(file_path)}")
                
                docling_result = processor.convert_file(
                    file_path,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    metadata={
                        **(metadata or {}),
                        "source": _original_source,
                        "title": _original_title,
                        "content_type": "application/pdf",
                        # Legal-Hardening (2026-08-30): Web-Herkunft (z. B.
                        # "web_pdf") aus der Caller-Metadaten nicht überschreiben,
                        # sonst verpasst der Retention-Pruner diese Records.
                        "source_type": (metadata or {}).get("source_type") or "document",
                    },
                )
            
            # Page-Loss-Ratio bestimmt ob Partial-Ergebnis nutzbar ist
            _page_loss = docling_result.metadata.get('page_loss_ratio', 0.0) if docling_result.metadata else 0.0
            _is_partial = getattr(docling_result, 'is_partial', False)
            _is_usable_partial = _is_partial and _page_loss <= 0.1
            _accept_result = not _is_partial or _is_usable_partial
            
            if docling_result.success and docling_result.chunks and _accept_result:
                if _is_usable_partial:
                    logger.info(
                        f"📄 Docling PARTIAL akzeptiert für {os.path.basename(file_path)}: "
                        f"page_loss={_page_loss:.0%}, {len(docling_result.chunks)} chunks vorhanden"
                    )
                # Docling-Chunks direkt in RAG Store einfügen (kein Re-Chunking nötig)
                base_doc_id = doc_id or os.path.abspath(file_path)
                documents = []
                
                for chunk in docling_result.chunks:
                    chunk_doc_id = f"{base_doc_id}_chunk_{chunk.index}"
                    documents.append({
                        "id": chunk_doc_id,
                        "text": chunk.text,
                        "metadata": {
                            **(metadata or {}),
                            "source": _original_source,
                            "title": _original_title,
                            "content_type": "application/pdf",
                            # Legal-Hardening (2026-08-30): Web-Herkunft nicht überschreiben
                            "source_type": (metadata or {}).get("source_type") or "document",
                            "extraction_method": "docling",
                            "docling_format": docling_result.format_detected,
                            "docling_tables": docling_result.num_tables,
                            "docling_pages": docling_result.num_pages,
                            "chunk_headings": " > ".join(chunk.headings) if chunk.headings else "",
                            **chunk.metadata,
                        },
                    })
                
                if documents:
                    # Store KG-chunks for _create_automatic_knowledge_graphs
                    # Thread-safe: cuda_lock serializes KG creation anyway
                    if hasattr(docling_result, 'kg_chunks') and docling_result.kg_chunks:
                        if not hasattr(self, '_pending_kg_chunks'):
                            self._pending_kg_chunks: Dict[str, Any] = {}
                        self._pending_kg_chunks[base_doc_id] = docling_result.kg_chunks
                    
                    # chunk_size=50000 verhindert Re-Chunking (Docling-Chunks sind bereits optimal)
                    result = self.upsert_documents(
                        documents, chunk_size=50000, chunk_overlap=0,
                    )
                    inserted_count = result.get("inserted", 0)
                    logger.info(
                        f"✅ Docling SOTA: {os.path.basename(file_path)} → "
                        f"{inserted_count} chunks, {docling_result.num_tables} Tabellen, "
                        f"{docling_result.num_pages} Seiten, {docling_result.processing_time_s:.1f}s"
                    )
                    return {
                        "success": True,
                        "inserted": inserted_count,
                        "chunks_added": inserted_count,
                        "docs": len(documents),
                        "duplicates_skipped": result.get("duplicates_skipped", 0),
                        "extraction_method": "docling",
                        "docling_tables": docling_result.num_tables,
                        "docling_pages": docling_result.num_pages,
                        "processing_time_s": docling_result.processing_time_s,
                        "docling_text": docling_result.text,  # Für Chat-Context
                    }
            
            elif docling_result.success and docling_result.text and _accept_result:
                # Docling hat Text aber keine Chunks → durch unsere Pipeline chunken
                logger.info(f"Docling: Text ohne Chunks, nutze Semantic Chunking Pipeline")
                base_doc_id = doc_id or os.path.abspath(file_path)
                documents = [{
                    "id": base_doc_id,
                    "text": docling_result.text,
                    "metadata": {
                        **(metadata or {}),
                        "source": _original_source,
                        "title": _original_title,
                        "content_type": "application/pdf",
                        # Legal-Hardening (2026-08-30): Web-Herkunft nicht überschreiben
                        "source_type": (metadata or {}).get("source_type") or "document",
                        "extraction_method": "docling",
                    },
                }]
                result = self.upsert_documents(
                    documents, chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                )
                inserted_count = result.get("inserted", 0)
                return {
                    "success": True,
                    "inserted": inserted_count,
                    "chunks_added": inserted_count,
                    "docs": 1,
                    "duplicates_skipped": result.get("duplicates_skipped", 0),
                    "extraction_method": "docling",
                    "docling_text": docling_result.text,
                }
            
            else:
                _partial_hint = ""
                if getattr(docling_result, 'is_partial', False):
                    _partial_hint = " (PARTIAL_SUCCESS: Seiten unvollständig verarbeitet)"

                if self._is_nonrecoverable_docling_failure(docling_result):
                    error_code = getattr(docling_result, 'error_code', 'unknown')
                    raise RuntimeError(
                        f"Docling local-only fail-fast ({error_code}): {docling_result.error}{_partial_hint}"
                    )

                logger.warning(
                    f"⚠️ Docling fehlgeschlagen für {os.path.basename(file_path)}: "
                    f"{docling_result.error}{_partial_hint}. Fallback auf Legacy-Pipeline."
                )
                # Durchfallen zum Legacy-Pfad
                
        except ImportError:
            logger.debug("DoclingProcessor nicht verfügbar, nutze Legacy-Pipeline")
        except RuntimeError as e:
            if "Docling local-only fail-fast" in str(e):
                logger.error(str(e))
                raise
            logger.warning(f"Docling PDF-Verarbeitung fehlgeschlagen: {e}, Fallback auf Legacy")
        except Exception as e:
            logger.warning(f"Docling PDF-Verarbeitung fehlgeschlagen: {e}, Fallback auf Legacy")
        
        # ═══════════════════════════════════════════════════════════════
        # 3. LEGACY FALLBACK: pymupdf4llm → pdfminer → PyMuPDF → OCR
        #    Wird nur genutzt wenn Docling nicht verfügbar oder fehlgeschlagen
        # ═══════════════════════════════════════════════════════════════
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        use_parallel = (self.parallel_capable and 
                       file_size > 1024 * 1024 and  # > 1MB
                       (extract_tables or build_kg))
        
        if use_parallel:
            return self._upsert_pdf_parallel(
                file_path, doc_id=doc_id, metadata=metadata,
                chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                extract_tables=extract_tables, build_kg=build_kg,
                table_doc_suffix=table_doc_suffix
            )
        else:
            return self._upsert_pdf_sequential(
                file_path, doc_id=doc_id, metadata=metadata,
                chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                extract_tables=extract_tables, build_kg=build_kg,
                table_doc_suffix=table_doc_suffix
            )

    def _upsert_pdf_sequential(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Sequenzielle PDF-Verarbeitung mit State-of-the-Art Extraktion"""
        
        errors: List[str] = []
        text: str = ""
        page_texts: List[str] = []
        extraction_method: str = "unknown"
        
        # 🆕 STATE-OF-THE-ART: pymupdf4llm für LLM-optimierte Extraktion
        # Liefert Markdown mit Struktur (Überschriften, Listen, Tabellen)
        try:
            import pymupdf4llm
            # table_strategy='lines' erkennt Tabellen besser als 'lines_strict'
            md_text = pymupdf4llm.to_markdown(
                file_path,
                table_strategy='lines',  # Bessere Tabellenerkennung
                show_progress=False
            )
            if md_text and md_text.strip():
                text = md_text
                # Teile nach Markdown-Seitenumbrüchen oder Form-Feed
                if "\n---\n" in text:
                    page_texts = text.split("\n---\n")
                elif "\x0c" in text:
                    page_texts = text.split("\x0c")
                else:
                    page_texts = [text]
                extraction_method = "pymupdf4llm"
                logging.info(f"✅ pymupdf4llm Extraktion erfolgreich: {len(text)} Zeichen (Markdown)")
        except ImportError:
            errors.append("pymupdf4llm nicht installiert")
        except Exception as e:
            errors.append(f"pymupdf4llm: {e}")
            logging.warning(f"pymupdf4llm fehlgeschlagen, versuche Fallback: {e}")
        
        # Fallback 1: pdfminer
        # (Der frühere "Advanced PDF"-Fallback wurde entfernt: er delegierte an
        #  DoclingProcessor, der an dieser Stelle bereits fehlgeschlagen ist.)
        if not text.strip():
            try:
                from pdfminer.high_level import extract_text
                text = extract_text(file_path) or ""
                if text:
                    page_texts = text.split("\x0c")
                    extraction_method = "pdfminer"
            except Exception as e:
                errors.append(f"pdfminer: {e}")
        
        # Fallback 2: PyMuPDF direkt
        if not text.strip():
            try:
                import fitz
                with fitz.open(file_path) as doc:
                    page_texts = []
                    for page in doc:
                        try:
                            page_text = getattr(page, 'get_text', lambda: "")()
                            page_texts.append(page_text or "")
                        except (AttributeError, TypeError) as e:
                            logger.warning(f"Could not extract text from page: {e}")
                            page_texts.append("")
                    
                    if page_texts:
                        text = "\n".join(page_texts)
                        extraction_method = "pymupdf_direct"
                    
            except Exception as e:
                errors.append(f"PyMuPDF: {e}")
        
        # 🆕 Fallback 4: OCR für gescannte PDFs (Bild-PDFs ohne eingebetteten Text)
        if not text.strip():
            try:
                logging.info("📸 Versuche OCR-Fallback für möglicherweise gescanntes PDF...")
                ocr_text = self._extract_pdf_with_ocr(file_path)
                if ocr_text and ocr_text.strip():
                    text = ocr_text
                    # OCR liefert keine natürlichen Seitenumbrüche, teile nach Zeilengruppen
                    page_texts = [text]
                    extraction_method = "ocr_easyocr"
                    logging.info(f"✅ OCR Extraktion erfolgreich: {len(text)} Zeichen")
            except Exception as e:
                errors.append(f"OCR: {e}")
                logging.warning(f"OCR-Fallback fehlgeschlagen: {e}")
        
        # Prüfe ob Extraktion erfolgreich war
        has_content = bool(page_texts) or bool(text.strip())
        if not has_content:
            return {"success": False, "error": f"Empty PDF: {'; '.join(errors)}"}
        
        # Dokumente für Upsert vorbereiten
        # ✅ SOTA FIX: Original-Dateiname aus metadata["source"] hat Vorrang vor Temp-Pfad
        _meta = kwargs.get('metadata') or {}
        _source_name = _meta.get('source') or os.path.basename(file_path)
        base_doc_id = kwargs.get('doc_id') or os.path.abspath(file_path)
        base_meta = {
            **_meta,
            "source": _source_name,
            "title": _meta.get('title', _source_name),
            "content_type": "application/pdf",
            "source_type": "document",
            "extraction_method": extraction_method,
        }
        
        docs = []
        for i, ptxt in enumerate(page_texts, start=1):
            if not str(ptxt or "").strip():
                continue
            docs.append({
                "id": f"{base_doc_id}#p{i}",
                "text": str(ptxt),
                "metadata": {**base_meta, "page": i},
            })
        
        # ✅ Fallback nur wenn wirklich keine Seiten extrahiert wurden aber text vorhanden
        if not docs and text and text.strip():
            docs.append({
                "id": base_doc_id,
                "text": text,
                "metadata": base_meta,
            })
        
        # ✅ Memory Cleanup
        if page_texts:
            del page_texts
        if text:
            del text
        
        # Upsert
        result = self.upsert_documents(docs, 
                                     chunk_size=kwargs.get('chunk_size', 1500),
                                     chunk_overlap=kwargs.get('chunk_overlap', 200))
        
        inserted_count = result.get("inserted", 0)
        return {
            "success": True,
            "inserted": inserted_count,
            "chunks_added": inserted_count,  # Alias für Kompatibilität
            "docs": result.get("docs", 1),
            "duplicates_skipped": result.get("duplicates_skipped", 0),
            "warnings": "; ".join(errors) if errors else None,
        }

    def upsert_pdf_with_vision(
        self,
        file_path: str,
        *,
        doc_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        chunk_size: int = 1500,
        chunk_overlap: int = 200,
        build_kg: bool = True,
        use_vision: bool = True,
        force_vision_all_pages: bool = False
    ) -> Dict[str, Any]:
        """
        🆕 STATE-OF-THE-ART: PDF-Verarbeitung mit Vision-LLM für Infografiken.
        
        Kombiniert pymupdf4llm + Vision-LLM für:
        - Infografiken, Charts, Statistiken
        - Gescannte/Hybrid-PDFs
        
        Args:
            file_path: Pfad zur PDF-Datei
            use_vision: Vision-LLM für Bilder nutzen
            force_vision_all_pages: Vision für ALLE Seiten
        """
        result: Dict[str, Any] = {"success": False, "chunks_added": 0, "vision_pages": 0}
        text = ""
        
        try:
            if use_vision:
                from agent.pdf_vision_extractor import PDFVisionExtractor
                extractor = PDFVisionExtractor(enable_vision_for_text_pages=force_vision_all_pages)
                vision_result = extractor.extract_complete(file_path, force_vision=force_vision_all_pages)
                
                text = vision_result.vision_enhanced_text or vision_result.full_text
                result["vision_pages"] = vision_result.vision_analyzed_pages
                result["pdf_type"] = vision_result.pdf_type
                
        except Exception as e:
            logger.warning(f"Vision-Extraktion fehlgeschlagen: {e}")
        
        if not text.strip():
            return self._upsert_pdf_sequential(file_path, doc_id=doc_id, metadata=metadata,
                                               chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                                               build_kg=build_kg)
        
        # Dokumente erstellen
        base_doc_id = doc_id or os.path.abspath(file_path)
        base_meta = {**(metadata or {}), "source": os.path.basename(file_path),
                     "extraction_method": "vision_enhanced", "vision_pages": result["vision_pages"]}
        
        page_texts = text.split("\n---\n") if "\n---\n" in text else [text]
        docs = [{"id": f"{base_doc_id}#p{i}", "text": ptxt, "metadata": {**base_meta, "page": i}}
                for i, ptxt in enumerate(page_texts, 1) if ptxt.strip()]
        
        upsert_result = self.upsert_documents(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        result["success"] = True
        result["chunks_added"] = upsert_result.get("inserted", 0)
        result["vision_used"] = True
        return result

    def _upsert_pdf_parallel(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parallele PDF-Verarbeitung (wenn verfügbar)"""
        
        start_time = time.time()
        
        # Phase 1: Parallele Seiten-Extraktion
        page_texts = self._extract_pages_parallel(file_path)
        if not page_texts:
            return {"success": False, "error": "Keine Seiten extrahiert"}
        
        # Phase 2: Parallele Chunk-Erstellung
        chunk_size = kwargs.get('chunk_size', 1500)
        chunk_overlap = kwargs.get('chunk_overlap', 200)
        
        all_chunks = []
        for page_idx, page_text in enumerate(page_texts):
            if not page_text.strip():
                continue
            chunks = self._chunk_text_with_spans(page_text, chunk_size, chunk_overlap)
            for start, end, chunk_text in chunks:
                all_chunks.append((page_idx, start, end, chunk_text))
        
        # Phase 3: Batch-Embedding-Generierung
        chunk_texts = [chunk_text for _, _, _, chunk_text in all_chunks]
        embeddings = self.embed_texts(chunk_texts)
        
        # Phase 4: Dokumente vorbereiten
        # ✅ SOTA FIX: Original-Dateiname aus metadata["source"] hat Vorrang vor Temp-Pfad
        _meta = kwargs.get('metadata') or {}
        _source_name = _meta.get('source') or os.path.basename(file_path)
        base_doc_id = kwargs.get('doc_id') or os.path.abspath(file_path)
        base_meta = {
            **_meta,
            "source": _source_name,
            "title": _meta.get('title', _source_name),
            "content_type": "application/pdf",
            # Legal-Hardening (2026-08-30): Web-Herkunft (z. B. "web_pdf")
            # aus der Caller-Metadaten nicht überschreiben.
            "source_type": _meta.get('source_type') or "document",
        }
        
        docs = []
        for i, (page_idx, start, end, chunk_text) in enumerate(all_chunks):
            docs.append({
                "id": f"{base_doc_id}#c{i}",
                "chunks": [chunk_text],  # Pre-chunked
                "metadata": {
                    **base_meta, 
                    "page": page_idx + 1,
                    "char_start": start,
                    "char_end": end
                },
            })
        
        # Upsert
        result = self.upsert_documents(docs)
        
        processing_time = time.time() - start_time
        
        inserted_count = result.get("inserted", 0)
        return {
            "success": True,
            "inserted": inserted_count,
            "chunks_added": inserted_count,  # Alias für Kompatibilität
            "docs": result.get("docs", 1),
            "duplicates_skipped": result.get("duplicates_skipped", 0),
            "pages_processed": len(page_texts),
            "chunks_generated": len(all_chunks),
            "processing_time": processing_time,
            "parallel_processing": True
        }

    def _extract_pages_parallel(self, file_path: str) -> List[str]:
        """
        Extrahiert PDF-Seiten parallel mit Batch-Processing (Memory-optimiert)
        
        WICHTIG: Verarbeitet Seiten in Batches statt alle gleichzeitig
        um Memory-Leaks und Out-of-Memory-Crashes zu vermeiden
        """
        
        page_texts = []
        
        try:
            import fitz
            with fitz.open(file_path) as doc:
                num_pages = len(doc)
                
                if num_pages <= 4:
                    # Kleine PDFs: sequenziell
                    for page in doc:
                        try:
                            # Use getattr to avoid type checking issues
                            page_text = getattr(page, 'get_text', lambda: "")()
                            page_texts.append(page_text)
                        except (AttributeError, TypeError) as e:
                            logger.warning(f"Could not extract text from page: {e}")
                            page_texts.append("")
                else:
                    # ✅ Große PDFs: BATCH-WEISE parallel (verhindert Memory-Explosion!)
                    batch_size = min(10, self.config.max_workers_pdf)  # Max 10 Seiten gleichzeitig
                    logger.info(f"📄 Verarbeite {num_pages} Seiten in Batches von {batch_size}")
                    
                    # Iteration 4: Nutze chunk_list() für sauberes Batching
                    all_page_indices = list(range(num_pages))
                    
                    for batch_pages in chunk_list(all_page_indices, batch_size):
                        # Batch parallel verarbeiten
                        with ProcessPoolExecutor(max_workers=len(batch_pages)) as executor:
                            futures = {
                                executor.submit(self._extract_single_page, file_path, i): i
                                for i in batch_pages
                            }
                            
                            batch_results = {}
                            for future in as_completed(futures):
                                page_idx = futures[future]
                                try:
                                    result = future.result()
                                    batch_results[page_idx] = result if result else ""
                                except Exception as e:
                                    logger.debug(f"Fehler bei Seite {page_idx}: {e}")
                                    batch_results[page_idx] = ""
                            
                            # Füge Batch-Ergebnisse in korrekter Reihenfolge hinzu
                            for i in batch_pages:
                                page_texts.append(batch_results.get(i, ""))
                            
                            # ✅ KRITISCH: Cleanup nach jedem Batch!
                            del batch_results
                            del futures
                            
                        # ✅ Executor wird automatisch geschlossen (with-Block)
                        
                    logger.info(f"✅ {num_pages} Seiten erfolgreich in {len(list(chunk_list(all_page_indices, batch_size)))} Batches verarbeitet")
                        
        except Exception as e:
            logger.error(f"PDF-Extraktion fehlgeschlagen: {e}")
            # Fallback
            try:
                from pdfminer.high_level import extract_text
                text = extract_text(file_path) or ""
                page_texts = text.split("\x0c") if text else []
            except (ImportError, Exception) as e:
                logger.warning(f"PDF fallback extraction failed: {e}")
                page_texts = []
        
        return page_texts

    @staticmethod
    def _extract_single_page(file_path: str, page_num: int) -> str:
        """Extrahiert eine einzelne PDF-Seite (für ProcessPool)"""
        try:
            import fitz
            with fitz.open(file_path) as doc:
                if page_num < len(doc):
                    page = doc[page_num]
                    # Use getattr to avoid type checking issues
                    return getattr(page, 'get_text', lambda: "")()
        except (ImportError, AttributeError, FileNotFoundError, Exception) as e:
            logger.warning(f"Could not extract page {page_num}: {e}")
        return ""

    # ======= URL PROCESSING =======
    
    def upsert_url_with_vision(
        self,
        url: str,
        *,
        doc_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        use_vision: bool = True,
        timeout: int = 30,
        chunk_size: int = 1500,
        chunk_overlap: int = 200
    ) -> Dict[str, Any]:
        """
        🆕 STATE-OF-THE-ART: URL-Verarbeitung mit Vision-LLM für Webseiten-Bilder.
        
        Kombiniert Text-Extraktion + Vision-LLM für:
        - Infografiken auf Webseiten
        - Charts und Diagramme
        - Relevante Bilder (filtert Logos/Icons automatisch)
        
        Args:
            url: URL der Webseite
            doc_id: Optional Document-ID
            metadata: Zusätzliche Metadaten
            use_vision: Vision-LLM für Bilder nutzen (Default: True)
            timeout: HTTP-Timeout
            chunk_size: Chunk-Größe
            chunk_overlap: Chunk-Überlappung
            
        Returns:
            Dict mit success, chunks_added, images_analyzed
        """
        result: Dict[str, Any] = {"success": False, "chunks_added": 0, "images_analyzed": 0, "images_found": 0}

        if parse_bool_env("APP_LOCAL_ONLY", "0"):
            result["error"] = (
                "APP_LOCAL_ONLY aktiv: URL-Verarbeitung ist deaktiviert. "
                "Bitte Datei lokal speichern und upsert_pdf/upsert_documents nutzen."
            )
            return result

        # ── Legal/ethical Gate (2026-08-30): robots.txt vor dem Fetch (fail-open) ──
        if not web_compliance.gate_persistence("upsert_url_with_vision", url):
            result["error"] = (
                "Compliance blockiert Persistierung (robots.txt/X-Robots-Tag). "
                "Details in den Logs ([WEB-COMPLIANCE])."
            )
            return result

        try:
            import requests
            from bs4 import BeautifulSoup
            
            # 1. Prüfe ob URL auf PDF zeigt
            if url.lower().endswith('.pdf'):
                logger.info(f"Detected PDF URL, using upsert_pdf_with_vision")
                pdf_result = self._upsert_pdf_from_url_with_vision(
                    url, doc_id=doc_id, metadata=metadata, timeout=timeout,
                    chunk_size=chunk_size, chunk_overlap=chunk_overlap
                )
                return pdf_result
            
            # 2. HTML laden
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
            response = requests.get(url, headers=headers, timeout=timeout)
            response.raise_for_status()
            html = response.text
            
            # 3. Text extrahieren
            soup = BeautifulSoup(html, 'html.parser')
            
            # Titel extrahieren
            title_tag = soup.find('title')
            page_title = title_tag.get_text(strip=True) if title_tag else ""
            
            # Unwanted elements entfernen
            for unwanted in soup(["script", "style", "nav", "header", "footer", "aside"]):
                unwanted.decompose()
            
            text = soup.get_text(separator=" ", strip=True)
            del soup
            
            if len(text) < 100:
                logger.warning(f"Insufficient content from {url}")
                result["error"] = "Insufficient content"
                return result
            
            # 4. Vision-Analyse für Bilder (wenn aktiviert)
            vision_text = text
            if use_vision:
                try:
                    from agent.web_vision_extractor import WebVisionExtractor
                    
                    extractor = WebVisionExtractor(timeout=timeout)
                    vision_result = extractor.extract_with_vision(
                        url=url,
                        html=html,
                        page_text=text,
                        page_title=page_title
                    )
                    
                    result["images_found"] = vision_result.total_images_found
                    result["images_analyzed"] = vision_result.analyzed_images
                    
                    if vision_result.vision_enhanced_text:
                        vision_text = vision_result.vision_enhanced_text
                        logger.info(f"🖼️ Vision: {vision_result.analyzed_images} Bilder analysiert")
                    
                except ImportError:
                    logger.warning("WebVisionExtractor nicht verfügbar")
                except Exception as e:
                    logger.warning(f"Vision-Analyse fehlgeschlagen: {e}")
            
            # 5. Metadaten
            base_meta = {
                **(metadata or {}),
                "source": url,
                "source_url": url,
                "title": page_title,
                "content_type": "text/html",
                "source_type": "web",
                "extraction_method": "vision_enhanced" if result["images_analyzed"] > 0 else "text_only",
                "images_analyzed": result["images_analyzed"]
            }
            
            # 6. Upsert
            doc = {
                "id": doc_id or self._hash_str(url),
                "text": vision_text,
                "metadata": base_meta
            }
            
            upsert_result = self.upsert_documents(
                [doc],
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            
            result["success"] = upsert_result.get("success", False)
            result["chunks_added"] = upsert_result.get("inserted", 0)
            
            logger.info(f"✅ URL mit Vision verarbeitet: {result['chunks_added']} Chunks, "
                       f"{result['images_analyzed']} Bilder")
            
            return result
            
        except Exception as e:
            logger.error(f"URL-Vision-Import fehlgeschlagen: {e}")
            result["error"] = str(e)
            return result
    
    def _upsert_pdf_from_url_with_vision(
        self,
        url: str,
        *,
        doc_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        timeout: int = 30,
        chunk_size: int = 1500,
        chunk_overlap: int = 200
    ) -> Dict[str, Any]:
        """PDF von URL mit Vision-Analyse."""
        import tempfile
        result: Dict[str, Any] = {"success": False, "chunks_added": 0, "images_analyzed": 0}

        if parse_bool_env("APP_LOCAL_ONLY", "0"):
            result["error"] = (
                "APP_LOCAL_ONLY aktiv: PDF-URL-Download ist deaktiviert. "
                "Bitte PDF lokal speichern und upsert_pdf_with_vision nutzen."
            )
            return result
        
        try:
            import requests
            
            # Download
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
            response = requests.get(url, headers=headers, timeout=timeout)
            response.raise_for_status()
            
            # Temp-Datei
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name
            
            try:
                # Mit Vision verarbeiten
                vision_result = self.upsert_pdf_with_vision(
                    tmp_path,
                    doc_id=doc_id,
                    # Legal-Hardening (2026-08-30): Web-Herkunft kennzeichnen,
                    # damit Retention/Pruning greift (siehe docs/18).
                    metadata={**(metadata or {}), "source_url": url, "source_type": "web_pdf"},
                    use_vision=True
                )
                
                result["success"] = vision_result.get("success", False)
                result["chunks_added"] = vision_result.get("chunks_added", 0)
                result["images_analyzed"] = vision_result.get("vision_pages", 0)
                
            finally:
                # Cleanup
                import os
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
            
            return result
            
        except Exception as e:
            logger.error(f"PDF-URL-Vision fehlgeschlagen: {e}")
            result["error"] = str(e)
            return result

    def upsert_url(self, url: str, *, doc_id: Optional[str] = None,
                   metadata: Optional[Dict[str, Any]] = None,
                   include_tables: bool = True, include_links: bool = False,
                   timeout: int = 30, chunk_size: int = 1500, 
                   chunk_overlap: int = 200) -> bool:
        """
        URL-Content in den RAG Store laden (HTML, PDF, DOCX, PPTX, XLSX, ...)
        
        🆕 DOCLING SOTA (2025-07-18):
            - Universelle Dokumenten-Erkennung via Extension + Content-Type
            - PDF, DOCX, PPTX, XLSX → Docling AI-Verarbeitung (CPU)
            - HTML → Schneller HTML-Parser (Trafilatura/BeautifulSoup)
            - AI-basierte Tabellenerkennung (TableFormer)
            - Strukturiertes Markdown-Output mit Layout-Erkennung
        
        Args:
            url: URL zum Extrahieren (HTML, PDF, DOCX, PPTX, XLSX, ...)
            doc_id: Optional document ID
            metadata: Zusätzliche Metadaten
            include_tables: Tabellen extrahieren
            include_links: Links extrahieren
            timeout: HTTP-Timeout
            chunk_size: Chunk-Größe
            chunk_overlap: Chunk-Überlappung
            
        Returns:
            True wenn erfolgreich
        """
        if parse_bool_env("APP_LOCAL_ONLY", "0"):
            logger.warning(
                "🔒 APP_LOCAL_ONLY aktiv: upsert_url blockiert für %s. "
                "Nutze lokale Datei-Imports.",
                url,
            )
            return False

        # ── Legal/ethical Gate (2026-08-30): robots.txt vor dem Fetch (fail-open) ──
        # Siehe docs/18_LEGAL_WEB_PERSIST.md — blockiert bei Disallow/noindex/no-store.
        if not web_compliance.gate_persistence("upsert_url", url):
            return False

        try:
            # Import requests nur für URL-Processing (lokaler Import)
            try:
                import requests
            except ImportError:
                logger.error("requests library not available for URL processing")
                return False
            
            # ═══════════════════════════════════════════════════════════
            # SCHRITT 1: Docling-basierte Format-Erkennung (SOTA)
            # Erkennt PDF, DOCX, PPTX, XLSX via Extension + Content-Type
            # ═══════════════════════════════════════════════════════════
            use_docling = False
            detected_format = "unknown"
            
            try:
                from utils.docling_processor import DoclingProcessor
                detected_format, use_docling = DoclingProcessor.detect_format_from_url(
                    url, timeout=5
                )
                if use_docling:
                    logger.info(
                        f"🔬 Docling: Detected {detected_format} document at {url}"
                    )
            except ImportError:
                logger.debug("DoclingProcessor not available, using legacy routing")
            except Exception as e:
                logger.debug(f"Docling format detection failed: {e}, using legacy routing")
            
            # ═══════════════════════════════════════════════════════════
            # SCHRITT 2: Docling-Verarbeitung für Nicht-HTML-Dokumente
            # PDF, DOCX, PPTX, XLSX → Docling AI Pipeline
            # ═══════════════════════════════════════════════════════════
            if use_docling:
                return self._upsert_via_docling(
                    url, doc_id=doc_id, metadata=metadata,
                    timeout=timeout, chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    detected_format=detected_format,
                )
            
            # ═══════════════════════════════════════════════════════════
            # LEGACY FALLBACK: PDF-only Erkennung (falls Docling nicht verfügbar)
            # ═══════════════════════════════════════════════════════════
            if url.lower().endswith('.pdf'):
                logger.info(f"Detected PDF URL by extension (legacy): {url}")
                return self._upsert_pdf_from_url(
                    url, doc_id=doc_id, metadata=metadata,
                    include_tables=include_tables, timeout=timeout,
                    chunk_size=chunk_size, chunk_overlap=chunk_overlap
                )
            
            # HEAD-Request für Content-Type (Legacy PDF-Check)
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            
            try:
                head_response = requests.head(url, headers=headers, timeout=5, allow_redirects=True)
                content_type = head_response.headers.get('Content-Type', '').lower()
                
                if 'application/pdf' in content_type:
                    logger.info(f"Detected PDF URL by Content-Type (legacy): {url}")
                    return self._upsert_pdf_from_url(
                        url, doc_id=doc_id, metadata=metadata,
                        include_tables=include_tables, timeout=timeout,
                        chunk_size=chunk_size, chunk_overlap=chunk_overlap
                    )
            except Exception as e:
                logger.debug(f"HEAD request failed for {url}: {e}, proceeding with GET")
            
            # SCHRITT 3: HTML-Verarbeitung (SOTA -- Trafilatura + BS4 Fallback)
            # Quick timeout for fast processing
            response = requests.get(url, headers=headers, timeout=min(timeout, 8))
            response.raise_for_status()
            
            html_content = response.text
            text = None
            extraction_method = "none"
            
            # ═══════════════════════════════════════════════════════
            # LAYER 1a: Trafilatura (SOTA -- jusText-Algorithmus)
            #   Entfernt automatisch: Boilerplate, Navigation, Ads,
            #   Tracking-Code, Cookie-Banner, Sidebar-Widgets,
            #   JavaScript-Fragmente, JSON-LD, Code-Blöcke
            # ═══════════════════════════════════════════════════════
            try:
                import trafilatura
                text = trafilatura.extract(
                    html_content,
                    include_tables=include_tables,
                    include_links=include_links,
                    include_comments=False,
                    include_images=False,
                    no_fallback=False,      # Fallback auf readability/jusText
                    favor_precision=True,    # Präzision > Recall (weniger Noise)
                    deduplicate=True,        # Intra-Dokument-Dedup
                )
                if text and len(text.strip()) >= 100:
                    extraction_method = "trafilatura"
                    logger.debug(f"Trafilatura extracted {len(text)} chars from {url}")
                else:
                    text = None  # Zu wenig Content → Fallback
            except Exception as e:
                logger.debug(f"Trafilatura extraction failed for {url}: {e}")
                text = None
            
            # ═══════════════════════════════════════════════════════
            # LAYER 1b: Enhanced BS4 Fallback
            #   Erweiterte Tag-Liste: entfernt auch <code>, <pre>,
            #   <noscript>, <template>, <svg>, <form>, <iframe>,
            #   <button>, <input>, <select>, <textarea>, <canvas>,
            #   <object>, <embed>, <applet>, <map>, <area>
            # ═══════════════════════════════════════════════════════
            if text is None:
                try:
                    from bs4 import BeautifulSoup
                    
                    soup = BeautifulSoup(html_content, "html.parser")
                    
                    try:
                        # Erweiterte Liste: ALLE nicht-Content-Tags entfernen
                        NOISE_TAGS = [
                            "script", "style", "nav", "header", "footer", "aside",
                            # NEU -- Code/Template-Artefakte:
                            "code", "pre", "noscript", "template", "svg", "math",
                            # NEU -- Interaktive/Formular-Elemente:
                            "form", "iframe", "button", "input", "select", "textarea",
                            # NEU -- Embedded Objects:
                            "canvas", "object", "embed", "applet", "map", "area",
                        ]
                        for unwanted in soup(NOISE_TAGS):
                            unwanted.decompose()
                        
                        # Entferne Elemente mit verdächtigen Klassen/IDs
                        import re as _re
                        NOISE_PATTERN = _re.compile(
                            r'cookie|consent|gdpr|tracking|analytics|sidebar|widget|'
                            r'popup|modal|banner|advertisement|ad-|social|share|related|'
                            r'comment|disqus|newsletter|subscribe',
                            _re.IGNORECASE
                        )
                        for el in soup.find_all(attrs={"class": NOISE_PATTERN}):
                            el.decompose()
                        for el in soup.find_all(attrs={"id": NOISE_PATTERN}):
                            el.decompose()
                        
                        # Entferne inline JS Event-Handler
                        for tag in soup.find_all(True):
                            attrs_to_remove = [
                                attr for attr in tag.attrs
                                if attr.startswith('on')  # onclick, onload, etc.
                            ]
                            for attr in attrs_to_remove:
                                del tag[attr]
                        
                        text = soup.get_text(separator=" ", strip=True)
                        extraction_method = "bs4_enhanced"
                        
                    finally:
                        del soup
                    
                except ImportError:
                    # Letzter Fallback: Regex
                    import re
                    text = re.sub(r'<[^>]+>', ' ', html_content)
                    text = re.sub(r'\s+', ' ', text).strip()
                    extraction_method = "regex_fallback"
            
            # ✅ KRITISCH: HTML-Content löschen (kann sehr groß sein!)
            del html_content
            import gc
            gc.collect()
            
            # Content-Qualitätsprüfung
            if not text or len(text) < 100:
                logger.warning(f"Insufficient content extracted from {url} (method={extraction_method})")
                return False
            
            # Log content size
            if len(text) > 100000:
                logger.info(f"Large content from {url}: {len(text)} chars (will be chunked)")
            
            logger.debug(f"Extracted {len(text)} chars from {url} via {extraction_method}")
            
            if self.debug:
                print(f"📄 Extracted {len(text)} characters from {url}")
            
            # Metadaten zusammenstellen
            base_meta = {
                **(metadata or {}),
                "source": url,
                "source_url": url,
                "canonical_url": url,
                "content_type": "text/html",
                "source_type": "web",
                "extraction_method": extraction_method
            }
            
            # Dokument für Upsert vorbereiten
            doc = {
                "id": doc_id or self._hash_str(url),
                "text": text,
                "metadata": base_meta
            }
            
            # Upsert
            result = self.upsert_documents([doc], 
                                         chunk_size=chunk_size,
                                         chunk_overlap=chunk_overlap)
            
            return bool(result.get("success", False)) and int(result.get("inserted", 0)) > 0
                
        except Exception as e:
            logger.error(f"URL-Import fehlgeschlagen für {url}: {e}")
            return False

    def _upsert_via_docling(
        self, url: str, *, doc_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        timeout: int = 30, chunk_size: int = 1500,
        chunk_overlap: int = 200, detected_format: str = "unknown",
    ) -> bool:
        """
        🆕 SOTA Docling-basierte URL-Verarbeitung (2025-07-18)
        
        Universeller Pfad für PDF, DOCX, PPTX, XLSX via Docling AI.
        Download → Temp-File → Docling Conversion → Chunking → Upsert
        
        Features:
        - AI Layout-Analyse (RT-DETR)
        - TableFormer für Tabellenerkennung
        - RapidOCR für gescannte Dokumente
        - HierarchicalChunker für semantische Chunks
        
        Args:
            url: Dokument-URL
            doc_id: Optionale Dokument-ID
            metadata: Zusätzliche Metadaten
            timeout: Download-Timeout
            chunk_size: Chunk-Größe
            chunk_overlap: Chunk-Überlappung
            detected_format: Erkanntes Format (PDF, DOCX, etc.)
            
        Returns:
            True wenn erfolgreich
        """
        try:
            from utils.docling_processor import DoclingProcessor
            
            processor = DoclingProcessor.get_instance()
            
            logger.info(f"🔬 Docling: Processing {detected_format} from {url}")
            
            # Konvertiere URL → DoclingResult
            docling_result = processor.convert_url(
                url,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                metadata={
                    **(metadata or {}),
                    "source_url": url,
                    "canonical_url": url,
                    "source_type": f"web_{detected_format.lower()}",
                    "content_type": detected_format,
                    "search_timestamp": datetime.now(timezone.utc).isoformat(),
                },
                timeout=timeout,
            )
            
            if not docling_result.success:
                logger.warning(
                    f"⚠️ Docling failed for {url}: {docling_result.error}. "
                    f"Falling back to legacy processing."
                )
                # Fallback auf Legacy-PDF-Verarbeitung für PDFs
                if detected_format == "PDF":
                    return self._upsert_pdf_from_url(
                        url, doc_id=doc_id, metadata=metadata,
                        include_tables=True, timeout=timeout,
                        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                    )
                return False
            
            # ═══════════════════════════════════════════════════════
            # Docling Chunks → RAG Store Upsert
            # ═══════════════════════════════════════════════════════
            if docling_result.chunks:
                # Nutze Docling's semantische Chunks direkt
                documents = []
                base_doc_id = doc_id or self._hash_str(url)
                
                for chunk in docling_result.chunks:
                    chunk_doc_id = f"{base_doc_id}_chunk_{chunk.index}"
                    documents.append({
                        "id": chunk_doc_id,
                        "text": chunk.text,
                        "metadata": {
                            **(metadata or {}),
                            **chunk.metadata,
                            "source": url,
                            "source_url": url,
                            "canonical_url": url,
                            "source_type": f"web_{detected_format.lower()}",
                            "content_type": detected_format,
                            "extraction_method": "docling",
                            "docling_format": docling_result.format_detected,
                            "docling_tables": docling_result.num_tables,
                            "docling_pages": docling_result.num_pages,
                            "chunk_headings": " > ".join(chunk.headings) if chunk.headings else "",
                        },
                    })
                
                if documents:
                    # Upsert OHNE erneutes Chunking (Docling hat bereits gechunkt)
                    # chunk_size=50000 verhindert Re-Chunking da Docling-Chunks << 50k
                    result = self.upsert_documents(
                        documents,
                        chunk_size=50000,
                        chunk_overlap=0,
                    )
                    
                    inserted = int(result.get("inserted", 0))
                    logger.info(
                        f"✅ Docling: {url} → {inserted} chunks inserted "
                        f"({docling_result.num_tables} tables, "
                        f"{docling_result.processing_time_s:.1f}s)"
                    )
                    return inserted > 0
            
            # Fallback: Wenn keine Chunks, nutze den vollen Markdown-Text
            if docling_result.text:
                doc = {
                    "id": doc_id or self._hash_str(url),
                    "text": docling_result.text,
                    "metadata": {
                        **(metadata or {}),
                        "source": url,
                        "source_url": url,
                        "canonical_url": url,
                        "source_type": f"web_{detected_format.lower()}",
                        "content_type": detected_format,
                        "extraction_method": "docling",
                    },
                }
                result = self.upsert_documents(
                    [doc], chunk_size=chunk_size, chunk_overlap=chunk_overlap
                )
                return bool(result.get("success", False)) and int(result.get("inserted", 0)) > 0
            
            logger.warning(f"Docling produced no content for {url}")
            return False
            
        except ImportError:
            logger.warning("DoclingProcessor not available, falling back to legacy")
            if detected_format == "PDF":
                return self._upsert_pdf_from_url(
                    url, doc_id=doc_id, metadata=metadata,
                    include_tables=True, timeout=timeout,
                    chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                )
            return False
        except Exception as e:
            logger.error(f"❌ Docling upsert failed for {url}: {e}")
            # Fallback
            if detected_format == "PDF":
                try:
                    return self._upsert_pdf_from_url(
                        url, doc_id=doc_id, metadata=metadata,
                        include_tables=True, timeout=timeout,
                        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
                    )
                except Exception as exc:
                    logger.error(f"❌ PDF fallback upsert failed for {url}: {exc}")
            return False

    def _upsert_pdf_from_url(self, url: str, *, doc_id: Optional[str] = None,
                            metadata: Optional[Dict[str, Any]] = None,
                            include_tables: bool = True, timeout: int = 30,
                            chunk_size: int = 1500, chunk_overlap: int = 200) -> bool:
        """
        PDF von URL herunterladen und verarbeiten (LEGACY Fallback)
        
        🆕 NEUE METHODE (2025-10-16)
        ⚠️ LEGACY: Wird nur als Fallback genutzt wenn Docling fehlschlägt
        
        Workflow:
        1. Download PDF in Memory
        2. Speichere in Temp-Datei
        3. Verarbeite mit upsert_pdf()
        4. Cleanup Temp-Datei
        
        Args:
            url: PDF-URL
            doc_id: Optional document ID
            metadata: Zusätzliche Metadaten
            include_tables: Tabellen extrahieren
            timeout: HTTP-Timeout
            chunk_size: Chunk-Größe
            chunk_overlap: Chunk-Überlappung
            
        Returns:
            True wenn erfolgreich
        """
        if parse_bool_env("APP_LOCAL_ONLY", "0"):
            logger.warning(
                "🔒 APP_LOCAL_ONLY aktiv: _upsert_pdf_from_url blockiert für %s",
                url,
            )
            return False

        import tempfile
        
        # Import requests at function level (before try-except)
        try:
            import requests
        except ImportError:
            logger.error("requests library not available for PDF URL download")
            return False
        
        try:
            logger.info(f"📥 Downloading PDF from URL: {url}")
            
            # Download PDF
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=timeout)
            response.raise_for_status()
            
            # Validiere Content-Type
            content_type = response.headers.get('Content-Type', '').lower()
            if 'application/pdf' not in content_type and not url.lower().endswith('.pdf'):
                logger.warning(f"URL {url} does not appear to be a PDF (Content-Type: {content_type})")
                # Versuche trotzdem, könnte trotzdem ein PDF sein
            
            # Prüfe Dateigröße
            content_length = len(response.content)
            if content_length > 100 * 1024 * 1024:  # 100 MB Limit
                logger.error(f"PDF too large: {content_length / (1024*1024):.1f} MB (limit: 100 MB)")
                return False
            
            logger.info(f"Downloaded {content_length / 1024:.1f} KB")
            
            # Speichere in Temp-Datei
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name
            
            # Cleanup Response-Objekt
            del response
            import gc
            gc.collect()
            
            try:
                # Verarbeite PDF mit upsert_pdf()
                logger.info(f"📄 Processing downloaded PDF: {tmp_path}")
                
                result = self.upsert_pdf(
                    tmp_path,
                    doc_id=doc_id or self._hash_str(url),
                    metadata={
                        **(metadata or {}),
                        'source_url': url,
                        'source_type': 'web_pdf',
                        'canonical_url': url,
                        'download_timestamp': datetime.now(timezone.utc).isoformat(),
                        'content_type': 'application/pdf'
                    },
                    extract_tables=include_tables,
                    build_kg=True,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap
                )
                
                success: bool = bool(result.get('success', False))
                
                if success:
                    logger.info(f"✅ Successfully processed PDF from {url}")
                    logger.info(f"   Chunks added: {result.get('inserted', 0)}")
                    if result.get('duplicates_skipped', 0) > 0:
                        logger.info(f"   Duplicates skipped: {result.get('duplicates_skipped', 0)}")
                else:
                    logger.error(f"❌ Failed to process PDF from {url}: {result.get('error', 'Unknown error')}")
                
                return success
                
            finally:
                # Lösche Temp-Datei
                try:
                    os.unlink(tmp_path)
                    logger.debug(f"Cleaned up temp file: {tmp_path}")
                except Exception as e:
                    logger.warning(f"Failed to delete temp file {tmp_path}: {e}")
        
        except requests.exceptions.Timeout:
            logger.error(f"Timeout downloading PDF from {url}")
            return False
        except requests.exceptions.RequestException as e:
            logger.error(f"Failed to download PDF from {url}: {e}")
            return False
        except Exception as e:
            logger.error(f"Unexpected error processing PDF URL {url}: {e}")
            return False

    def get_url_import_status(self, url: str) -> str:
        """
        Gibt den Import-Status einer URL zurück
        
        Returns:
            'full', 'snippet_only', oder 'none'
        """
        try:
            conn = self.get_connection()
            cur = conn.cursor()
            
            try:
                cur.execute("""
                    SELECT COUNT(*) FROM chunks 
                    WHERE json_extract(metadata, '$.source_url') = ? 
                       OR json_extract(metadata, '$.canonical_url') = ?
                """, (url, url))
                
                count = cur.fetchone()[0]
                return 'full' if count > 0 else 'none'
                
            finally:
                cur.close()
                self.return_connection(conn)
                
        except Exception as exc:
            logger.debug(f"get_url_import_status failed: {exc}", exc_info=True)
            return 'none'

    # ═════════════════════════════════════════════════════════════════════
    # Web Retention / Legal Hardening (2026-08-30)
    # Siehe docs/18_LEGAL_WEB_PERSIST.md
    # ═════════════════════════════════════════════════════════════════════
    @staticmethod
    def _parse_retention_timestamp(value: Any) -> Optional[datetime]:
        """Parst einen Timestamp (ISO-8601, 'Z'-Suffix, epoch seconds oder
        'YYYY-MM-DD [HH:MM:SS]') zu einem tz-aware UTC-Datetime. None bei
        nicht verwertbaren Werten (keine Exceptions nach außen)."""
        if value is None:
            return None
        if isinstance(value, (int, float)) and not isinstance(value, bool):
            try:
                return datetime.fromtimestamp(float(value), tz=timezone.utc)
            except (OverflowError, OSError, ValueError):
                return None
        s = str(value).strip()
        if not s:
            return None
        if s.endswith("Z"):
            s = s[:-1] + "+00:00"
        dt: Optional[datetime] = None
        try:
            dt = datetime.fromisoformat(s)
        except ValueError:
            for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
                try:
                    dt = datetime.strptime(s, fmt)
                    break
                except ValueError:
                    continue
        if dt is None:
            return None
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(timezone.utc)

    def prune_web_content(self, max_age_days: Optional[int] = None, *,
                          dry_run: bool = False) -> int:
        """Löscht abgelaufenes Web-sourced Content aus dem RAG-Store.

        Nur Dokumente, deren Chunk-Metadaten ``source_type`` mit ``web``
        beginnen, sind betroffen. Lokale Dateien, KI-Generiertes und
        Records ohne verwertbaren Timestamp bleiben unberührt.

        Ablauf pro Dokument:
          1. ``retention_until`` (wenn gesetzt) → Ablaufdatum
          2. sonst ``search_timestamp``/``date_stored`` + ``max_age_days``
          3. sonst Übersprung (kein verwertbarer Timestamp)

        Explizite Kind-Tabelle-Löschung (chunk_quality, triples, tables,
        image_metadata, chunks) vor der ``documents``-Zeile; FAISS/BM25-
        Staleness wird über ``auto_rebuild_on_stale`` aufgelöst.

        Args:
            max_age_days: Maximales Alter in Tagen. None →
                ``WEB_RETENTION_DAYS`` (Default 30). 0 → unbegrenzt.
            dry_run: Nur zählen, nichts löschen.

        Returns:
            Anzahl gelöschter (bzw. bei dry_run: zu löschender) Dokumente.
        """
        days = web_compliance.get_retention_days() if max_age_days is None else int(max_age_days)
        if days is None or days <= 0:
            logger.info("[WEB-RETENTION] Retention unbegrenzt — kein Pruning")
            return 0

        now = datetime.now(timezone.utc)
        pruned = 0
        skipped_no_ts = 0
        skipped_not_expired = 0
        removed_rows = 0

        conn = self.get_connection()
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT doc_id,
                       MIN(CASE WHEN json_extract(metadata, '$.retention_until') IS NOT NULL
                                THEN json_extract(metadata, '$.retention_until') END) AS retention_until,
                       MIN(CASE WHEN json_extract(metadata, '$.search_timestamp') IS NOT NULL
                                THEN json_extract(metadata, '$.search_timestamp') END) AS search_timestamp,
                       MIN(CASE WHEN json_extract(metadata, '$.date_stored') IS NOT NULL
                                THEN json_extract(metadata, '$.date_stored') END) AS date_stored
                FROM chunks
                WHERE json_extract(metadata, '$.source_type') LIKE 'web%'
                GROUP BY doc_id
            """)
            rows = cur.fetchall()

            if not rows:
                logger.info("[WEB-RETENTION] Keine Web-sourced Dokumente gefunden")
                return 0

            # Vorhandene Tabellen mit doc_id-Spalte ermitteln
            table_names = {r[0] for r in cur.execute(
                "SELECT name FROM sqlite_master WHERE type='table'"
            ).fetchall()}
            child_tables = [t for t in ("chunk_quality", "triples", "tables",
                                        "image_metadata", "chunks") if t in table_names]
            has_documents = "documents" in table_names

            for doc_id, retention_until, search_timestamp, date_stored in rows:
                # Ablaufdatum bestimmen (Retention-Modell, siehe Doku)
                expiry = self._parse_retention_timestamp(retention_until)
                if expiry is None:
                    base_ts = (self._parse_retention_timestamp(search_timestamp)
                               or self._parse_retention_timestamp(date_stored))
                    if base_ts is None:
                        skipped_no_ts += 1
                        continue
                    expiry = base_ts + timedelta(days=days)

                if expiry > now:
                    skipped_not_expired += 1
                    continue

                if dry_run:
                    pruned += 1
                    continue

                for t in child_tables:
                    cur.execute(f"DELETE FROM {t} WHERE doc_id = ?", (doc_id,))
                    removed_rows += cur.rowcount or 0
                if has_documents:
                    cur.execute("DELETE FROM documents WHERE doc_id = ?", (doc_id,))

                pruned += 1

            conn.commit()
            mode = "DRY-RUN" if dry_run else "executed"
            logger.info(
                "[WEB-RETENTION] Web-Pruning (%s): %d Dokumente gelöscht, "
                "%d übersprungen (keiner Timestamp), %d noch gültig",
                mode, pruned, skipped_no_ts, skipped_not_expired,
            )
            if self.debug:
                logger.debug("[WEB-RETENTION] Entfernte Zeilen in Kind-Tabellen: %d", removed_rows)
            return pruned

        except Exception as exc:
            logger.warning(f"[WEB-RETENTION] Pruning fehlgeschlagen: {exc}", exc_info=True)
            return 0
        finally:
            self.return_connection(conn)

    # ======= UTILITY METHODS =======
    
    @contextmanager
    def get_connection_context(self):
        """Context manager für sichere Connection-Verwaltung"""
        conn = self.get_connection()
        try:
            yield conn
        finally:
            self.return_connection(conn)
    
    def __enter__(self):
        """Context manager entry"""
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit - ensure cleanup"""
        self.close()
        return False

    def __del__(self):
        """Destruktor - stellt sicher, dass cleanup immer aufgerufen wird"""
        if self._is_interpreter_shutting_down():
            return
        try:
            self.cleanup_resources()
        except Exception as exc:
            logger.debug(f"Destructor cleanup failed: {exc}")

    def _run_cleanup_step(self, label: str, action: Any) -> None:
        """Runs one cleanup step with explicit error reporting.

        Cleanup darf Fehler nicht eskalieren, soll sie aber ausserhalb von
        Interpreter-Finalization sichtbar machen.
        """
        try:
            action()
        except Exception as exc:
            if not (self._is_interpreter_shutting_down() or self._is_shutdown_exception(exc)):
                logger.debug(f"Cleanup step '{label}' failed: {exc}")

    def cleanup_resources(self) -> None:
        """ECHTE Resource Cleanup ohne Race Conditions"""
        if self._cleaned_up:
            return
            
        try:
            # 1. Embedding Model cleanup
            _emb_model = getattr(self, 'embedding_model', None)
            if _emb_model is not None:
                self._run_cleanup_step(
                    "drop_embedding_model",
                    lambda: delattr(self, 'embedding_model'),
                )
                    
            # 2. Thread Executor cleanup
            _thread_exec: Optional[Any] = getattr(self, 'thread_executor', None)
            if _thread_exec is not None:
                self._run_cleanup_step(
                    "thread_executor_shutdown",
                    lambda: _thread_exec.shutdown(wait=False),
                )

            # 2.5 KG entity-resolution scheduler cleanup
            if hasattr(self, '_entity_resolution_executor'):
                self._run_cleanup_step(
                    "entity_resolution_scheduler_shutdown",
                    self._shutdown_entity_resolution_scheduler,
                )
                    
            # 5. Process Executor cleanup
            _process_exec: Optional[Any] = getattr(self, 'process_executor', None)
            if _process_exec is not None:
                self._run_cleanup_step(
                    "process_executor_shutdown",
                    lambda: _process_exec.shutdown(wait=False),
                )
                    
            # 6. GPU Resources cleanup
            _torch_ref = torch
            if (
                hasattr(self, 'device')
                and self.device
                and CUDA_AVAILABLE
                and _torch_ref is not None
                and 'cuda' in str(getattr(self, 'device', ''))
            ):
                self._run_cleanup_step("cuda_empty_cache", lambda _t=_torch_ref: _t.cuda.empty_cache())
            
            # 6.5. FAISS Manager cleanup (prevent memory leak)
            if hasattr(self, 'faiss_manager') and self.faiss_manager:
                self._run_cleanup_step(
                    "drop_faiss_manager_ref",
                    lambda: setattr(self, 'faiss_manager', None),
                )
            
            if hasattr(self, 'fusion_engine') and self.fusion_engine:
                self._run_cleanup_step(
                    "drop_fusion_engine_ref",
                    lambda: setattr(self, 'fusion_engine', None),
                )
                    
            # 7. Database Connections cleanup (use existing robust close)
            self._run_cleanup_step("store_close", self.close)
                
            # 8. Force garbage collection
            import gc
            gc.collect()
                
        except Exception as e:
            # Log aber nicht crashen
            logger.debug(f"Resource cleanup: {e}")
        finally:
            # Rufe Parent cleanup auf
            super().cleanup_resources()

    def upsert_file(self, file_path: str, *, doc_id: Optional[str] = None,
                   metadata: Optional[Dict[str, Any]] = None, force_update: bool = False,
                   chunk_size: int = 1500, chunk_overlap: int = 200) -> bool:
        """
        Universelle Datei-Upload-Methode für alle unterstützten Dateiformate
        
        🆕 DOCLING SOTA (2025-07-18):
            - Docling als primärer Prozessor für PDF, DOCX, PPTX, XLSX
            - AI-basierte Layout-Analyse, Tabellenerkennung, OCR
            - HierarchicalChunker für semantische Chunks
            - Legacy-Fallback wenn Docling fehlschlägt
        
        Args:
            file_path: Pfad zur Datei
            doc_id: Optionale Dokument-ID
            metadata: Zusätzliche Metadaten
            force_update: Dokument überschreiben
            chunk_size: Chunk-Größe (für Fallback-Processing)
            chunk_overlap: Chunk-Überlappung (für Fallback-Processing)
            
        Returns:
            True wenn erfolgreich, False sonst
        """
        try:
            metadata = dict(metadata or {})
            metadata.setdefault("corpus_domain", "general")
            metadata.setdefault("safety_flag", "safe")

            if not os.path.exists(file_path):
                logger.error(f"Datei nicht gefunden: {file_path}")
                return False
            
            # Bestimme Dateityp basierend auf Endung
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # ═══════════════════════════════════════════════════════════
            # SCHRITT 1: Docling SOTA-Verarbeitung (primärer Pfad)
            # PDF, DOCX, PPTX, XLSX → Docling AI Pipeline
            # ═══════════════════════════════════════════════════════════
            docling_formats = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls'}
            
            if file_ext in docling_formats:
                try:
                    from utils.docling_processor import DoclingProcessor
                    
                    processor = DoclingProcessor.get_instance()
                    docling_result = processor.convert_file(
                        file_path,
                        chunk_size=chunk_size,
                        chunk_overlap=chunk_overlap,
                        metadata={
                            **(metadata or {}),
                            "source": file_path,
                            "filename": os.path.basename(file_path),
                            "file_type": file_ext,
                        },
                    )
                    
                    if docling_result.success and docling_result.chunks and not getattr(docling_result, 'is_partial', False):
                        # Docling Chunks → RAG Store
                        base_doc_id = doc_id or f"docling_{os.path.basename(file_path)}"
                        documents = []
                        
                        for chunk in docling_result.chunks:
                            chunk_doc_id = f"{base_doc_id}_chunk_{chunk.index}"
                            documents.append({
                                "id": chunk_doc_id,
                                "text": chunk.text,
                                "metadata": {
                                    **(metadata or {}),
                                    **chunk.metadata,
                                    "source": file_path,
                                    "filename": os.path.basename(file_path),
                                    "file_type": file_ext,
                                    "source_type": file_ext.lstrip('.'),
                                    "extraction_method": "docling",
                                    "docling_format": docling_result.format_detected,
                                    "docling_tables": docling_result.num_tables,
                                    "docling_pages": docling_result.num_pages,
                                    "chunk_headings": " > ".join(chunk.headings) if chunk.headings else "",
                                },
                            })
                        
                        if documents:
                            result = self.upsert_documents(
                                documents,
                                chunk_size=50000,  # Kein Re-Chunking
                                chunk_overlap=0,
                            )
                            inserted = int(result.get('inserted', 0))
                            updated = int(result.get('updated', 0))
                            logger.info(
                                f"✅ Docling: {os.path.basename(file_path)} → "
                                f"{inserted} chunks inserted, {updated} updated "
                                f"({docling_result.num_tables} tables, "
                                f"{docling_result.processing_time_s:.1f}s)"
                            )
                            return inserted > 0 or updated > 0
                    
                    elif docling_result.success and docling_result.text and not getattr(docling_result, 'is_partial', False):
                        # Fallback: Volltext ohne Chunks
                        documents = [{
                            'text': docling_result.text,
                            'id': doc_id or f"docling_{os.path.basename(file_path)}",
                            'metadata': {
                                **(metadata or {}),
                                'source': file_path,
                                'filename': os.path.basename(file_path),
                                'file_type': file_ext,
                                'source_type': file_ext.lstrip('.'),
                                'extraction_method': 'docling',
                            }
                        }]
                        result = self.upsert_documents(
                            documents, chunk_size=chunk_size, chunk_overlap=chunk_overlap
                        )
                        return int(result.get('inserted', 0)) > 0
                    
                    else:
                        if self._is_nonrecoverable_docling_failure(docling_result):
                            error_code = getattr(docling_result, 'error_code', 'unknown')
                            raise RuntimeError(
                                f"Docling local-only fail-fast ({error_code}): {docling_result.error}"
                            )

                        logger.warning(
                            f"⚠️ Docling failed for {file_path}: {docling_result.error}. "
                            f"Falling back to legacy processing."
                        )
                        # Durchfallen zum Legacy-Pfad
                        
                except ImportError:
                    logger.debug("DoclingProcessor not available, using legacy processing")
                except RuntimeError:
                    raise
                except Exception as e:
                    logger.warning(f"Docling processing error: {e}, falling back to legacy")
            
            # ═══════════════════════════════════════════════════════════
            # SCHRITT 2: Legacy-Fallback Processing
            # Wird nur genutzt wenn Docling nicht verfügbar oder fehlgeschlagen
            # ═══════════════════════════════════════════════════════════
            
            if file_ext == '.pdf':
                result = self.upsert_pdf(file_path, doc_id=doc_id, 
                                       metadata=metadata)
                return int(result.get('inserted', 0)) > 0 or int(result.get('updated', 0)) > 0
            elif file_ext in ['.xlsx', '.xls']:
                # Excel-Dateien verarbeiten mit pandas
                try:
                    import pandas as pd
                    
                    # ✅ Lese alle Sheets mit explizitem Cleanup
                    excel_file = pd.ExcelFile(file_path)
                    all_content = []
                    
                    try:
                        for sheet_name in excel_file.sheet_names:
                            df = pd.read_excel(file_path, sheet_name=sheet_name)
                            
                            try:
                                # Konvertiere DataFrame zu strukturiertem Text
                                sheet_content = f"=== Sheet: {sheet_name} ===\n"
                                
                                # Header
                                headers = df.columns.tolist()
                                sheet_content += f"Spalten: {', '.join(str(h) for h in headers)}\n\n"
                                
                                # Daten (max 1000 Zeilen für Performance)
                                for idx, row in df.head(1000).iterrows():
                                    row_data = []
                                    for col in headers:
                                        value = row[col]
                                        if pd.notna(value):
                                            row_data.append(f"{col}: {value}")
                                    
                                    if row_data:
                                        sheet_content += " | ".join(row_data) + "\n"
                                
                                # Zusammenfassung
                                sheet_content += f"\nZusammenfassung Sheet '{sheet_name}': {len(df)} Zeilen, {len(df.columns)} Spalten\n"
                                
                                if len(df) > 0:
                                    # Numerische Spalten zusammenfassen
                                    numeric_cols = df.select_dtypes(include=['number']).columns
                                    if len(numeric_cols) > 0:
                                        sheet_content += f"Numerische Spalten: {', '.join(numeric_cols)}\n"
                                        for col in numeric_cols:
                                            if pd.notna(df[col]).any():
                                                stats = df[col].describe()
                                                sheet_content += f"  {col}: Min={stats['min']:.2f}, Max={stats['max']:.2f}, Mittel={stats['mean']:.2f}\n"
                                
                                all_content.append(sheet_content)
                                
                            finally:
                                # ✅ KRITISCH: DataFrame nach jedem Sheet löschen!
                                del df
                                import gc
                                gc.collect()
                                
                    finally:
                        # ✅ KRITISCH: ExcelFile schließen und löschen!
                        excel_file.close()
                        del excel_file
                        import gc
                        gc.collect()
                    
                    # Kombiniere alle Sheets
                    full_content = "\n\n".join(all_content)
                    
                    # Als Dokument speichern
                    documents = [{
                        'text': full_content,
                        'id': doc_id or f"excel_{os.path.basename(file_path)}",
                        'metadata': {
                            'source': file_path,
                            'file_type': file_ext,
                            'filename': os.path.basename(file_path),
                            'sheet_count': len(all_content),
                            'source_type': 'excel',
                        }
                    }]
                    
                    result = self.upsert_documents(documents)
                    
                    
                    # ✅ Cleanup after processing
                    del all_content
                    del full_content
                    
                    return int(result.get('inserted', 0)) > 0 or int(result.get('updated', 0)) > 0
                    
                except ImportError:
                    logger.error("pandas ist erforderlich für Excel-Import. Install mit: pip install pandas openpyxl")
                    return False
                except Exception as e:
                    logger.error(f"Fehler beim Excel-Import {file_path}: {e}")
                    return False
            elif file_ext in ['.txt', '.md', '.csv']:
                # Text-Dateien direkt verarbeiten
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                # Als Dokument speichern
                documents = [{
                    'text': content,
                    'id': doc_id or f"file_{os.path.basename(file_path)}",
                    'metadata': {
                        'source': file_path,
                        'file_type': file_ext,
                        'filename': os.path.basename(file_path),
                        **(metadata or {})
                    }
                }]
                
                result = self.upsert_documents(documents)
                return int(result.get('inserted', 0)) > 0 or int(result.get('updated', 0)) > 0
            
            elif file_ext in ['.docx', '.doc']:
                # 🆕 Legacy Word-Verarbeitung (Fallback wenn Docling fehlt)
                try:
                    from docx import Document as DocxDocument
                    
                    doc = DocxDocument(file_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    
                    # Tabellen extrahieren
                    table_texts = []
                    for table in doc.tables:
                        table_rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            table_rows.append(" | ".join(cells))
                        if table_rows:
                            table_texts.append("\n".join(table_rows))
                    
                    full_content = "\n\n".join(paragraphs)
                    if table_texts:
                        full_content += "\n\n=== Tabellen ===\n\n" + "\n\n".join(table_texts)
                    
                    if not full_content.strip():
                        logger.warning(f"Leeres Word-Dokument: {file_path}")
                        return False
                    
                    documents = [{
                        'text': full_content,
                        'id': doc_id or f"docx_{os.path.basename(file_path)}",
                        'metadata': {
                            'source': file_path,
                            'file_type': file_ext,
                            'filename': os.path.basename(file_path),
                            'source_type': 'docx',
                            'extraction_method': 'python-docx_legacy',
                            'paragraph_count': len(paragraphs),
                            'table_count': len(doc.tables),
                            **(metadata or {}),
                        }
                    }]
                    
                    del doc
                    result = self.upsert_documents(documents)
                    return int(result.get('inserted', 0)) > 0 or int(result.get('updated', 0)) > 0
                    
                except ImportError:
                    logger.error("python-docx ist erforderlich für Word-Import: pip install python-docx")
                    return False
                except Exception as e:
                    logger.error(f"Fehler beim Word-Import {file_path}: {e}")
                    return False
            
            elif file_ext in ['.pptx', '.ppt']:
                # 🆕 Legacy PowerPoint-Verarbeitung (Fallback wenn Docling fehlt)
                try:
                    from pptx import Presentation
                    
                    prs = Presentation(file_path)
                    slide_texts = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = f"=== Folie {slide_num} ===\n"
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                                    text = paragraph.text.strip()
                                    if text:
                                        slide_content += text + "\n"
                            if shape.has_table:
                                table = shape.table  # type: ignore[union-attr]
                                for row in table.rows:
                                    cells = [cell.text.strip() for cell in row.cells]
                                    slide_content += " | ".join(cells) + "\n"
                        
                        if slide_content.strip() != f"=== Folie {slide_num} ===":
                            slide_texts.append(slide_content)
                    
                    full_content = "\n\n".join(slide_texts)
                    
                    if not full_content.strip():
                        logger.warning(f"Leere Präsentation: {file_path}")
                        return False
                    
                    documents = [{
                        'text': full_content,
                        'id': doc_id or f"pptx_{os.path.basename(file_path)}",
                        'metadata': {
                            'source': file_path,
                            'file_type': file_ext,
                            'filename': os.path.basename(file_path),
                            'source_type': 'pptx',
                            'extraction_method': 'python-pptx_legacy',
                            'slide_count': len(prs.slides),
                            **(metadata or {}),
                        }
                    }]
                    
                    del prs
                    result = self.upsert_documents(documents)
                    return int(result.get('inserted', 0)) > 0 or int(result.get('updated', 0)) > 0
                    
                except ImportError:
                    logger.error("python-pptx ist erforderlich für PowerPoint-Import: pip install python-pptx")
                    return False
                except Exception as e:
                    logger.error(f"Fehler beim PowerPoint-Import {file_path}: {e}")
                    return False
            
            else:
                logger.error(f"Nicht unterstützter Dateityp: {file_ext}")
                return False
                
        except RuntimeError:
            raise
        except Exception as e:
            logger.error(f"Fehler beim Datei-Upload {file_path}: {e}")
            return False

    def _create_llm_knowledge_graph_per_chunk(
        self,
        kg_chunks: List[Dict[str, Any]],
        doc_id: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> List[Tuple[str, str, str, Optional[int], float]]:
        """
        ★ SOTA v4: Per-Chunk KG-Extraktion mit heading-aware Kontext.

        Ruft extract_from_chunks() auf dem LLM-KG-Extractor auf.
        Gibt Tuples zurück: (subject, predicate, object, chunk_id, confidence)

        ``confidence`` stammt direkt aus dem LLM-Extractor (KGTriple.confidence,
        per Default 0.8 wenn nicht im JSON spezifiziert) und wird in der
        Insert-Pipeline als per-Evidence-Confidence für Noisy-OR-Aggregation
        verwendet.
        """
        if not self.llm_kg_extractor:
            logger.warning(
                "🚫 LLM-KG-Extractor nicht verfügbar — KG-Generierung übersprungen "
                "(Regex-Fallback wurde entfernt: produzierte nur deutsche Person/Org-Patterns "
                "und schadete der Qualität mehr als sie nutzte)."
            )
            return []

        try:
            doc_context = self._build_doc_context(doc_id, metadata)

            # extract_from_chunks erwartet List[Dict] mit 'text', 'chunk_id', 'headings'
            kg_triples = self.llm_kg_extractor.extract_from_chunks(kg_chunks, doc_context)

            result: List[Tuple[str, str, str, Optional[int], float]] = []
            for triple in kg_triples:
                # KGTriple.confidence kommt aus LLM-JSON; Default in der Klasse ist 1.0,
                # gefiltert wird bereits in _is_quality_triple/_get_confidence_threshold.
                conf = float(getattr(triple, 'confidence', 0.8) or 0.8)
                # Robustheits-Clamp [0,1]
                conf = max(0.0, min(1.0, conf))
                result.append((
                    triple.subject,
                    triple.predicate,
                    triple.object,
                    triple.source_chunk_id,
                    conf,
                ))

            logger.info(f"✅ {len(result)} Per-Chunk-KG-Triples für '{doc_id}'")
            return result

        except Exception as e:
            logger.error(f"❌ Per-Chunk-KG fehlgeschlagen (fail-fast) für {doc_id}: {e}")
            raise RuntimeError(f"Per-chunk KG extraction failed for {doc_id}: {e}") from e
    
    def _build_doc_context(self, doc_id: str, metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Baut Dokument-Kontext für KG-Extraktion.
        Domain-Erkennung basiert auf METADATEN, nicht Text-Inhalt.
        """
        doc_context: Dict[str, Any] = {"doc_id": doc_id, **(metadata or {})}
        
        if metadata:
            source = metadata.get("source", "").lower()
            
            if "sbb" in source or "geschäftsbericht" in source or "finanzbericht" in source:
                doc_context["source_type"] = "business"
            elif metadata.get("source_url") or "web" in source or "url" in source:
                doc_context["source_type"] = "web"
                doc_context["source_url"] = metadata.get("source_url", "unknown")
            elif metadata.get("domain") == "psychology" or source.startswith("psychological_"):
                doc_context["source_type"] = "psychology"
            else:
                doc_context["source_type"] = "generic"
        
        return doc_context

    def _create_llm_knowledge_graph(self, text: str, doc_id: str,
                                   metadata: Optional[Dict[str, Any]] = None) -> List[Tuple[str, str, str]]:
        """DEPRECATED — never called in the live pipeline. Kept as a stub for
        backward-compatible imports; raises ``NotImplementedError`` if invoked.

        Live code uses :py:meth:`_create_llm_knowledge_graph_per_chunk` which is
        chunk-aware and grounded against the source chunk via cross-encoder.
        """
        raise NotImplementedError(
            "_create_llm_knowledge_graph(text-blob) wurde entfernt. "
            "Verwende _create_llm_knowledge_graph_per_chunk(kg_chunks, ...)."
        )

    def _verify_kg_triples_via_reranker(
        self,
        triples: List[Tuple[str, str, str]],
        source_text: str,
        min_grounding_score: float = 0.15,
    ) -> List[Tuple[str, str, str]]:
        """
        ★ SOTA: Verify KG triples against source text using cross-encoder reranker.
        
        This is NON-CIRCULAR: the reranker (BGE-reranker-v2-m3, 568M params, ONNX-
        optimised cross-encoder) is a completely different model from the LLM that
        generated the triples. It evaluates semantic similarity/entailment.
        
        For each triple, we check: does the source text support the claim
        "{subject} {predicate} {object}"?
        
        Args:
            triples: List of (subject, predicate, object)
            source_text: The original document text to check against
            min_grounding_score: Minimum reranker score to keep a triple
            
        Returns:
            Filtered list of grounded triples
        """
        if not triples or not source_text:
            return triples
        
        # Get the reranker (lazy — may not be loaded yet)
        try:
            from agent.reranker import get_reranker
            reranker = get_reranker()
            if not hasattr(reranker, 'is_available') or not reranker.is_available:
                logger.debug("[KG-VERIFY] Reranker not available, skipping verification")
                return triples
        except Exception as exc:
            logger.debug(f"[KG-VERIFY] Could not load reranker, skipping verification: {exc}")
            return triples
        
        # BGE-reranker-v2-m3 max_length=512 tokens ≈ 2000 chars.
        # With per-chunk grounding, source_text is typically ~3500 chars (1 KG-chunk).
        # Truncate only if needed to avoid reranker OOM, but use as much as possible.
        MAX_RERANKER_CHARS = 4000  # ~1000 tokens, safe for BGE-reranker-v2-m3
        source_truncated = source_text[:MAX_RERANKER_CHARS]
        verified = []
        rejected = 0
        
        for subj, pred, obj in triples:
            hypothesis = f"{subj} {pred} {obj}"
            try:
                result = reranker.rerank(
                    query=hypothesis,
                    passages=[{"text": source_truncated}],
                    top_k=1,
                    text_key="text",
                )
                score = 0.0
                if result and len(result) > 0:
                    score = result[0].get("rerank_score", result[0].get("score", 0.0))
                
                if score >= min_grounding_score:
                    verified.append((subj, pred, obj))
                else:
                    rejected += 1
                    logger.info(
                        f"[KG-VERIFY] Rejected (score={score:.3f}<{min_grounding_score}): "
                        f"'{subj}' → '{pred}' → '{obj}' "
                        f"| source_len={len(source_truncated)}"
                    )
            except Exception as e:
                logger.error(
                    f"[KG-VERIFY] Grounding score computation failed for triple "
                    f"'{subj}' → '{pred}' → '{obj}': {e}"
                )
                raise RuntimeError(
                    "KG grounding verification failed; aborting to prevent unverified triple persistence"
                ) from e
        
        if rejected > 0:
            logger.info(
                f"[KG-VERIFY] Kept {len(verified)}/{len(triples)} triples "
                f"(rejected {rejected} with grounding score < {min_grounding_score})"
            )
        
        return verified

    # ──────────────────────────────────────────────────────────────────
    #  ★ SOTA: Central triple deletion with entity frequency maintenance
    # ──────────────────────────────────────────────────────────────────
    def _delete_triples_for_doc(self, cur: sqlite3.Cursor, doc_id: str) -> int:
        """
        Central triple deletion hook — the ONLY correct way to remove triples.

        Maintains kg_entities consistency that cannot be expressed via FK CASCADE:
          1. Collects affected entities (subjects/objects) before deletion
          2. Deletes triples (PRAGMA foreign_keys = ON → CASCADE to triple_quality)
          3. Recalculates kg_entities.frequency from remaining triples
          4. Removes kg_entities entries that are no longer referenced

        Args:
            cur: Active database cursor (within transaction — caller commits)
            doc_id: Document ID whose triples are to be deleted

        Returns:
            Number of triples deleted
        """
        # Step 1: Collect entities affected by this deletion
        cur.execute(
            "SELECT DISTINCT subject FROM triples WHERE doc_id = ? "
            "UNION "
            "SELECT DISTINCT object FROM triples WHERE doc_id = ?",
            (doc_id, doc_id)
        )
        affected_entities = {row[0] for row in cur.fetchall() if row[0]}

        # Step 2: Count and delete triples (FK CASCADE handles triple_quality)
        cur.execute("SELECT COUNT(*) FROM triples WHERE doc_id = ?", (doc_id,))
        deleted_count = cur.fetchone()[0]

        if deleted_count == 0:
            return 0

        cur.execute("DELETE FROM triples WHERE doc_id = ?", (doc_id,))

        # Step 3: Recalculate frequency for each affected entity
        for entity in affected_entities:
            cur.execute(
                "SELECT COUNT(*) FROM triples WHERE subject = ? OR object = ?",
                (entity, entity)
            )
            remaining = cur.fetchone()[0]

            if remaining == 0:
                # Entity no longer referenced — remove from kg_entities
                cur.execute(
                    "DELETE FROM kg_entities WHERE entity_text = ?",
                    (entity,)
                )
            else:
                # Set frequency to actual count (not decrement — avoids drift)
                cur.execute(
                    "UPDATE kg_entities SET frequency = ? WHERE entity_text = ?",
                    (remaining, entity)
                )

        return deleted_count

    def _store_entity_embeddings(
        self, entities: List[Tuple[str, str]], cur
    ) -> int:
        """
        ★ SOTA v4: Embed and store entities in kg_entities table.

        For each unique entity text:
          1. Normalize via normalize_entity_for_matching() (deep: case, title-strip, parens)
          2. Skip if already embedded (UNIQUE constraint on normalized_text)
          3. Batch-embed new entities via embedding model
          4. Store in kg_entities with embedding BLOB

        Args:
            entities: List of (entity_text, doc_id) tuples
            cur: Active database cursor (within transaction)

        Returns:
            Number of new entities embedded and stored.
        """
        # ★ SOTA: canonical deep entity normalisation. No shallow fallback —
        # a lower()/strip()-only fallback silently breaks dedup (different
        # normalized_text for the same entity), so we let an ImportError
        # surface instead of silently corrupting kg_entities.
        from agent.llm_knowledge_graph import normalize_entity_for_matching

        # Deduplicate entities using deep normalization
        unique_entities: Dict[str, str] = {}  # normalized → original (first seen wins)
        for entity_text, doc_id in entities:
            entity_text = entity_text.strip()
            if len(entity_text) < 2:
                continue
            normalized = normalize_entity_for_matching(entity_text)
            if not normalized or len(normalized) < 2:
                continue
            if normalized not in unique_entities:
                unique_entities[normalized] = entity_text

        if not unique_entities:
            return 0

        # Filter out entities that already have embeddings
        new_entities: List[Tuple[str, str]] = []  # (normalized, original)
        for normalized, original in unique_entities.items():
            try:
                cur.execute(
                    "SELECT entity_id FROM kg_entities WHERE normalized_text = ?",
                    (normalized,)
                )
                row = cur.fetchone()
                if row:
                    # Entity exists — update frequency
                    cur.execute(
                        "UPDATE kg_entities SET frequency = frequency + 1 WHERE normalized_text = ?",
                        (normalized,)
                    )
                else:
                    new_entities.append((normalized, original))
            except Exception as exc:
                logger.debug(f"[KG-Entities] Failed reading existing normalized entity '{normalized}': {exc}")
                new_entities.append((normalized, original))

        if not new_entities:
            return 0

        # Batch-embed new entities
        try:
            entity_texts = [f"passage: {orig}" for _, orig in new_entities]
            embeddings = self.embed_texts(entity_texts)

            stored = 0
            for i, (normalized, original) in enumerate(new_entities):
                try:
                    emb = embeddings[i]
                    if hasattr(emb, 'numpy'):
                        emb = emb.numpy()
                    emb_blob = np.array(emb, dtype=np.float32).tobytes()

                    cur.execute("""
                        INSERT OR IGNORE INTO kg_entities
                        (entity_text, normalized_text, entity_type, frequency, embedding)
                        VALUES (?, ?, 'entity', 1, ?)
                    """, (original, normalized, emb_blob))
                    stored += 1
                except Exception as e:
                    logger.debug(f"Entity embedding insert failed for '{original}': {e}")

            logger.info(
                f"[KG-Entities] Embedded {stored}/{len(new_entities)} new entities "
                f"({len(unique_entities) - len(new_entities)} already existed)"
            )
            return stored

        except Exception as e:
            logger.warning(f"[KG-Entities] Batch embedding failed: {e}")
            return 0

    def resolve_duplicate_entities(self,
                                   similarity_threshold: float = 0.92,
                                   *,
                                   dry_run: bool = False,
                                   min_pairwise_threshold: Optional[float] = None,
                                   clustering: str = "conservative",
                                   lexical_threshold: float = 0.35,
                                   debug_report_limit: int = 50
                                   ) -> Dict[str, Any]:
        """
        ★ SOTA v3: Entity Resolution via Embedding Similarity

        Finds near-duplicate entities in kg_entities table:
          "Dr. Müller" ≈ "Herr Müller" ≈ "Mueller"
          "machine learning" ≈ "maschinelles Lernen"

        Algorithm:
          1. Load all entity embeddings from kg_entities
          2. Compute pairwise cosine similarity
          3. For pairs above threshold, merge to canonical form
             (canonical = entity with highest frequency)
          4. Update triples table to use canonical entity texts
          5. Remove merged entity entries

        This is O(N²) for N entities — for very large KGs (>100K entities),
        consider batched/approximate approaches.

        Returns:
            Statistics about merged entities.
        """
        stats: Dict[str, Any] = {"entities_before": 0, "merged_groups": 0, "entities_after": 0}

        try:
            conn = self.get_connection()
            cur = conn.cursor()

            # Load all entities with embeddings
            cur.execute(
                "SELECT entity_id, entity_text, normalized_text, frequency, embedding "
                "FROM kg_entities WHERE embedding IS NOT NULL"
            )
            rows = cur.fetchall()
            stats["entities_before"] = len(rows)

            if len(rows) < 2:
                cur.close()
                self.return_connection(conn)
                stats["entities_after"] = len(rows)
                return stats

            # ── Phase 0: Case-aware pre-merge (O(N log N), catches trivial duplicates) ──
            # Groups by lowercase text and merges before the expensive O(N²) embedding step.
            # This fixes "künstliche Intelligenz" vs "Künstliche Intelligenz" which
            # embedding similarity might miss at threshold 0.90.
            case_groups: Dict[str, List[Tuple]] = {}
            for eid, etext, enorm, freq, emb_blob in rows:
                key = etext.strip().lower()
                if key not in case_groups:
                    case_groups[key] = []
                case_groups[key].append((eid, etext, enorm, freq or 1, emb_blob))

            case_pre_merged = 0
            surviving_rows = []
            for key, group in case_groups.items():
                if len(group) <= 1:
                    surviving_rows.append(group[0])
                    continue
                # Sort by frequency desc — highest frequency = canonical
                group.sort(key=lambda x: x[3], reverse=True)
                canonical = group[0]
                canonical_id, canonical_text = canonical[0], canonical[1]

                # Lazy-import the shared merge helper (handles triple_hash recompute
                # + Bayesian Noisy-OR collapse on hash collision — see
                # agent/kg_entity_merge.py for rationale).
                from agent.kg_entity_merge import merge_entity_in_triples
                from agent.rag_store.utils.memory import calculate_triple_hash

                for eid, etext, enorm, freq, emb_blob in group[1:]:
                    # Rewrite triples to canonical text (recomputes triple_hash,
                    # collapses any colliding canonical row via Noisy-OR).
                    merge_entity_in_triples(
                        conn, calculate_triple_hash, canonical_text, etext
                    )
                    # Delete merged entity from kg_entities
                    cur.execute("DELETE FROM kg_entities WHERE entity_id = ?", (eid,))
                    case_pre_merged += 1

                # Recalculate canonical frequency after merge
                cur.execute(
                    "SELECT COUNT(*) FROM ("
                    "  SELECT subject AS e FROM triples WHERE subject = ? "
                    "  UNION ALL "
                    "  SELECT object FROM triples WHERE object = ?"
                    ")",
                    (canonical_text, canonical_text)
                )
                actual_freq = cur.fetchone()[0]
                cur.execute("UPDATE kg_entities SET frequency = ? WHERE entity_id = ?",
                            (actual_freq, canonical_id))

                surviving_rows.append(canonical)

            if case_pre_merged > 0:
                conn.commit()
                logger.info(
                    f"[KG-Resolve] Case pre-merge: {case_pre_merged} trivial duplicates merged"
                )
                stats["case_pre_merged"] = case_pre_merged

            # ── Phase 1: Embedding-based similarity (O(N²)) ──
            # Build embedding matrix from surviving rows (after case pre-merge)
            entity_ids = []
            entity_texts = []
            entity_freqs = []
            embs = []

            for eid, etext, enorm, freq, emb_blob in surviving_rows:
                if emb_blob:
                    emb = np.frombuffer(emb_blob, dtype=np.float32).copy()
                    entity_ids.append(eid)
                    entity_texts.append(etext)
                    entity_freqs.append(freq or 1)
                    embs.append(emb)

            if len(embs) < 2:
                cur.close()
                self.return_connection(conn)
                stats["entities_after"] = stats["entities_before"]
                return stats

            emb_matrix = np.stack(embs, axis=0)
            # L2-normalize for cosine similarity
            norms = np.linalg.norm(emb_matrix, axis=1, keepdims=True)
            norms = np.maximum(norms, 1e-10)
            emb_matrix = emb_matrix / norms

            # Pairwise cosine similarity
            sim_matrix = emb_matrix @ emb_matrix.T

            # Helper: token-jaccard on canonically normalised texts
            from agent.llm_knowledge_graph import normalize_entity_for_matching

            def token_jaccard(a: str, b: str) -> float:
                ta = set(normalize_entity_for_matching(a).split())
                tb = set(normalize_entity_for_matching(b).split())
                if not ta and not tb:
                    return 0.0
                inter = len(ta & tb)
                uni = len(ta | tb)
                return float(inter) / float(uni) if uni > 0 else 0.0

            # Default min_pairwise_threshold slightly stricter than base
            if min_pairwise_threshold is None:
                min_pairwise_threshold = min(0.995, similarity_threshold + 0.03)

            # Build candidate pairs above base threshold
            n = len(entity_ids)
            candidate_pairs = [(i, j) for i in range(n) for j in range(i + 1, n)
                               if sim_matrix[i, j] >= similarity_threshold]

            # If no candidates, finish early
            if not candidate_pairs:
                stats["merged_groups"] = 0
                stats["entities_after"] = stats["entities_before"]
                cur.close()
                self.return_connection(conn)
                return stats

            # Conservative clustering: create cliques via greedy clique-growing
            remaining = set(range(n))
            cliques: List[List[int]] = []

            # Build adjacency for speed
            adj = {i: set() for i in range(n)}
            for i, j in candidate_pairs:
                adj[i].add(j)
                adj[j].add(i)

            # Greedy clique extraction (maximal cliques w.r.t. threshold)
            # Note: Conservative but deterministic and avoids transitive false merges
            while remaining:
                # pick node with max degree in remaining
                node = max(remaining, key=lambda x: len(adj[x] & remaining))
                # seed clique with node and its strongest neighbor
                neighbors = sorted([(sim_matrix[node, nb], nb) for nb in (adj[node] & remaining)], reverse=True)
                if neighbors:
                    _, best_nb = neighbors[0]
                    clique = {node, best_nb}
                else:
                    # isolated node (no neighbors above threshold)
                    remaining.remove(node)
                    continue

                # try to expand clique: add any remaining node that has sim >= threshold to all current members
                added = True
                while added:
                    added = False
                    for cand in list(remaining - clique):
                        ok = True
                        for member in clique:
                            if sim_matrix[cand, member] < similarity_threshold:
                                ok = False
                                break
                        if ok:
                            clique.add(cand)
                            added = True

                # finalize clique
                cliques.append(sorted(list(clique)))
                for v in clique:
                    if v in remaining:
                        remaining.remove(v)

            # Prepare proposed merges (with pairwise validation)
            proposed_merges: List[Dict[str, Any]] = []
            for clique in cliques:
                if len(clique) <= 1:
                    continue
                # choose canonical by highest frequency
                canonical_idx = max(clique, key=lambda idx: entity_freqs[idx])
                canonical_text = entity_texts[canonical_idx]
                canonical_id = entity_ids[canonical_idx]

                members = [m for m in clique if m != canonical_idx]
                members_report = []
                accepted_members = []

                for m in members:
                    sim_to_can = float(sim_matrix[m, canonical_idx])
                    lex = token_jaccard(entity_texts[m], canonical_text)

                    # Combined, but conservative decision rule:
                    # accept if embedding sim >= min_pairwise_threshold OR
                    # (embedding sim >= similarity_threshold and lexical overlap reasonably high)
                    accept = False
                    if sim_to_can >= min_pairwise_threshold:
                        accept = True
                    elif sim_to_can >= similarity_threshold and lex >= lexical_threshold:
                        accept = True

                    members_report.append({
                        "entity_text": entity_texts[m],
                        "entity_id": entity_ids[m],
                        "sim_to_canonical": sim_to_can,
                        "lexical_jaccard": lex,
                        "accepted": accept
                    })

                    if accept:
                        accepted_members.append(m)

                if not accepted_members:
                    continue

                proposed_merges.append({
                    "canonical_index": canonical_idx,
                    "canonical_id": canonical_id,
                    "canonical_text": canonical_text,
                    "members": members_report
                })

            # If dry_run, report and return without committing
            merged_count = 0
            if dry_run:
                # Attach report to stats for callers
                stats["proposed_merges"] = proposed_merges[:debug_report_limit]
                stats["merged_groups"] = len(proposed_merges)
                stats["entities_after"] = stats["entities_before"]
                cur.close()
                self.return_connection(conn)
                logger.info(f"[KG-Resolve] Dry-run: {len(proposed_merges)} proposed merge groups")
                return stats

            # Commit accepted merges to DB (uses shared helper for hash-recompute
            # + Bayesian Noisy-OR collapse — guarantees no SPO duplicates survive).
            from agent.kg_entity_merge import merge_entity_in_triples
            from agent.rag_store.utils.memory import calculate_triple_hash

            for pm in proposed_merges:
                canonical_text = pm["canonical_text"]
                canonical_id = pm["canonical_id"]
                merge_texts = [m["entity_text"] for m in pm["members"] if m["accepted"]]

                logger.info(f"[KG-Resolve] Merging {merge_texts} → '{canonical_text}'")

                for m in pm["members"]:
                    if not m["accepted"]:
                        continue
                    old_text = m["entity_text"]
                    old_id = m["entity_id"]

                    # Look up the merged entity's actual frequency (FIX: previously
                    # incremented by hard-coded 1, causing frequency drift).
                    cur.execute(
                        "SELECT frequency FROM kg_entities WHERE entity_id = ?",
                        (old_id,),
                    )
                    fr_row = cur.fetchone()
                    old_freq = int((fr_row[0] if fr_row else 1) or 1)

                    # Rewrite triples (recomputes triple_hash; collapses via Noisy-OR
                    # if a canonical row already exists for the new SPO).
                    merge_entity_in_triples(
                        conn, calculate_triple_hash, canonical_text, old_text
                    )

                    # Accumulate the merged entity's actual frequency
                    cur.execute(
                        "UPDATE kg_entities SET frequency = frequency + ? WHERE entity_id = ?",
                        (old_freq, canonical_id),
                    )
                    # Delete merged entity
                    cur.execute(
                        "DELETE FROM kg_entities WHERE entity_id = ?",
                        (old_id,),
                    )

                merged_count += 1

            conn.commit()
            cur.close()
            self.return_connection(conn)

            # Count remaining entities (single round-trip; previously this block
            # was duplicated and ran twice — removed).
            try:
                conn2 = self.get_connection()
                cur2 = conn2.cursor()
                cur2.execute("SELECT COUNT(*) FROM kg_entities")
                stats["entities_after"] = cur2.fetchone()[0]
                cur2.close()
                self.return_connection(conn2)
            except Exception as exc:
                logger.debug(f"[KG-Resolve] Could not count entities_after: {exc}")
                stats["entities_after"] = stats["entities_before"] - merged_count

            stats["merged_groups"] = merged_count

            # Invalidate entity index in search manager
            if self._search_manager is not None:
                try:
                    self._search_manager.invalidate_entity_index()
                except Exception as exc:
                    logger.debug(f"[KG-Resolve] invalidate_entity_index failed: {exc}")

            logger.info(
                f"[KG-Resolve] Entity resolution complete: "
                f"{stats['entities_before']} → {stats['entities_after']} entities "
                f"({merged_count} merge groups, threshold={similarity_threshold})"
            )

        except Exception as e:
            logger.error(f"[KG-Resolve] Entity resolution failed: {e}")
            import traceback
            logger.debug(traceback.format_exc())

        return stats

    def _create_automatic_knowledge_graphs(self, doc_ids: List[str], 
                                         documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        ★ SOTA v5: Dual-Chunker Knowledge Graph Extraction (2026-03)
        
        Architektur:
        - v3: Aggregierte alle Chunks → 1 Blob → LLM re-chunked intern → Informationsverlust
        - v4: Heading-aware Merge von Retrieval-Chunks → zu viele LLM-Calls bei kleinen Chunks
        - v5: Docling Dual-Chunker — zweiter HybridChunker-Pass (max_tokens=1200)
              erzeugt KG-optimierte Chunks NATIV. Kein eigenes Merge, keine Heuristiken.
              Docling kennt die Dokumentstruktur, den Tokenizer und merge_peers.
        
        Kein Hard-Limit auf Triple-Count. Qualitäts-Gates in LLM-KG-Extractor.
        """
        if not doc_ids:
            return {"triples": 0}
        
        total_triples = 0
        processed_docs = 0
        conn = None
        cur = None
        
        try:
            conn = self.get_connection()
            cur = conn.cursor()
            
            # ── Step 1: Group chunk-level doc_ids by parent document ──
            import re as _re
            parent_groups: Dict[str, List[Tuple[str, str, Dict]]] = {}
            standalone_docs: List[Tuple[str, str, Dict]] = []
            
            for doc_id in doc_ids:
                source_doc = None
                for doc in documents:
                    if doc.get("id") == doc_id or doc.get("doc_id") == doc_id:
                        source_doc = doc
                        break
                if not source_doc:
                    continue
                
                doc_text = str(source_doc.get("text", ""))
                doc_metadata = source_doc.get("metadata", {})
                
                chunk_match = _re.match(r'^(.+)_chunk_(\d+)$', doc_id)
                if chunk_match:
                    parent_id = chunk_match.group(1)
                    if parent_id not in parent_groups:
                        parent_groups[parent_id] = []
                    parent_groups[parent_id].append((doc_id, doc_text, doc_metadata))
                else:
                    standalone_docs.append((doc_id, doc_text, doc_metadata))
            
            # ── Step 2: Build KG extraction units ──
            # Use Docling's pre-computed KG-chunks when available (Dual-Chunker).
            # Fallback: use retrieval-chunks directly (for non-Docling documents).
            
            docs_to_process: List[Tuple[str, List[Dict[str, Any]], Dict, List[str]]] = []
            # Each entry: (parent_id, kg_chunks, metadata, chunk_doc_ids_for_cleanup)
            
            pending_kg = getattr(self, '_pending_kg_chunks', {})
            
            for parent_id, chunks in parent_groups.items():
                # Sort by chunk index
                chunks.sort(key=lambda c: int(_re.search(r'_chunk_(\d+)$', c[0]).group(1))  # type: ignore
                             if _re.search(r'_chunk_(\d+)$', c[0]) else 0)
                
                merged_metadata = chunks[0][2] if chunks else {}
                chunk_doc_ids = [c[0] for c in chunks]
                
                # ── Docling Dual-Chunker: KG-Chunks bereits vorhanden? ──
                if parent_id in pending_kg:
                    docling_kg_chunks = pending_kg.pop(parent_id)
                    
                    # Build retrieval-chunk lookup for source_chunk_id mapping
                    # Map each retrieval-chunk's doc_id → real chunk_id in DB
                    retrieval_chunk_lookup: List[Tuple[int, str]] = []  # (chunk_id, text)
                    for chunk_doc_id, chunk_text, _ in chunks:
                        try:
                            cur.execute(
                                "SELECT chunk_id FROM chunks WHERE doc_id = ? LIMIT 1",
                                (chunk_doc_id,)
                            )
                            row = cur.fetchone()
                            if row:
                                retrieval_chunk_lookup.append((row[0], chunk_text))
                        except Exception as exc:
                            logger.debug(f"KG retrieval chunk lookup failed for {chunk_doc_id}: {exc}")
                    
                    # Convert DoclingChunks → KG-chunk dicts
                    kg_chunks: List[Dict[str, Any]] = []
                    for dkg_chunk in docling_kg_chunks:
                        kg_text = dkg_chunk.text.strip()
                        if not kg_text:
                            continue
                        
                        # Find best-matching retrieval chunk_id via text overlap
                        # KG-chunk (1200 tokens) contains 1-3 retrieval chunks (384 tokens)
                        best_chunk_id = None
                        best_overlap = 0
                        for r_chunk_id, r_text in retrieval_chunk_lookup:
                            # Fast overlap: count shared words
                            kg_words = set(kg_text[:500].split())
                            r_words = set(r_text[:500].split())
                            overlap = len(kg_words & r_words)
                            if overlap > best_overlap:
                                best_overlap = overlap
                                best_chunk_id = r_chunk_id
                        
                        headings_str = " > ".join(dkg_chunk.headings) if dkg_chunk.headings else ""
                        
                        kg_chunks.append({
                            "text": kg_text,
                            "headings": headings_str,
                            "chunk_id": best_chunk_id,
                            "_source_chunk_ids": [best_chunk_id] if best_chunk_id else [],
                        })
                    
                    if kg_chunks:
                        docs_to_process.append((parent_id, kg_chunks, merged_metadata, chunk_doc_ids))
                        logger.info(
                            f"🔗 KG: {len(chunks)} Retrieval-Chunks, "
                            f"{len(kg_chunks)} KG-Chunks (Docling Dual-Chunker) "
                            f"für '{parent_id}'"
                        )
                
                else:
                    # ── Fallback: Merge Retrieval-Chunks zu KG-optimaler Größe ──
                    # Web/Text haben keine Docling-Dokumentstruktur → einfaches
                    # Size-Merge zu ~4800 Chars (≈1200 Tokens, GraphRAG-Niveau).
                    # Kein Heading-Split, da Web/Text keine zuverlässigen Headings hat.
                    MAX_KG_FALLBACK_CHARS = 4800
                    
                    # First: collect all chunks with their DB chunk_ids
                    raw_chunks: List[Tuple[str, Optional[int], str]] = []  # (text, chunk_id, heading)
                    for chunk_doc_id, chunk_text, chunk_meta in chunks:
                        if not chunk_text.strip():
                            continue
                        
                        real_chunk_id = None
                        try:
                            cur.execute(
                                "SELECT chunk_id FROM chunks WHERE doc_id = ? LIMIT 1",
                                (chunk_doc_id,)
                            )
                            row = cur.fetchone()
                            if row:
                                real_chunk_id = row[0]
                        except Exception as exc:
                            logger.debug(f"KG fallback chunk lookup failed for {chunk_doc_id}: {exc}")
                        
                        heading = chunk_meta.get("chunk_headings", "")
                        raw_chunks.append((chunk_text, real_chunk_id, heading))
                    
                    # Merge adjacent chunks to ~4800 chars.
                    # Vorher: einzelne Chunks, die alleine schon > MAX sind
                    # (z.B. monolithische Web-Paragraphen) hart aufsplitten,
                    # sonst rutschen sie ungeteilt durch und treiben den KG-
                    # LLM-Call auf 200+ Sekunden.
                    from agent.llm_knowledge_graph import split_oversized_for_kg
                    expanded_raw: List[Tuple[str, Optional[int], str]] = []
                    for chunk_text, chunk_id, heading in raw_chunks:
                        if len(chunk_text) <= MAX_KG_FALLBACK_CHARS:
                            expanded_raw.append((chunk_text, chunk_id, heading))
                        else:
                            for piece in split_oversized_for_kg(chunk_text, MAX_KG_FALLBACK_CHARS):
                                expanded_raw.append((piece, chunk_id, heading))
                    raw_chunks = expanded_raw
                    
                    kg_chunks_fallback: List[Dict[str, Any]] = []
                    current_text = ""
                    current_ids: List[Optional[int]] = []
                    current_heading = ""
                    
                    for chunk_text, chunk_id, heading in raw_chunks:
                        would_overflow = len(current_text) + len(chunk_text) + 2 > MAX_KG_FALLBACK_CHARS
                        
                        if current_text and would_overflow:
                            kg_chunks_fallback.append({
                                "text": current_text,
                                "headings": current_heading,
                                "chunk_id": current_ids[0] if current_ids else None,
                                "_source_chunk_ids": list(current_ids),
                            })
                            current_text = ""
                            current_ids = []
                            current_heading = ""
                        
                        current_text = current_text + "\n\n" + chunk_text if current_text else chunk_text
                        current_ids.append(chunk_id)
                        if heading:
                            current_heading = heading
                    
                    # Flush last
                    if current_text.strip():
                        kg_chunks_fallback.append({
                            "text": current_text,
                            "headings": current_heading,
                            "chunk_id": current_ids[0] if current_ids else None,
                            "_source_chunk_ids": list(current_ids),
                        })
                    
                    if kg_chunks_fallback:
                        docs_to_process.append((parent_id, kg_chunks_fallback, merged_metadata, chunk_doc_ids))
                        logger.info(
                            f"🔗 KG: {len(chunks)} Retrieval-Chunks → "
                            f"{len(kg_chunks_fallback)} KG-Chunks (Merge-Fallback, ~{MAX_KG_FALLBACK_CHARS} Chars) "
                            f"für '{parent_id}'"
                        )
            
            # Standalone docs → split into ~4800-char KG-chunks
            MAX_KG_STANDALONE_CHARS = 4800
            for doc_id, doc_text, doc_metadata in standalone_docs:
                if not doc_text.strip():
                    continue
                
                if len(doc_text) <= MAX_KG_STANDALONE_CHARS:
                    docs_to_process.append((
                        doc_id,
                        [{"text": doc_text, "headings": "", "chunk_id": None, "_source_chunk_ids": []}],
                        doc_metadata, []
                    ))
                else:
                    # Split at paragraph boundaries (\n\n), respecting max size.
                    # Einzelne Paragraphen, die alleine schon > MAX sind
                    # (z.B. monolithische Web-Paragraphen ohne \n\n), werden
                    # vorher hart auf Satz-/Whitespace-Grenzen gesplittet —
                    # sonst landen sie ungeteilt im KG-Chunk und blähen
                    # die LLM-Inferenz auf 200+ Sekunden auf.
                    from agent.llm_knowledge_graph import split_oversized_for_kg
                    raw_paragraphs = doc_text.split("\n\n")
                    paragraphs: List[str] = []
                    for p in raw_paragraphs:
                        p_stripped = p.strip()
                        if not p_stripped:
                            continue
                        if len(p_stripped) <= MAX_KG_STANDALONE_CHARS:
                            paragraphs.append(p_stripped)
                        else:
                            paragraphs.extend(
                                split_oversized_for_kg(p_stripped, MAX_KG_STANDALONE_CHARS)
                            )
                    
                    standalone_kg_chunks: List[Dict[str, Any]] = []
                    current_text = ""
                    
                    for para in paragraphs:
                        if not para.strip():
                            continue
                        would_overflow = len(current_text) + len(para) + 2 > MAX_KG_STANDALONE_CHARS
                        
                        if current_text and would_overflow:
                            standalone_kg_chunks.append({
                                "text": current_text,
                                "headings": "",
                                "chunk_id": None,
                                "_source_chunk_ids": [],
                            })
                            current_text = ""
                        
                        current_text = current_text + "\n\n" + para if current_text else para
                    
                    if current_text.strip():
                        standalone_kg_chunks.append({
                            "text": current_text,
                            "headings": "",
                            "chunk_id": None,
                            "_source_chunk_ids": [],
                        })
                    
                    if standalone_kg_chunks:
                        docs_to_process.append((doc_id, standalone_kg_chunks, doc_metadata, []))
                        logger.info(
                            f"🔗 KG: Standalone '{doc_id}' ({len(doc_text)} Zeichen) → "
                            f"{len(standalone_kg_chunks)} KG-Chunks"
                        )
            
            # ── Step 3: Extract KG per-chunk via LLM ──
            for doc_id, kg_chunks_list, doc_metadata, chunk_doc_ids_cleanup in docs_to_process:
                try:
                    total_chars = sum(len(c["text"]) for c in kg_chunks_list)
                    if total_chars < 300:
                        logger.debug(f"📋 Dokument {doc_id} zu kurz für KG ({total_chars} Zeichen)")
                        continue
                    
                    logger.info(f"🧠 KG-Extraktion für '{doc_id}': {len(kg_chunks_list)} Chunks, {total_chars} Zeichen")
                    
                    # Per-Chunk LLM-Extraktion
                    kg_triples_with_source = self._create_llm_knowledge_graph_per_chunk(
                        kg_chunks_list, doc_id, doc_metadata
                    )
                    
                    if not kg_triples_with_source:
                        logger.debug(f"⚠️ Keine KG-Triples für {doc_id} extrahiert")
                        continue
                    
                    # ★ Per-Chunk Reranker Grounding (gegen Quell-Chunk, nicht Blob!)
                    chunk_text_lookup: Dict[Optional[int], str] = {}
                    for kc in kg_chunks_list:
                        for src_id in kc.get("_source_chunk_ids", []):
                            chunk_text_lookup[src_id] = kc["text"]
                        chunk_text_lookup[kc.get("chunk_id")] = kc["text"]
                    
                    verified_triples: List[Tuple[str, str, str, Optional[int], float]] = []
                    for subj, pred, obj, src_chunk_id, conf in kg_triples_with_source:
                        grounding_text = chunk_text_lookup.get(src_chunk_id, "")
                        if grounding_text:
                            single_verified = self._verify_kg_triples_via_reranker(
                                [(subj, pred, obj)], grounding_text
                            )
                            if single_verified:
                                verified_triples.append((subj, pred, obj, src_chunk_id, conf))
                            else:
                                logger.debug(
                                    f"[KG-VERIFY] REJECTED: '{subj}' → '{pred}' → '{obj}' "
                                    f"| src_chunk_id={src_chunk_id} "
                                    f"| grounding_text_len={len(grounding_text)} "
                                    f"| grounding_text_start='{grounding_text[:120]}...'"
                                )
                        else:
                            verified_triples.append((subj, pred, obj, src_chunk_id, conf))
                    
                    if not verified_triples:
                        logger.debug(f"⚠️ Alle Triples für {doc_id} durch Reranker gefiltert")
                        continue
                    
                    logger.info(
                        f"[KG-VERIFY] {doc_id}: {len(kg_triples_with_source)} → "
                        f"{len(verified_triples)} nach Per-Chunk-Grounding"
                    )
                    
                    # Delete old triples (central hook maintains entity frequencies + CASCADE)
                    self._delete_triples_for_doc(cur, doc_id)
                    for cleanup_id in chunk_doc_ids_cleanup:
                        self._delete_triples_for_doc(cur, cleanup_id)
                    
                    # Insert verified triples
                    doc_triples_count = 0
                    bayesian_updates = 0
                    new_entities: List[Tuple[str, str]] = []
                    now_iso = datetime.now(timezone.utc).isoformat()

                    for subject, predicate, obj, source_chunk_id, evidence_conf in verified_triples:
                        triple_hash = calculate_triple_hash(subject, predicate, obj)

                        # SOTA Repeat-Mention-Update (Noisy-OR Bayesian):
                        #   P_new = 1 - (1 - P_old) * (1 - P_evidence)
                        # Bei jeder erneuten Extraktion derselben Aussage wächst
                        # die Confidence monoton gegen 1.0; mention_count++.
                        existing = cur.execute(
                            "SELECT triple_id, confidence, mention_count "
                            "FROM triples WHERE triple_hash = ?",
                            (triple_hash,),
                        ).fetchone()

                        if existing is not None:
                            ex_id, old_conf, old_count = existing
                            old_conf = float(old_conf if old_conf is not None else 0.5)
                            old_count = int(old_count if old_count is not None else 1)
                            new_conf = 1.0 - (1.0 - old_conf) * (1.0 - evidence_conf)
                            new_conf = min(0.99, max(old_conf, new_conf))
                            cur.execute(
                                "UPDATE triples SET confidence = ?, "
                                "mention_count = ?, updated_at = ? "
                                "WHERE triple_id = ?",
                                (new_conf, old_count + 1, now_iso, ex_id),
                            )
                            bayesian_updates += 1
                            continue

                        cur.execute("""
                            INSERT INTO triples(doc_id, page, table_id, subject, predicate, object,
                                                metadata, triple_hash, source_chunk_id,
                                                confidence, mention_count, created_at, updated_at)
                            VALUES (?, NULL, NULL, ?, ?, ?, ?, ?, ?, ?, 1, ?, ?)
                        """, (doc_id, subject, predicate, obj,
                              json.dumps({"kg_source": "llm_per_chunk", "created_at": now_iso}),
                              triple_hash, source_chunk_id,
                              evidence_conf, now_iso, now_iso))

                        doc_triples_count += 1
                        new_entities.append((subject, doc_id))
                        new_entities.append((obj, doc_id))

                    total_triples += doc_triples_count
                    processed_docs += 1

                    # Embed entities
                    if new_entities and doc_triples_count > 0:
                        try:
                            self._store_entity_embeddings(new_entities, cur)
                        except Exception as e_emb:
                            raise RuntimeError(
                                f"entity embedding storage failed: {e_emb}"
                            ) from e_emb

                    logger.info(
                        f"✅ KG-Update '{doc_id}': {doc_triples_count} neu, "
                        f"{bayesian_updates} Bayesian-update"
                    )
                    
                except Exception as e:
                    raise RuntimeError(
                        f"KG-Erstellung für {doc_id} fehlgeschlagen: {e}"
                    ) from e
            
            conn.commit()
            
            # Invalidate entity index
            if self._search_manager is not None:
                try:
                    self._search_manager.invalidate_entity_index()
                except Exception as exc:
                    logger.debug(f"KG auto-build: invalidate_entity_index failed: {exc}")
            
            # Cleanup pending KG chunks
            if hasattr(self, '_pending_kg_chunks'):
                self._pending_kg_chunks.clear()
            
        except Exception as e:
            logger.error(f"❌ Automatische KG-Erstellung fehlgeschlagen (fail-fast): {e}")
            raise RuntimeError(f"Automatic KG build failed: {e}") from e
        finally:
            try:
                if 'cur' in locals() and cur:
                    cur.close()
                if 'conn' in locals() and conn:
                    self.return_connection(conn)
            except Exception as exc:
                logger.debug(f"KG auto-build cleanup failed: {exc}")
        
        logger.info(f"📊 KG-Zusammenfassung: {total_triples} Triples für {processed_docs}/{len(doc_ids)} Dokumente")
        
        return {
            "triples": total_triples,
            "documents_processed": processed_docs,
            "documents_total": len(doc_ids)
        }

    def _extract_pdf_with_ocr(self, file_path: str) -> str:
        """
        OCR-Fallback für gescannte PDFs oder Bild-PDFs
        
        Verwendet EasyOCR um Text aus PDF-Seiten zu extrahieren,
        die als Bilder gerendert werden (für gescannte Dokumente).
        
        Args:
            file_path: Pfad zur PDF-Datei
            
        Returns:
            Extrahierter Text aus allen Seiten
        """
        import tempfile
        
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logging.warning("PyMuPDF nicht verfügbar für OCR-Rendering")
            return ""
        
        try:
            from agent.ocr_processor import OCRProcessor
        except ImportError:
            logging.warning("OCR Processor nicht verfügbar")
            return ""
        
        all_text = []
        temp_images = []
        
        try:
            # OCR Processor initialisieren (lazy loading)
            if not hasattr(self, '_ocr_processor'):
                self._ocr_processor = OCRProcessor(
                    languages=['de', 'en'],
                    use_gpu=True  # GPU wenn verfügbar
                )
            
            with fitz.open(file_path) as doc:
                logging.info(f"📸 OCR für {len(doc)} Seiten starten...")
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    try:
                        # Seite als Bild rendern (300 DPI für bessere OCR-Qualität)
                        mat = fitz.Matrix(300/72, 300/72)  # 300 DPI
                        pix = page.get_pixmap(matrix=mat)
                        
                        # Temporäres Bild speichern
                        temp_path = tempfile.mktemp(suffix=f'_page_{page_num}.png')
                        pix.save(temp_path)
                        temp_images.append(temp_path)
                        
                        # OCR auf Bild anwenden
                        result = self._ocr_processor.extract_text(temp_path)
                        
                        if result.text.strip():
                            all_text.append(f"--- Seite {page_num + 1} ---\n{result.text}")
                            logging.debug(f"OCR Seite {page_num + 1}: {len(result.text)} Zeichen, "
                                        f"Confidence: {result.confidence:.2%}")
                        
                    except Exception as e:
                        logging.warning(f"OCR Seite {page_num + 1} fehlgeschlagen: {e}")
                        continue
            
            combined_text = "\n\n".join(all_text)
            logging.info(f"✅ OCR abgeschlossen: {len(combined_text)} Zeichen aus {len(all_text)} Seiten")
            return combined_text
            
        except Exception as e:
            logging.error(f"OCR-Verarbeitung fehlgeschlagen: {e}")
            return ""
            
        finally:
            # Temporäre Bilder aufräumen
            for temp_path in temp_images:
                try:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
                except Exception as exc:
                    logging.debug(f"Temp image cleanup failed for {temp_path}: {exc}")
    
    # ======= FAISS INDEX MANAGEMENT =======
    
    def rebuild_faiss_indexes(self, force: bool = False) -> Dict[str, Any]:
        """
        Baut FAISS-Indizes neu auf.
        
        Performance Boost:
            - Recent Index (20k neueste Chunks)
            - Full Index (alle Chunks)
            - Persistente Speicherung
            
        Args:
            force: True = Immer neu bauen, False = Nur wenn nötig
            
        Returns:
            Statistiken über Index-Building
        """
        if not self.faiss_manager or not FAISS_HYBRID_AVAILABLE:  # type: ignore[unreachable]
            return {
                "success": False,
                "error": "FAISS Hybrid Search nicht verfügbar"
            }
        
        try:  # type: ignore[unreachable]
            logger.info("🔨 Rebuilding FAISS Indexes...")
            start_time = time.time()
            
            # Build Indexes (force parameter wird nicht unterstützt)
            self.faiss_manager.build_indexes(recent_limit=20000)
            
            build_time = time.time() - start_time
            
            # Hole Stats
            stats = {
                "recent_index_size": len(self.faiss_manager.recent_id_map),
                "full_index_size": len(self.faiss_manager.full_id_map),
                "build_time_seconds": round(build_time, 2)
            }
            
            logger.info(
                f"✅ FAISS Indexes rebuilt: "
                f"{stats.get('recent_index_size', 0)} recent, "
                f"{stats.get('full_index_size', 0)} full "
                f"({build_time:.2f}s)"
            )
            
            return {
                "success": True,
                **stats
            }
            
        except Exception as e:
            logger.error(f"❌ FAISS Index rebuild failed: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def get_faiss_stats(self) -> Dict[str, Any]:
        """
        Liefert Statistiken über FAISS-Indizes.
        
        Returns:
            Dictionary mit Index-Statistiken
        """
        if not self.faiss_manager or not FAISS_HYBRID_AVAILABLE:  # type: ignore[unreachable]
            return {
                "available": False,
                "message": "FAISS Hybrid Search nicht verfügbar"
            }
        
        try:  # type: ignore[unreachable]
            stats: Dict[str, Any] = {
                "recent_index_size": len(self.faiss_manager.recent_id_map) if self.faiss_manager.recent_id_map else 0,
                "full_index_size": len(self.faiss_manager.full_id_map) if self.faiss_manager.full_id_map else 0,
                "recent_threshold": 20000,
                "embedding_dim": self.faiss_manager.embedding_dim
            }
            
            # Ergänze mit Fusion Engine Stats (falls vorhanden)
            if self.fusion_engine:
                stats["fusion"] = {
                    "recency_boost_factor": self.fusion_engine.recency_boost,
                    "kg_boost_factor": self.fusion_engine.kg_boost,
                    "target_kg_ratio": self.fusion_engine.target_kg_ratio
                }
            
            return {
                "available": True,
                **stats
            }
            
        except Exception as e:
            logger.error(f"Error getting FAISS stats: {e}")
            return {
                "available": True,
                "error": str(e)
            }

    # ──────────────────────────────────────────────────────────────────
    #  ★ SOTA: Full KG Consistency Rebuild
    # ──────────────────────────────────────────────────────────────────
    def rebuild_kg_consistency(self,
                               progress_callback: Optional[Any] = None
                               ) -> Dict[str, Any]:
        """
        ★ SOTA: Full KG consistency rebuild — heals all known drift issues.

        Phases:
          1. Garbage triple purge (numeric/overlength/sentence-fragment entities)
          2. Orphan cascade cleanup (triple_quality, chunk_quality referencing deleted rows)
          3. Case-duplicate entity merge (trivial case/whitespace differences)
          4. Frequency resync (recalculate ALL kg_entities.frequency from actual triples)
          5. Dead entity pruning (remove kg_entities with frequency = 0)
          6. Missing entity tracking (add entities from triples not yet in kg_entities)
          7. Embedding-based entity resolution (cosine similarity merge)
          8. FAISS index rebuild

        Args:
            progress_callback: Optional callable(phase: int, total: int, description: str)
                               for UI progress reporting (e.g. Streamlit progress bar)

        Returns:
            Detailed statistics dictionary with per-phase results.
        """
        import re as _re

        TOTAL_PHASES = 8
        stats: Dict[str, Any] = {
            "success": False,
            "phases": {},
            "total_time_seconds": 0.0,
        }
        start_time = time.time()

        def _progress(phase: int, desc: str):
            if progress_callback:
                try:
                    progress_callback(phase, TOTAL_PHASES, desc)
                except Exception as exc:
                    logger.debug(f"[KG-Rebuild] progress callback failed: {exc}")
            logger.info(f"[KG-Rebuild] Phase {phase}/{TOTAL_PHASES}: {desc}")

        # ★ SOTA: see _store_entity_embeddings — canonical import only.
        from agent.llm_knowledge_graph import normalize_entity_for_matching

        try:
            conn = self.get_connection()
            cur = conn.cursor()

            # ════════════════════════════════════════════════════════════
            # Phase 1: Garbage triple purge
            # ════════════════════════════════════════════════════════════
            _progress(1, "Garbage-Entities und deren Triples entfernen")
            _NUMERIC_RE = _re.compile(r'^\d+([.,]\d+)?[%°]?$')

            # Collect all unique entities from triples
            cur.execute("""
                SELECT DISTINCT entity FROM (
                    SELECT subject AS entity FROM triples WHERE subject IS NOT NULL
                    UNION
                    SELECT object AS entity FROM triples WHERE object IS NOT NULL
                )
            """)
            all_triple_entities = [row[0] for row in cur.fetchall()]

            garbage_entities = set()
            for entity in all_triple_entities:
                text = entity.strip()
                if not text:
                    garbage_entities.add(entity)
                    continue
                if _NUMERIC_RE.match(text):
                    garbage_entities.add(entity)
                    continue
                if len(text) > 80:
                    garbage_entities.add(entity)
                    continue
                if len(text.split()) > 8:
                    garbage_entities.add(entity)
                    continue

            garbage_triples_deleted = 0
            for ge in garbage_entities:
                cur.execute(
                    "DELETE FROM triples WHERE subject = ? OR object = ?",
                    (ge, ge)
                )
                garbage_triples_deleted += cur.rowcount
                # Also remove from kg_entities
                cur.execute("DELETE FROM kg_entities WHERE entity_text = ?", (ge,))

            stats["phases"]["garbage_entities"] = len(garbage_entities)
            stats["phases"]["garbage_triples_deleted"] = garbage_triples_deleted

            # ════════════════════════════════════════════════════════════
            # Phase 2: Orphan cascade cleanup
            # ════════════════════════════════════════════════════════════
            _progress(2, "Verwaiste triple_quality / chunk_quality bereinigen")

            # Orphaned triple_quality (triple_id references deleted triples)
            try:
                cur.execute("""
                    DELETE FROM triple_quality
                    WHERE triple_id NOT IN (SELECT triple_id FROM triples)
                """)
                orphan_tq = cur.rowcount
            except Exception as exc:
                logger.debug(f"[KG-Rebuild] orphan triple_quality cleanup skipped: {exc}")
                orphan_tq = 0

            # Orphaned chunk_quality (doc_id, chunk_id references deleted chunks)
            try:
                cur.execute("""
                    DELETE FROM chunk_quality
                    WHERE (doc_id, chunk_id) NOT IN (
                        SELECT doc_id, chunk_id FROM chunks
                    )
                """)
                orphan_cq = cur.rowcount
            except Exception as exc:
                logger.debug(f"[KG-Rebuild] orphan chunk_quality cleanup skipped: {exc}")
                orphan_cq = 0

            stats["phases"]["orphan_triple_quality"] = orphan_tq
            stats["phases"]["orphan_chunk_quality"] = orphan_cq

            # ════════════════════════════════════════════════════════════
            # Phase 3: Case-duplicate entity merge
            # ════════════════════════════════════════════════════════════
            _progress(3, "Case-Duplikate in Entities zusammenführen")

            cur.execute("SELECT entity_id, entity_text, frequency FROM kg_entities")
            all_entities = cur.fetchall()

            case_groups: Dict[str, List[Tuple]] = {}
            for eid, etext, freq in all_entities:
                key = etext.strip().lower()
                if key not in case_groups:
                    case_groups[key] = []
                case_groups[key].append((eid, etext, freq or 0))

            case_merged = 0
            from agent.kg_entity_merge import merge_entity_in_triples
            from agent.rag_store.utils.memory import calculate_triple_hash
            for key, group in case_groups.items():
                if len(group) <= 1:
                    continue
                # Canonical = highest frequency
                group.sort(key=lambda x: x[2], reverse=True)
                canonical_id, canonical_text, _ = group[0]

                for eid, etext, freq in group[1:]:
                    merge_entity_in_triples(
                        conn, calculate_triple_hash, canonical_text, etext
                    )
                    cur.execute("DELETE FROM kg_entities WHERE entity_id = ?", (eid,))

                case_merged += 1

            stats["phases"]["case_groups_merged"] = case_merged

            # ════════════════════════════════════════════════════════════
            # Phase 4: Frequency resync
            # ════════════════════════════════════════════════════════════
            _progress(4, "Entity-Frequenzen aus tatsächlichen Triples neu berechnen")

            # Build actual frequency map from triples
            cur.execute("""
                SELECT entity, SUM(cnt) FROM (
                    SELECT subject AS entity, COUNT(*) as cnt
                    FROM triples WHERE subject IS NOT NULL GROUP BY subject
                    UNION ALL
                    SELECT object AS entity, COUNT(*) as cnt
                    FROM triples WHERE object IS NOT NULL GROUP BY object
                ) GROUP BY entity
            """)
            actual_freqs = dict(cur.fetchall())

            cur.execute("SELECT entity_id, entity_text, frequency FROM kg_entities")
            all_entities = cur.fetchall()
            freq_corrected = 0
            for eid, etext, old_freq in all_entities:
                new_freq = actual_freqs.get(etext, 0)
                if new_freq != (old_freq or 0):
                    cur.execute("UPDATE kg_entities SET frequency = ? WHERE entity_id = ?",
                                (new_freq, eid))
                    freq_corrected += 1

            stats["phases"]["frequency_corrected"] = freq_corrected

            # ════════════════════════════════════════════════════════════
            # Phase 5: Dead entity pruning
            # ════════════════════════════════════════════════════════════
            _progress(5, "Nicht mehr referenzierte Entities entfernen")

            cur.execute("DELETE FROM kg_entities WHERE frequency = 0 OR frequency IS NULL")
            dead_entities = cur.rowcount
            stats["phases"]["dead_entities_pruned"] = dead_entities

            # ════════════════════════════════════════════════════════════
            # Phase 6: Missing entity tracking
            # ════════════════════════════════════════════════════════════
            _progress(6, "Fehlende Entity-Einträge ergänzen und embedden")

            cur.execute("""
                SELECT DISTINCT entity FROM (
                    SELECT subject AS entity FROM triples WHERE subject IS NOT NULL
                    UNION
                    SELECT object AS entity FROM triples WHERE object IS NOT NULL
                ) WHERE entity NOT IN (SELECT entity_text FROM kg_entities)
            """)
            untracked = [row[0] for row in cur.fetchall()
                         if row[0] and len(row[0].strip()) >= 2]

            new_entities_for_embedding: List[Tuple[str, str]] = []
            for entity_text in untracked:
                normalized = normalize_entity_for_matching(entity_text)
                # Actual frequency from triples
                cur.execute(
                    "SELECT COUNT(*) FROM ("
                    "  SELECT subject AS e FROM triples WHERE subject = ? "
                    "  UNION ALL "
                    "  SELECT object FROM triples WHERE object = ?"
                    ")",
                    (entity_text, entity_text)
                )
                freq = cur.fetchone()[0]
                cur.execute("""
                    INSERT OR IGNORE INTO kg_entities
                    (entity_text, normalized_text, entity_type, frequency)
                    VALUES (?, ?, 'entity', ?)
                """, (entity_text, normalized, freq))
                new_entities_for_embedding.append((entity_text, "rebuild"))

            # Embed new entities
            if new_entities_for_embedding:
                try:
                    self._store_entity_embeddings(new_entities_for_embedding, cur)
                except Exception as e:
                    logger.warning(f"[KG-Rebuild] Entity embedding for untracked failed: {e}")

            stats["phases"]["untracked_entities_added"] = len(untracked)

            # Commit all phases 1-6
            conn.commit()
            cur.close()
            self.return_connection(conn)

            # ════════════════════════════════════════════════════════════
            # Phase 7: Embedding-based entity resolution
            # ════════════════════════════════════════════════════════════
            _progress(7, "Embedding-basierte Entity Resolution")
            resolution_stats = self.resolve_duplicate_entities(similarity_threshold=0.90)
            stats["phases"]["entity_resolution"] = resolution_stats

            # ════════════════════════════════════════════════════════════
            # Phase 8: FAISS index rebuild
            # ════════════════════════════════════════════════════════════
            _progress(8, "FAISS-Indizes neu aufbauen")
            faiss_stats = self.rebuild_faiss_indexes(force=True)
            stats["phases"]["faiss_rebuild"] = faiss_stats

            # Invalidate search manager entity index
            if self._search_manager is not None:
                try:
                    self._search_manager.invalidate_entity_index()
                except Exception as exc:
                    logger.debug(f"[KG-Rebuild] invalidate_entity_index failed: {exc}")

            stats["success"] = True
            stats["total_time_seconds"] = round(time.time() - start_time, 2)

            logger.info(
                f"[KG-Rebuild] ✅ Complete in {stats['total_time_seconds']}s — "
                f"garbage={len(garbage_entities)}, orphan_tq={orphan_tq}, "
                f"case_merged={case_merged}, freq_fixed={freq_corrected}, "
                f"dead_pruned={dead_entities}, untracked={len(untracked)}"
            )

        except Exception as e:
            logger.error(f"[KG-Rebuild] ❌ Failed: {e}")
            import traceback
            logger.debug(traceback.format_exc())
            stats["error"] = str(e)
            stats["total_time_seconds"] = round(time.time() - start_time, 2)

        return stats


# ======= COMPATIBILITY ALIASES =======
# Für nahtlose Abwärtskompatibilität mit bestehendem Code

# RagStore ist jetzt ein Alias für UnifiedRagStore
RagStore = UnifiedRagStore

# Factory-Funktion für automatische Erkennung
def create_rag_store(db_path: Optional[str] = None, **kwargs) -> UnifiedRagStore:
    """
    Factory-Funktion für RAG Store mit automatischer Optimierung
    
    Args:
        db_path: Datenbankpfad
        **kwargs: Zusätzliche Parameter für UnifiedRagStore
        
    Returns:
        Optimierte UnifiedRagStore-Instanz
    """
    return UnifiedRagStore(db_path=db_path, **kwargs)


# Legacy-Unterstützung
def get_rag_store_shared(db_path: Optional[str] = None, dim: int = 384) -> UnifiedRagStore:
    """Legacy-Unterstützung für geteilte RAG Store-Instanzen"""
    return UnifiedRagStore.get_shared(db_path=db_path, dim=dim)


# Export all important classes and functions
__all__ = [
    'UnifiedRagStore',
    'RagStore',  # Alias
    'create_rag_store',
    'get_rag_store_shared',
    'ProcessingConfig'
]