"""
PowerPoint (PPTX) Document Processor
Extrahiert Text, Bilder, Tabellen und Metadaten aus PowerPoint-Präsentationen.
"""

import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx nicht verfügbar. PPTX-Verarbeitung deaktiviert.")

logger = logging.getLogger(__name__)


class PowerPointProcessor:
    """Verarbeitet PowerPoint-Dateien (PPTX) für RAG-Indexierung."""
    
    def __init__(self):
        self.supported_extensions = ['.pptx']
        
    def can_process(self, file_path: str) -> bool:
        """Prüft, ob die Datei verarbeitet werden kann."""
        if not PPTX_AVAILABLE:
            return False
        return Path(file_path).suffix.lower() in self.supported_extensions
    
    def process(self, file_path: str) -> Dict[str, Any]:
        """
        Verarbeitet eine PowerPoint-Datei.
        
        Returns:
            Dict mit:
            - text: Extrahierter Text (alle Slides)
            - metadata: Metadaten (Autor, Datum, Slide-Count, etc.)
            - slides: Liste von Slide-Daten
            - images: Extrahierte Bilder (optional)
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx ist nicht installiert")
            
        try:
            prs = Presentation(file_path)
            
            # Metadaten extrahieren
            metadata = self._extract_metadata(prs, file_path)
            
            # Slides verarbeiten
            slides_data = []
            all_text = []
            
            for idx, slide in enumerate(prs.slides, start=1):
                slide_data = self._process_slide(slide, idx)
                slides_data.append(slide_data)
                all_text.append(f"\n\n--- Slide {idx} ---\n{slide_data['text']}")
            
            # Kombinierter Text für RAG
            combined_text = "\n".join(all_text)
            
            return {
                'text': combined_text,
                'metadata': metadata,
                'slides': slides_data,
                'slide_count': len(prs.slides),
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Fehler beim Verarbeiten von {file_path}: {e}")
            return {
                'text': '',
                'metadata': {'error': str(e)},
                'slides': [],
                'success': False
            }
    
    def _extract_metadata(self, prs: Presentation, file_path: str) -> Dict[str, Any]:
        """Extrahiert Metadaten aus der Präsentation."""
        core_props = prs.core_properties
        
        metadata = {
            'file_name': Path(file_path).name,
            'file_path': str(file_path),
            'file_size': Path(file_path).stat().st_size,
            'slide_count': len(prs.slides),
            'title': core_props.title or 'Unbekannt',
            'author': core_props.author or 'Unbekannt',
            'subject': core_props.subject or '',
            'created': core_props.created.isoformat() if core_props.created else None,
            'modified': core_props.modified.isoformat() if core_props.modified else None,
            'processed_at': datetime.now().isoformat(),
            'processor': 'PowerPointProcessor'
        }
        
        return metadata
    
    def _process_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """Verarbeitet einen einzelnen Slide."""
        slide_text = []
        tables_data = []
        images_data = []
        
        for shape in slide.shapes:
            # Text extrahieren
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Tabellen extrahieren
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = self._extract_table(shape.table)
                if table_data:
                    tables_data.append(table_data)
                    slide_text.append(f"\n{table_data}\n")
            
            # Bilder extrahieren (Metadaten nur)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data = self._extract_image_metadata(shape)
                if image_data:
                    images_data.append(image_data)
        
        # Notes extrahieren (Speaker Notes)
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text.append(f"\nNotes: {notes_text}")
        
        return {
            'slide_number': slide_number,
            'text': "\n".join(slide_text),
            'notes': notes_text,
            'tables': tables_data,
            'images': images_data,
            'shape_count': len(slide.shapes)
        }
    
    def _extract_table(self, table) -> Optional[str]:
        """Extrahiert Tabellendaten als Markdown."""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            
            if not rows:
                return None
            
            # Als Markdown-Tabelle formatieren
            markdown_table = []
            
            # Header
            if len(rows) > 0:
                markdown_table.append("| " + " | ".join(rows[0]) + " |")
                markdown_table.append("|" + "|".join(["---"] * len(rows[0])) + "|")
            
            # Data rows
            for row in rows[1:]:
                markdown_table.append("| " + " | ".join(row) + " |")
            
            return "\n".join(markdown_table)
            
        except Exception as e:
            logger.warning(f"Fehler beim Extrahieren einer Tabelle: {e}")
            return None
    
    def _extract_image_metadata(self, shape) -> Optional[Dict[str, Any]]:
        """Extrahiert Metadaten von Bildern (nicht das Bild selbst)."""
        try:
            image = shape.image
            return {
                'content_type': image.content_type,
                'size': len(image.blob),
                'width': shape.width,
                'height': shape.height,
                'name': shape.name
            }
        except Exception as e:
            logger.warning(f"Fehler beim Extrahieren von Bild-Metadaten: {e}")
            return None


# Singleton-Instanz
_pptx_processor = None

def get_pptx_processor() -> PowerPointProcessor:
    """Liefert Singleton-Instanz des PowerPoint-Processors."""
    global _pptx_processor
    if _pptx_processor is None:
        _pptx_processor = PowerPointProcessor()
    return _pptx_processor
